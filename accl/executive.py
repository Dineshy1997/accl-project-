import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.table import Table
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import math
import traceback
import uuid
from datetime import datetime
from dateutil.relativedelta import relativedelta
import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize session state for file uploads
if 'sales_file' not in st.session_state:
    st.session_state.sales_file = None
if 'ly_sales_file' not in st.session_state:
    st.session_state.ly_sales_file = None
if 'budget_file' not in st.session_state:
    st.session_state.budget_file = None
if 'os_jan_file' not in st.session_state:
    st.session_state.os_jan_file = None
if 'os_feb_file' not in st.session_state:
    st.session_state.os_feb_file = None
if 'logo_file' not in st.session_state:
    st.session_state.logo_file = None


def extract_executive_name(executive):
    if pd.isna(executive) or str(executive).strip() == '':
        return 'BLANK'
    return str(executive).strip().upper()

def get_excel_sheets(file):
    try:
        xl = pd.ExcelFile(file)
        return xl.sheet_names
    except Exception as e:
        st.error(f"Error reading Excel sheets: {e}")
        return []

def create_title_slide(prs, title, logo_file=None):
    """Updated title slide creation with page number"""
    blank_slide_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(blank_slide_layout)
    
    # Company name
    company_name = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    company_frame = company_name.text_frame
    company_frame.text = "ASIA CRYSTAL COMMODITY LLP"
    p = company_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    # Logo
    if logo_file is not None:
        try:
            logo_buffer = BytesIO(logo_file.read())
            logo = title_slide.shapes.add_picture(logo_buffer, Inches(5.665), Inches(1.5), width=Inches(2), height=Inches(2))
            logo_file.seek(0)
        except Exception as e:
            logger.error(f"Error adding logo to slide: {e}")
    
    # Title
    title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title.upper()
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Times New Roman"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    # Subtitle removed - no longer adding ACCLLP text
    
    # Add page number (1) to title slide
    add_page_number(title_slide, 1)
    
    return title_slide



def add_page_number(slide, page_num, total_pages=None):
    """Add page number to bottom right corner of slide"""
    try:
        # Create text box for page number in bottom right corner
        left = Inches(11.5)  # Position from left
        top = Inches(6.8)    # Position from top (near bottom)
        width = Inches(1.5)  # Width of text box
        height = Inches(0.5) # Height of text box
        
        page_textbox = slide.shapes.add_textbox(left, top, width, height)
        page_frame = page_textbox.text_frame
        
        # Set page number text
        if total_pages:
            page_frame.text = f"{page_num}/{total_pages}"
        else:
            page_frame.text = str(page_num)
        
        # Format the page number
        p = page_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = "Times New Roman"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(128, 128, 128)  # Gray color
        
    except Exception as e:
        logger.error(f"Error adding page number to slide: {e}")

def add_table_slide(prs, df, title, percent_cols=None):
    """Updated table slide creation with page numbering and consistent formatting"""
    if percent_cols is None:
        percent_cols = []
    
    if df is None or df.empty:
        logger.warning(f"Skipping slide for {title}: DataFrame is None or empty")
        return
    
    # Constants
    MAX_EXECUTIVES_PER_SLIDE = 15
    
    # Remove any rows with 'ACCLP' if present
    df = df[df.iloc[:, 0] != "ACCLP"].copy()
    
    # Separate regular data from TOTAL/GRAND TOTAL row if present
    total_row = None
    total_keywords = ['TOTAL', 'GRAND TOTAL']
    
    # Check for any total row
    total_mask = df.iloc[:, 0].isin(total_keywords)
    if total_mask.any():
        total_row = df[total_mask]
        df = df[~total_mask]  # Remove total rows from main data
    
    num_executives = len(df)
    num_slides = math.ceil(num_executives / MAX_EXECUTIVES_PER_SLIDE)
    
    # Calculate starting page number (assuming title slide is page 1)
    current_slide_count = len(prs.slides)
    
    # Process each slide
    for slide_num in range(num_slides):
        start_idx = slide_num * MAX_EXECUTIVES_PER_SLIDE
        end_idx = start_idx + MAX_EXECUTIVES_PER_SLIDE
        slide_df = df.iloc[start_idx:end_idx].copy()
        
        # On last slide, try to include TOTAL row if there's space
        if slide_num == num_slides - 1 and total_row is not None:
            if len(slide_df) + len(total_row) <= MAX_EXECUTIVES_PER_SLIDE:
                slide_df = pd.concat([slide_df, total_row], ignore_index=True)
                total_row = None  # Mark as used
        
        # Create slide title - Use original title without slide numbers
        slide_title = title
        
        slide = _create_single_table_slide(prs, slide_df, slide_title, percent_cols)
        
        # Add page number to the slide
        page_number = current_slide_count + slide_num + 1
        add_page_number(slide, page_number)
    
    # Add final slide with TOTAL row if it wasn't included earlier
    if total_row is not None and not total_row.empty:
        slide = _create_single_table_slide(prs, total_row, title, percent_cols)
        page_number = len(prs.slides)
        add_page_number(slide, page_number)

def _create_single_table_slide(prs, df, title, percent_cols):
    """Updated helper function to create a single slide with consistent formatting"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with 15pt font
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title.upper()
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)  
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    
    # Prepare table data
    columns = df.columns.tolist()
    num_rows = len(df) + 1  # +1 for header row
    num_cols = len(columns)
    
    # Calculate table height - fixed row height of 0.3 inches
    row_height = Inches(0.3)
    table_height = row_height * num_rows
    
    # Add table (fixed width of 12 inches)
    table = slide.shapes.add_table(
        num_rows, num_cols, 
        Inches(0.5), Inches(1.5),  # Left, Top position
        Inches(12), table_height    # Width, Height
    ).table
    
    # Set column widths
    if num_cols > 0:
        table.columns[0].width = Inches(3.0)  # First column (Executive) wider
    remaining_width = 12.0 - 3.0
    if num_cols > 1:
        col_width = remaining_width / (num_cols - 1)
        for i in range(1, num_cols):
            table.columns[i].width = Inches(col_width)
    
    # Add header row with 15pt font and uppercase text
    for i, col_name in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = str(col_name).upper()
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.size = Pt(15)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add data rows with 14pt font and consistent formatting
    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        is_total_row = 'TOTAL' in str(row.iloc[0])
        for col_idx, col_name in enumerate(columns):
            cell = table.cell(row_idx, col_idx)
            value = row[col_name]
            
            # Format cells with consistent decimal places
            if col_idx in percent_cols and isinstance(value, (int, float)) and not pd.isna(value):
                cell.text = f"{value:.2f}%"
            else:
                # For numeric columns (not first column), ensure 2-decimal formatting
                if col_idx > 0 and isinstance(value, (int, float)) and not pd.isna(value):
                    if "%" not in str(df.columns[col_idx]):
                        cell.text = f"{float(value):.2f}"
                    else:
                        cell.text = f"{value:.2f}"
                elif col_idx == 0:
                    cell.text = str(value).upper()
                else:
                    cell.text = str(value)
            
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            
            if is_total_row:
                cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                cell.text_frame.paragraphs[0].font.bold = True
            else:
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    return slide
def create_table_image(df, title, percent_cols=None):
    if percent_cols is None:
        percent_cols = []
    fig, ax = plt.subplots(figsize=(12, len(df) * 0.5))
    ax.axis('off')
    columns = df.columns.tolist()
    rows = len(df)
    ncols = len(columns)
    table = Table(ax, bbox=[0, 0, 1, 1])
    for col_idx, col_name in enumerate(columns):
        table.add_cell(0, col_idx, 1.0/ncols, 1.0/rows, text=col_name, loc='center', facecolor='#0070C0')
        table[0, col_idx].set_text_props(weight='bold', color='white', fontsize=12)
    for row_idx in range(len(df)):
        for col_idx, col_name in enumerate(columns):
            value = df.iloc[row_idx, col_idx]
            if col_idx in percent_cols and isinstance(value, (int, float)) and not pd.isna(value):
               text = f"{value:.2f}%"  # Ensure 2 decimal places for percentages
            else:
                 if isinstance(value, (int, float)) and not pd.isna(value):
                  text = f"{value:.2f}"  # Ensure 2 decimal places for numeric values
                 else:
                   text = str(value)

            is_total_row = 'TOTAL' in str(df.iloc[row_idx, 0])
            if is_total_row:
                facecolor = '#D3D3D3'
                fontweight = 'bold'
            else:
                facecolor = '#DDEBF7' if row_idx % 2 == 0 else 'white'
                fontweight = 'normal'
            table.add_cell(row_idx + 1, col_idx, 1.0/ncols, 1.0/rows, text=text, loc='center', facecolor=facecolor)
            table[row_idx + 1, col_idx].set_text_props(fontsize=10, weight=fontweight)
    if ncols > 0:
        table[(0, 0)].width = 0.25
    if ncols > 1:
        remaining_width = 0.75
        col_width = remaining_width / (ncols - 1)
        for i in range(1, ncols):
            table[(0, i)].width = col_width
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    ax.add_table(table)
    plt.suptitle(title, fontsize=16, weight='bold', color='#0070C0', y=1.02)
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
    plt.close()
    img_buffer.seek(0)
    return img_buffer

def create_consolidated_ppt(dfs_info, logo_file=None, title="Consolidated Report"):
    """Updated consolidated PPT creation with page numbers and thank you slide"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide (page 1)
        create_title_slide(prs, title, logo_file)
        
        # Process each DataFrame
        for df_info in dfs_info:
            df = df_info['df']
            slide_title = df_info['title']
            percent_cols = df_info.get('percent_cols', [])
            add_table_slide(prs, df, slide_title, percent_cols)
        
        # Add thank you slide at the end
        create_thank_you_slide(prs, logo_file)
        
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer
    except Exception as e:
        logger.error(f"Error creating consolidated PPT: {e}")
        st.error(f"Error creating consolidated PPT: {e}")
        return None
    
def create_proof_of_calculation_excel(budget_df, sales_df, selected_month, 
                                      budget_exec_col, budget_exec_code_col, budget_area_col, 
                                      budget_sl_code_col, budget_product_group_col, 
                                      budget_qty_col, budget_value_col,
                                      sales_exec_col, sales_date_col, sales_area_col,
                                      sales_sl_code_col, sales_product_group_col,
                                      sales_qty_col, sales_value_col,
                                      selected_executives, selected_branches=None):
    """
    Creates a detailed Excel file showing the proof of calculation for budget vs billed mapping.
    
    Returns:
    - BytesIO buffer containing the Excel file
    """
    try:
        # Filter sales data for the selected month
        sales_df = sales_df.copy()
        budget_df = budget_df.copy()
        
        # Convert sales date and filter by month
        sales_df[sales_date_col] = pd.to_datetime(sales_df[sales_date_col], dayfirst=True, errors='coerce')
        filtered_sales_df = sales_df[sales_df[sales_date_col].dt.strftime('%b %y') == selected_month].copy()
        
        # Apply branch filter if provided
        if selected_branches:
            selected_branches_upper = [str(b).strip().upper() for b in selected_branches]
            filtered_sales_df = filtered_sales_df[
                filtered_sales_df[sales_area_col].astype(str).str.strip().str.upper().isin(selected_branches_upper)
            ]
            budget_df = budget_df[
                budget_df[budget_area_col].astype(str).str.strip().str.upper().isin(selected_branches_upper)
            ]
        
        # Filter by selected executives
        selected_executives_upper = [str(exec).strip().upper() for exec in selected_executives]
        filtered_sales_df = filtered_sales_df[
            filtered_sales_df[sales_exec_col].astype(str).str.strip().str.upper().isin(selected_executives_upper)
        ]
        budget_df = budget_df[
            budget_df[budget_exec_col].astype(str).str.strip().str.upper().isin(selected_executives_upper)
        ]
        
        # Standardize data for mapping
        budget_df['exec_std'] = budget_df[budget_exec_col].astype(str).str.strip().str.upper()
        budget_df['sl_code_std'] = budget_df[budget_sl_code_col].astype(str).str.strip()
        budget_df['product_std'] = budget_df[budget_product_group_col].astype(str).str.strip().str.upper()
        
        filtered_sales_df['exec_std'] = filtered_sales_df[sales_exec_col].astype(str).str.strip().str.upper()
        filtered_sales_df['sl_code_std'] = filtered_sales_df[sales_sl_code_col].astype(str).str.strip()
        filtered_sales_df['product_std'] = filtered_sales_df[sales_product_group_col].astype(str).str.strip().str.upper()
        
        # Convert numeric columns
        budget_df[budget_qty_col] = pd.to_numeric(budget_df[budget_qty_col], errors='coerce').fillna(0)
        budget_df[budget_value_col] = pd.to_numeric(budget_df[budget_value_col], errors='coerce').fillna(0)
        filtered_sales_df[sales_qty_col] = pd.to_numeric(filtered_sales_df[sales_qty_col], errors='coerce').fillna(0)
        filtered_sales_df[sales_value_col] = pd.to_numeric(filtered_sales_df[sales_value_col], errors='coerce').fillna(0)
        
        # Create detailed mapping
        detailed_records = []
        
        # Group budget data by executive, sl_code, and product_group
        budget_grouped = budget_df.groupby(['exec_std', 'sl_code_std', 'product_std']).agg({
            budget_exec_col: 'first',
            budget_exec_code_col: 'first',
            budget_area_col: 'first',
            budget_sl_code_col: 'first',
            budget_product_group_col: 'first',
            budget_qty_col: 'sum',
            budget_value_col: 'sum'
        }).reset_index()
        
        for _, budget_row in budget_grouped.iterrows():
            exec_name = budget_row['exec_std']
            sl_code = budget_row['sl_code_std']
            product = budget_row['product_std']
            budget_qty = budget_row[budget_qty_col]
            budget_value = budget_row[budget_value_col]
            
            # Check if budget conditions are met (qty > 0 AND value > 0)
            budget_conditions_met = budget_qty > 0 and budget_value > 0
            
            if budget_conditions_met:
                # Find matching sales records
                matching_sales = filtered_sales_df[
                    (filtered_sales_df['exec_std'] == exec_name) &
                    (filtered_sales_df['sl_code_std'] == sl_code) &
                    (filtered_sales_df['product_std'] == product) &
                    filtered_sales_df[sales_qty_col].notna() &
                    filtered_sales_df[sales_value_col].notna()
                ]
                
                # Calculate total sales for this combination
                sales_qty_total = matching_sales[sales_qty_col].sum() if not matching_sales.empty else 0
                sales_value_total = matching_sales[sales_value_col].sum() if not matching_sales.empty else 0
                
                # Apply min logic (same as calculate_budget_values)
                final_sales_qty = min(budget_qty, sales_qty_total) if sales_qty_total > 0 else 0
                final_sales_value = min(budget_value, sales_value_total) if sales_value_total > 0 else 0
                
                # Determine match status
                match_status = 'Mapped' if sales_qty_total > 0 or sales_value_total > 0 else 'No Sales Data'
            else:
                # Budget conditions not met, set sales to 0
                final_sales_qty = 0
                final_sales_value = 0
                match_status = 'Budget Invalid'
            
            # Create record for detailed view
            detailed_record = {
                'Executive Name': budget_row[budget_exec_col],
                'Executive Code': budget_row[budget_exec_code_col],
                'Branch': budget_row[budget_area_col],
                'SL Code': budget_row[budget_sl_code_col],
                'Product Group': budget_row[budget_product_group_col],
                'Budget Qty': round(budget_qty, 2),
                'Budget Value': round(budget_value, 2),
                'Sales Qty': round(final_sales_qty, 2),
                'Sales Value': round(final_sales_value, 2),
                'Match Status': match_status
            }
            
            detailed_records.append(detailed_record)
        
        # Create DataFrame
        proof_df = pd.DataFrame(detailed_records)
        
        # Sort by Executive Name, then by SL Code
        proof_df = proof_df.sort_values(['Executive Name', 'SL Code'])
        
        # Create Excel file
        excel_buffer = BytesIO()
        with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
            # Write the detailed data
            proof_df.to_excel(writer, sheet_name='Proof of Calculation', index=False)
            
            # Get the workbook and worksheet to apply formatting
            workbook = writer.book
            worksheet = writer.sheets['Proof of Calculation']
            
            # Apply formatting
            from openpyxl.styles import Font, PatternFill, Alignment
            
            # Header formatting
            header_font = Font(bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="0070C0", end_color="0070C0", fill_type="solid")
            
            for cell in worksheet[1]:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center')
            
            # Auto-adjust column widths
            for column in worksheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                worksheet.column_dimensions[column_letter].width = adjusted_width
        
        excel_buffer.seek(0)
        return excel_buffer
        
    except Exception as e:
        logger.error(f"Error creating proof of calculation Excel: {e}")
        st.error(f"Error creating proof of calculation Excel: {e}")
        return None
def calculate_budget_values(sales_df, budget_df, selected_month, sales_executives,
                           sales_date_col, sales_area_col, sales_value_col, sales_qty_col,
                           sales_product_group_col, sales_sl_code_col, sales_exec_col,
                           budget_area_col, budget_value_col, budget_qty_col,
                           budget_product_group_col, budget_sl_code_col, budget_exec_col,
                           selected_branches=None):
    """
    Calculates budget vs. billed and overall sales quantities and values.
    """
    try:
        sales_df = sales_df.copy()
        budget_df = budget_df.copy()

        # Log input data shapes (keep logging for debugging)
        logger.info(f"Sales DataFrame shape: {sales_df.shape}")
        logger.info(f"Budget DataFrame shape: {budget_df.shape}")

        # Validate input DataFrames
        if sales_df.empty:
            error_msg = "Sales DataFrame is empty. Please upload valid sales data."
            st.error(error_msg)
            logger.error(error_msg)
            return None, None, None, None
        if budget_df.empty:
            logger.warning("Budget DataFrame is empty. Proceeding with overall sales data only.")

        # Validate column existence
        required_sales_cols = [sales_date_col, sales_value_col, sales_qty_col, sales_exec_col,
                              sales_product_group_col, sales_sl_code_col, sales_area_col]
        required_budget_cols = [budget_value_col, budget_qty_col, budget_exec_col,
                               budget_product_group_col, budget_sl_code_col, budget_area_col]
        
        missing_sales_cols = [col for col in required_sales_cols if col not in sales_df.columns]
        missing_budget_cols = [col for col in required_budget_cols if col not in budget_df.columns and not budget_df.empty]
        
        if missing_sales_cols:
            error_msg = f"Missing columns in sales data: {', '.join(missing_sales_cols)}. Please check column mappings."
            st.error(error_msg)
            logger.error(error_msg)
            return None, None, None, None
        if missing_budget_cols:
            logger.warning(f"Missing columns in budget data: {', '.join(missing_budget_cols)}. Budget vs. billed reports will not be generated.")

        # Clean and convert sales DataFrame columns
        # Date column
        sales_df[sales_date_col] = pd.to_datetime(sales_df[sales_date_col], dayfirst=True, errors='coerce')
        invalid_dates = sales_df[sales_date_col].isna()
        if invalid_dates.any():
            logger.warning(f"Invalid dates in sales column '{sales_date_col}' at rows: {sales_df[invalid_dates].index.tolist()}")

        # Numeric columns - Clean data silently
        for col in [sales_value_col, sales_qty_col]:
            sales_df[col] = sales_df[col].astype(str).str.replace(r'[^\d.]', '', regex=True)
            sales_df[col] = pd.to_numeric(sales_df[col], errors='coerce').fillna(0)

        # String columns (standardize to uppercase)
        for col in [sales_area_col, sales_product_group_col, sales_exec_col]:
            sales_df[col] = sales_df[col].astype(str).str.strip().str.upper()

        # SL code (attempt numeric conversion, preserve invalid as strings)
        sales_df[sales_sl_code_col] = sales_df[sales_sl_code_col].astype(str).str.strip().str.replace(r'\.0$', '', regex=True)
        sales_df['temp_sl_code_numeric'] = pd.to_numeric(sales_df[sales_sl_code_col], errors='coerce')
        sales_df[sales_sl_code_col] = sales_df[sales_sl_code_col].where(sales_df['temp_sl_code_numeric'].isna(), sales_df['temp_sl_code_numeric'].astype(str))
        sales_df = sales_df.drop(columns=['temp_sl_code_numeric'])

        # Clean and convert budget DataFrame columns
        if not budget_df.empty:
            # Numeric columns - Clean data silently
            for col in [budget_value_col, budget_qty_col]:
                budget_df[col] = budget_df[col].astype(str).str.replace(r'[^\d.]', '', regex=True)
                budget_df[col] = pd.to_numeric(budget_df[col], errors='coerce').fillna(0)

            # String columns (standardize to uppercase)
            for col in [budget_area_col, budget_product_group_col, budget_exec_col]:
                budget_df[col] = budget_df[col].astype(str).str.strip().str.upper()

            # SL code (attempt numeric conversion, preserve invalid as strings)
            budget_df[budget_sl_code_col] = budget_df[budget_sl_code_col].astype(str).str.strip().str.replace(r'\.0$', '', regex=True)
            budget_df['temp_sl_code_numeric'] = pd.to_numeric(budget_df[budget_sl_code_col], errors='coerce')
            budget_df[budget_sl_code_col] = budget_df[budget_sl_code_col].where(budget_df['temp_sl_code_numeric'].isna(), budget_df['temp_sl_code_numeric'].astype(str))
            budget_df = budget_df.drop(columns=['temp_sl_code_numeric'])

        # Filter sales data for the selected month
        filtered_sales_df = sales_df[sales_df[sales_date_col].dt.strftime('%b %y') == selected_month].copy()
        if filtered_sales_df.empty:
            error_msg = f"No sales data found for {selected_month}. Please verify the selected month and date column."
            st.error(error_msg)
            logger.error(error_msg)
            return None, None, None, None

        # Apply branch filter if provided
        if selected_branches:
            selected_branches_upper = [str(b).strip().upper() for b in selected_branches]
            filtered_sales_df = filtered_sales_df[filtered_sales_df[sales_area_col].isin(selected_branches_upper)]
            if not budget_df.empty:
                budget_df = budget_df[budget_df[budget_area_col].isin(selected_branches_upper)]
            if filtered_sales_df.empty:
                error_msg = f"No sales data found for selected branches: {', '.join(selected_branches)}. Please verify branch names."
                st.error(error_msg)
                logger.error(error_msg)
                return None, None, None, None

        # Determine executives to display
        if selected_branches:
            branch_sales_df = filtered_sales_df[filtered_sales_df[sales_area_col].isin(selected_branches_upper)]
            branch_budget_df = budget_df[budget_df[budget_area_col].isin(selected_branches_upper)] if not budget_df.empty else pd.DataFrame()
            branch_executives = sorted(set(branch_sales_df[sales_exec_col].dropna().unique()) | 
                                     set(branch_budget_df[budget_exec_col].dropna().unique()))
            
            if sales_executives:
                selected_execs_upper = [str(exec).strip().upper() for exec in sales_executives]
                executives_to_display = [exec for exec in branch_executives if exec in selected_execs_upper]
            else:
                executives_to_display = branch_executives
        else:
            executives_to_display = [str(exec).strip().upper() for exec in sales_executives] if sales_executives else \
                                    sorted(set(filtered_sales_df[sales_exec_col].dropna().unique()) | 
                                          set(budget_df[budget_exec_col].dropna().unique() if not budget_df.empty else []))

        if not executives_to_display:
            error_msg = "No executives found to display. Please check executive selections or data."
            st.error(error_msg)
            logger.error(error_msg)
            return None, None, None, None

        # Filter by selected executives
        filtered_sales_df = filtered_sales_df[filtered_sales_df[sales_exec_col].isin(executives_to_display)].copy()
        if filtered_sales_df.empty:
            error_msg = f"No sales data found for selected executives: {', '.join(executives_to_display)}. Please verify executive names."
            st.error(error_msg)
            logger.error(error_msg)
            return None, None, None, None

        budget_filtered = budget_df[budget_df[budget_exec_col].isin(executives_to_display)].copy() if not budget_df.empty else pd.DataFrame()

        # Log filtered data shapes
        logger.info(f"Filtered Sales DataFrame shape: {filtered_sales_df.shape}")
        logger.info(f"Filtered Budget DataFrame shape: {budget_filtered.shape}")

        # Compute Overall Sales DataFrames
        # Exclude rows with NaN in numeric columns for aggregation
        overall_sales_data = filtered_sales_df[
            filtered_sales_df[sales_qty_col].notna() & 
            filtered_sales_df[sales_value_col].notna()
        ].groupby(sales_exec_col).agg({
            sales_qty_col: 'sum',
            sales_value_col: 'sum'
        }).reset_index()
        overall_sales_data.columns = ['EXECUTIVE', 'Overall_Sales_Qty', 'Overall_Sales_Value']

        overall_sales_qty_df = pd.DataFrame({'EXECUTIVE': executives_to_display})
        overall_sales_value_df = pd.DataFrame({'EXECUTIVE': executives_to_display})

        # Merge budget totals if available
        if not budget_filtered.empty:
            budget_totals = budget_filtered[
                budget_filtered[budget_qty_col].notna() & 
                budget_filtered[budget_value_col].notna()
            ].groupby(budget_exec_col).agg({
                budget_qty_col: 'sum',
                budget_value_col: 'sum'
            }).reset_index()
            budget_totals.columns = ['EXECUTIVE', 'Budget_Qty', 'Budget_Value']
        else:
            budget_totals = pd.DataFrame({'EXECUTIVE': executives_to_display, 'Budget_Qty': [0] * len(executives_to_display), 
                                        'Budget_Value': [0] * len(executives_to_display)})

        # UPDATED: Changed BUDGET to TARGET in column names
        overall_sales_qty_df = pd.merge(
            overall_sales_qty_df,
            budget_totals[['EXECUTIVE', 'Budget_Qty']].rename(columns={'Budget_Qty': 'TARGET QTY/MT'}),
            on='EXECUTIVE',
            how='left'
        ).fillna({'TARGET QTY/MT': 0})

        overall_sales_value_df = pd.merge(
            overall_sales_value_df,
            budget_totals[['EXECUTIVE', 'Budget_Value']].rename(columns={'Budget_Value': 'TARGET VALUE/L'}),
            on='EXECUTIVE',
            how='left'
        ).fillna({'TARGET VALUE/L': 0})

        # Merge sales data
        overall_sales_qty_df = pd.merge(
            overall_sales_qty_df,
            overall_sales_data[['EXECUTIVE', 'Overall_Sales_Qty']].rename(columns={'Overall_Sales_Qty': 'BILLED QTY/MT'}),
            on='EXECUTIVE',
            how='left'
        ).fillna({'BILLED QTY/MT': 0})

        overall_sales_value_df = pd.merge(
            overall_sales_value_df,
            overall_sales_data[['EXECUTIVE', 'Overall_Sales_Value']].rename(columns={'Overall_Sales_Value': 'BILLED VALUE/L'}),
            on='EXECUTIVE',
            how='left'
        ).fillna({'BILLED VALUE/L': 0})

        # Initialize budget vs. billed DataFrames
        budget_vs_billed_qty_df = None
        budget_vs_billed_value_df = None

        # Process Budget vs. Billed Data if budget data exists
        if not budget_filtered.empty:
            # Process Budget Data (exclude rows with NaN in numeric columns)
            budget_grouped = budget_filtered[
                budget_filtered[budget_qty_col].notna() & 
                budget_filtered[budget_value_col].notna() &
                budget_filtered[budget_sl_code_col].notna() &
                budget_filtered[budget_product_group_col].notna()
            ].groupby([
                budget_exec_col, 
                budget_sl_code_col, 
                budget_product_group_col
            ]).agg({
                budget_qty_col: 'sum',
                budget_value_col: 'sum'
            }).reset_index()
            
            budget_valid = budget_grouped[
                (budget_grouped[budget_qty_col] > 0) & 
                (budget_grouped[budget_value_col] > 0)
            ].copy()
            
            if budget_valid.empty:
                logger.warning("No valid budget data found (with qty > 0 and value > 0). Budget vs. billed reports will not be generated.")
            else:
                # Process Sales Data (exclude rows with NaN in numeric columns or SL code)
                final_results = []
                
                for _, budget_row in budget_valid.iterrows():
                    executive = budget_row[budget_exec_col]
                    sl_code = budget_row[budget_sl_code_col]
                    product = budget_row[budget_product_group_col]
                    budget_qty = budget_row[budget_qty_col]
                    budget_value = budget_row[budget_value_col]
                    
                    matching_sales = filtered_sales_df[
                        (filtered_sales_df[sales_exec_col] == executive) &
                        (filtered_sales_df[sales_sl_code_col] == sl_code) &
                        (filtered_sales_df[sales_product_group_col] == product) &
                        filtered_sales_df[sales_qty_col].notna() &
                        filtered_sales_df[sales_value_col].notna()
                    ]
                    
                    sales_qty_total = matching_sales[sales_qty_col].sum() if not matching_sales.empty else 0
                    sales_value_total = matching_sales[sales_value_col].sum() if not matching_sales.empty else 0
                    
                    final_qty = min(budget_qty, sales_qty_total) if pd.notna(budget_qty) and pd.notna(sales_qty_total) else 0
                    final_value = min(budget_value, sales_value_total) if pd.notna(budget_value) and pd.notna(sales_value_total) else 0
                    
                    final_results.append({
                        'EXECUTIVE': executive,
                        'SL_Code': sl_code,
                        'Product': product,
                        'Budget_Qty': budget_qty,
                        'Sales_Qty': sales_qty_total,
                        'Final_Qty': final_qty,
                        'Budget_Value': budget_value,
                        'Sales_Value': sales_value_total,
                        'Final_Value': final_value
                    })
                
                results_df = pd.DataFrame(final_results)
                
                # Aggregate by Executive - UPDATED: Changed column names
                exec_qty_summary = results_df.groupby('EXECUTIVE').agg({
                    'Budget_Qty': 'sum',
                    'Final_Qty': 'sum'
                }).reset_index()
                exec_qty_summary.columns = ['EXECUTIVE', 'TARGET QTY/MT', 'BILLED QTY/MT']
                
                exec_value_summary = results_df.groupby('EXECUTIVE').agg({
                    'Budget_Value': 'sum',
                    'Final_Value': 'sum'
                }).reset_index()
                exec_value_summary.columns = ['EXECUTIVE', 'TARGET VALUE/L', 'BILLED VALUE/L']
                
                # Build QUANTITY DataFrame
                qty_data = []
                for exec_name in executives_to_display:
                    exec_qty_row = exec_qty_summary[exec_qty_summary['EXECUTIVE'] == exec_name]
                    budget_val = round(float(exec_qty_row['TARGET QTY/MT'].iloc[0]), 2) if not exec_qty_row.empty else 0.00
                    billed_val = round(float(exec_qty_row['BILLED QTY/MT'].iloc[0]), 2) if not exec_qty_row.empty else 0.00
                    percentage = round((billed_val / budget_val * 100), 2) if budget_val > 0 else 0.00
                    
                    qty_data.append({
                        'EXECUTIVE': exec_name,
                        'TARGET QTY/MT': budget_val,
                        'BILLED QTY/MT': billed_val,
                        '%': percentage
                    })
                
                budget_vs_billed_qty_df = pd.DataFrame(qty_data)
                
                # Build VALUE DataFrame
                value_data = []
                for exec_name in executives_to_display:
                    exec_value_row = exec_value_summary[exec_value_summary['EXECUTIVE'] == exec_name]
                    budget_val = round(float(exec_value_row['TARGET VALUE/L'].iloc[0]), 2) if not exec_value_row.empty else 0.00
                    billed_val = round(float(exec_value_row['BILLED VALUE/L'].iloc[0]), 2) if not exec_value_row.empty else 0.00
                    percentage = round((billed_val / budget_val * 100), 2) if budget_val > 0 else 0.00
                    
                    value_data.append({
                        'EXECUTIVE': exec_name,
                        'TARGET VALUE/L': budget_val,
                        'BILLED VALUE/L': billed_val,
                        '%': percentage
                    })
                
                budget_vs_billed_value_df = pd.DataFrame(value_data)
                
                # Add total rows with consistent 2-decimal formatting
                total_budget_qty = float(f"{budget_vs_billed_qty_df['TARGET QTY/MT'].sum():.2f}")
                total_billed_qty = float(f"{budget_vs_billed_qty_df['BILLED QTY/MT'].sum():.2f}")
                total_percentage_qty = float(f"{(total_billed_qty / total_budget_qty * 100):.2f}") if total_budget_qty > 0 else 0.00
                
                total_row_qty = pd.DataFrame({
                    'EXECUTIVE': ['TOTAL'],
                    'TARGET QTY/MT': [total_budget_qty],
                    'BILLED QTY/MT': [total_billed_qty],
                    '%': [total_percentage_qty]
                })
                budget_vs_billed_qty_df = pd.concat([budget_vs_billed_qty_df, total_row_qty], ignore_index=True)
                
                total_budget_value = float(f"{budget_vs_billed_value_df['TARGET VALUE/L'].sum():.2f}")
                total_billed_value = float(f"{budget_vs_billed_value_df['BILLED VALUE/L'].sum():.2f}")
                total_percentage_value = float(f"{(total_billed_value / total_budget_value * 100):.2f}") if total_budget_value > 0 else 0.00
                
                total_row_value = pd.DataFrame({
                    'EXECUTIVE': ['TOTAL'],
                    'TARGET VALUE/L': [total_budget_value],
                    'BILLED VALUE/L': [total_billed_value],
                    '%': [total_percentage_value]
                })
                budget_vs_billed_value_df = pd.concat([budget_vs_billed_value_df, total_row_value], ignore_index=True)

        # Add total rows for overall sales with consistent formatting
        total_row_overall_qty = pd.DataFrame({
            'EXECUTIVE': ['TOTAL'],
            'TARGET QTY/MT': [float(f"{overall_sales_qty_df['TARGET QTY/MT'].sum():.2f}")],
            'BILLED QTY/MT': [float(f"{overall_sales_qty_df['BILLED QTY/MT'].sum():.2f}")]
        })
        overall_sales_qty_df = pd.concat([overall_sales_qty_df, total_row_overall_qty], ignore_index=True)

        total_row_overall_value = pd.DataFrame({
            'EXECUTIVE': ['TOTAL'],
            'TARGET VALUE/L': [float(f"{overall_sales_value_df['TARGET VALUE/L'].sum():.2f}")],
            'BILLED VALUE/L': [float(f"{overall_sales_value_df['BILLED VALUE/L'].sum():.2f}")]
        })
        overall_sales_value_df = pd.concat([overall_sales_value_df, total_row_overall_value], ignore_index=True)

        # Apply consistent 2-decimal formatting to all numeric columns
        for df in [budget_vs_billed_qty_df, budget_vs_billed_value_df, overall_sales_qty_df, overall_sales_value_df]:
            if df is not None:
                for col in df.columns:
                    if col != 'EXECUTIVE' and '%' not in str(col):
                        try:
                            df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)
                            df[col] = df[col].apply(lambda x: float(f"{x:.2f}") if pd.notna(x) else 0.00)
                        except:
                            pass

        return (budget_vs_billed_value_df, budget_vs_billed_qty_df, overall_sales_qty_df, overall_sales_value_df)

    except Exception as e:
        error_msg = f"Unexpected error in calculate_budget_values: {str(e)}. Please check your data and selections."
        st.error(error_msg)
        logger.error(error_msg, exc_info=True)
        # Return overall sales DataFrames if they were computed
        overall_sales_qty_df = locals().get('overall_sales_qty_df', None)
        overall_sales_value_df = locals().get('overall_sales_value_df', None)
        return None, None, overall_sales_qty_df, overall_sales_value_df

def create_budget_ppt(budget_vs_billed_value_df, budget_vs_billed_qty_df, overall_sales_qty_df, overall_sales_value_df, month_title=None, logo_file=None):
   try:
       prs = Presentation()
       prs.slide_width = Inches(13.33)
       prs.slide_height = Inches(7.5)
       
       # Create title slide (using original function)
       create_title_slide(prs, f"Monthly Review Meeting – {month_title}", logo_file)
       
       def process_df_for_slides(df, title_base, percent_cols=None, data_type=""):
           if percent_cols is None:
               percent_cols = []
           if df is None or df.empty:
               logger.warning(f"Skipping slide for {title_base}: DataFrame is None or empty")
               st.warning(f"Cannot generate slide for {title_base}: No data available")
               return
           
           df = df[df['EXECUTIVE'] != "ACCLP"].copy()
           num_executives = df[df['EXECUTIVE'] != 'TOTAL'].shape[0]
           
           # Increased split threshold from 12 to 15 to utilize bottom space
           split_threshold = 15
           
           if num_executives <= split_threshold:
               # UPDATED: Include month and data type in title
               slide_title = f"{title_base} - {data_type.upper()} - {month_title}"
               add_table_slide(prs, df, slide_title, percent_cols=percent_cols)
               return
           
           # Split logic for tables with more than 15 data rows
           data_rows = df[df['EXECUTIVE'] != 'TOTAL'].copy()
           total_row = df[df['EXECUTIVE'] == 'TOTAL'].copy()
           
           # Split into chunks of 15 rows each (no part totals)
           chunks = []
           for i in range(0, len(data_rows), split_threshold):
               chunks.append(data_rows.iloc[i:i+split_threshold].copy())
           
           # Process each chunk - add total row only to the last part
           for i, chunk in enumerate(chunks, 1):
               # Round numeric columns
               numeric_cols = chunk.select_dtypes(include=[np.number]).columns
               for col in numeric_cols:
                   chunk[col] = chunk[col].apply(lambda x: round(float(x), 2))
               
               # If this is the last chunk, add the total row at the bottom
               if i == len(chunks):
                   # Round total row numeric columns
                   total_numeric_cols = total_row.select_dtypes(include=[np.number]).columns
                   for col in total_numeric_cols:
                       total_row[col] = total_row[col].apply(lambda x: round(float(x), 2))
                   
                   # Combine last chunk with total row
                   chunk_with_total = pd.concat([chunk, total_row], ignore_index=True)
                   final_data = chunk_with_total
               else:
                   final_data = chunk
               
               # UPDATED: Include month and data type in title for all slides
               slide_title = f"{title_base} - {data_type.upper()} - {month_title}"
               add_table_slide(prs, final_data, slide_title, percent_cols=percent_cols)
       
       # Process each DataFrame with clean titles and appropriate data types
       dfs_to_process = [
           (budget_vs_billed_qty_df, "TARGET AGAINST BILLED", [3], "QUANTITY"),
           (budget_vs_billed_value_df, "TARGET AGAINST BILLED", [3], "VALUE"), 
           (overall_sales_qty_df, "OVERALL SALES", [], "QUANTITY"),
           (overall_sales_value_df, "OVERALL SALES", [], "VALUE")
       ]
       
       for df, title, percent_cols, data_type in dfs_to_process:
           process_df_for_slides(df, title, percent_cols, data_type)
       
       # Check if any slides were added (excluding title slide)
       if len(prs.slides) <= 1:
           st.error("No data slides generated. All DataFrames are empty or None. Please check your data.")
           logger.error("No data slides generated in Budget PPT")
           return None
       
       # Add thank you slide at the end
       create_thank_you_slide(prs, logo_file)
       
       ppt_buffer = BytesIO()
       prs.save(ppt_buffer)
       ppt_buffer.seek(0)
       return ppt_buffer
       
   except Exception as e:
       logger.error(f"Error creating Budget PPT: {e}")
       st.error(f"Error creating Budget PPT: {e}")
       return None

def determine_financial_year(date):
   year = date.year
   month = date.month
   if month >= 4:
       return f"{year % 100}-{year % 100 + 1}"
   else:
       return f"{year % 100 - 1}-{year % 100}"

def create_customer_table(sales_df, date_col, branch_col, customer_id_col, executive_col, selected_months=None, selected_branches=None, selected_executives=None):
    """
    Creates a table of unique customer counts per executive for selected months, branches, and executives.
    
    Parameters:
    - sales_df: DataFrame containing sales data.
    - date_col: Column name for dates.
    - branch_col: Column name for branch names.
    - customer_id_col: Column name for customer IDs (SL Code).
    - executive_col: Column name for executive names.
    - selected_months: List of months to filter (format: 'MMM YYYY').
    - selected_branches: List of branches to filter.
    - selected_executives: List of executives to filter.
    
    Returns:
    - Dictionary with formatted month period as key and (DataFrame, sorted_months) as value.
    """
    sales_df = sales_df.copy()
    
    # Validate columns
    for col in [date_col, branch_col, customer_id_col, executive_col]:
        if col not in sales_df.columns:
            st.error(f"Column '{col}' not found in sales data.")
            return None
    
    # Convert date column to datetime
    try:
        sales_df[date_col] = pd.to_datetime(sales_df[date_col], errors='coerce', dayfirst=True, format='mixed')
    except Exception as e:
        st.error(f"Error converting '{date_col}' to datetime: {e}. Ensure dates are in a valid format.")
        return None
    
    # Check for valid dates
    valid_dates = sales_df[date_col].notna()
    if not valid_dates.any():
        st.error(f"Column '{date_col}' contains no valid dates.")
        return None
    
    # Extract month-year for filtering
    sales_df['Month_Year'] = sales_df[date_col].dt.strftime('%b %Y')
    
    # Filter by selected months if provided
    if selected_months:
        sales_df = sales_df[sales_df['Month_Year'].isin(selected_months)]
        if sales_df.empty:
            st.error(f"No data found for selected months: {', '.join(selected_months)}")
            return None

    # Determine financial year for internal processing but use month for title
    def determine_financial_year(date):
        if pd.isna(date):
            return None
        year = date.year
        month = date.month
        if month >= 4:
            return f"{year}-{year+1}"
        else:
            return f"{year-1}-{year}"
    
    sales_df['Financial_Year'] = sales_df[date_col].apply(determine_financial_year)
    available_financial_years = sales_df['Financial_Year'].dropna().unique()
    if len(available_financial_years) == 0:
        st.error("No valid financial years found in the data.")
        return None
    
    result_dict = {}
    for fin_year in sorted(available_financial_years):
        fy_df = sales_df[sales_df['Financial_Year'] == fin_year].copy()
        if fy_df.empty:
            continue
        
        # Extract unique months in chronological order
        fy_df['Month_Year_Period'] = fy_df[date_col].dt.to_period('M')
        available_months = fy_df['Month_Year_Period'].unique()
        month_names = [pd.to_datetime(str(m) + '-01').strftime('%b %Y') for m in sorted(available_months)]
        
        if not month_names:
            continue
        
        # Filter by selected months if provided
        if selected_months:
            month_names = [m for m in month_names if m in selected_months]
            if not month_names:
                continue
        
        # Standardize branch and executive names
        try:
            fy_df['Branch'] = fy_df[branch_col].astype(str).str.strip().str.upper()
            fy_df['EXECUTIVE'] = fy_df[executive_col].astype(str).str.strip().str.upper()
        except Exception as e:
            st.error(f"Error processing columns: {e}")
            continue
        
        # Apply branch filter if provided
        if selected_branches:
            fy_df = fy_df[fy_df['Branch'].isin([b.upper() for b in selected_branches])]
            if fy_df.empty:
                continue
        
        # Determine executives to display based on both branch and executive selections
        if selected_branches:
            branch_df = fy_df[fy_df['Branch'].isin([b.upper() for b in selected_branches])]
            branch_executives = sorted(branch_df['EXECUTIVE'].dropna().unique())
            
            if selected_executives:
                selected_execs_upper = [str(e).upper() for e in selected_executives]
                executives_to_display = [exec for exec in branch_executives if exec in selected_execs_upper]
            else:
                executives_to_display = branch_executives
        else:
            executives_to_display = [str(e).upper() for e in selected_executives] if selected_executives else sorted(fy_df['EXECUTIVE'].dropna().unique())
        
        # Apply executive filter
        if executives_to_display:
            fy_df = fy_df[fy_df['EXECUTIVE'].isin(executives_to_display)]
            if fy_df.empty:
                continue
        
        if not executives_to_display:
            continue
        
        # Group by executive and month to count unique customer codes
        grouped_df = fy_df.groupby(['EXECUTIVE', 'Month_Year'])[customer_id_col].nunique().reset_index(name='Customer_Count')
        
        # Pivot to create table with months as columns
        pivot_df = grouped_df.pivot_table(
            values='Customer_Count',
            index='EXECUTIVE',
            columns='Month_Year',
            aggfunc='sum',
            fill_value=0
        ).reset_index()
        
        # Rename index column to uppercase
        pivot_df = pivot_df.rename(columns={'EXECUTIVE': 'EXECUTIVE NAME'})
        
        # Create result dataframe with all executives to display
        result_df = pd.DataFrame({'EXECUTIVE NAME': executives_to_display})
        result_df = pd.merge(
            result_df,
            pivot_df,
            on='EXECUTIVE NAME',
            how='left'
        ).fillna(0)
        
        # Keep only selected months
        columns_to_keep = ['EXECUTIVE NAME'] + month_names
        result_df = result_df[[col for col in columns_to_keep if col in result_df.columns]]
        
        # Convert counts to integers
        for col in result_df.columns[1:]:
            result_df[col] = result_df[col].astype(int)
        
        # Add total row
        total_row = {'EXECUTIVE NAME': 'GRAND TOTAL'}
        for col in month_names:
            if col in result_df.columns:
                total_row[col] = result_df[col].sum()
        
        result_df = pd.concat([result_df, pd.DataFrame([total_row])], ignore_index=True)
        
        # Format month period for title - CHANGED FROM FY TO MONTH
        if selected_months:
            month_period = format_month_for_title(selected_months)
        else:
            month_period = format_month_for_title(month_names)
        
        result_dict[month_period] = (result_df, month_names)
    
    return result_dict

def format_month_for_title(month_list):
    """Convert month list to title format like JUL-2025"""
    if not month_list:
        return "UNKNOWN"
    
    # If multiple months, use first and last
    if len(month_list) > 1:
        first_month = month_list[0]
        last_month = month_list[-1]
        # Convert "Jul 2025" to "JUL-2025" 
        first_formatted = first_month.replace(' ', '-').upper()
        last_formatted = last_month.replace(' ', '-').upper()
        return f"{first_formatted} TO {last_formatted}"
    else:
        # Single month: "Jul 2025" -> "JUL-2025"
        return month_list[0].replace(' ', '-').upper()
   

def create_customer_table_image(df, title, sorted_months, month_period):
    fig, ax = plt.subplots(figsize=(14, len(df) * 0.6))
    ax.axis('off')
    columns = list(df.columns)
    expected_columns = {'EXECUTIVE NAME'}.union(set(sorted_months))
    actual_columns = set(columns)
    if not {'EXECUTIVE NAME'}.issubset(actual_columns):
        st.warning(f"Missing essential columns in customer DataFrame for image: EXECUTIVE NAME")
        return BytesIO()
    rows = len(df)
    ncols = len(columns)
    width = 1.0 / ncols
    height = 1.0 / rows
    table = Table(ax, bbox=[0, 0, 1, 1])
    for col_idx, col_name in enumerate(columns):
        table.add_cell(0, col_idx, width, height, text=col_name, loc='center', facecolor='#0070C0')
        table[0, col_idx].set_text_props(weight='bold', color='white', fontsize=10)
    for row_idx in range(rows):
        for col_idx in range(ncols):
            value = df.iloc[row_idx, col_idx]
            text = str(value)
            facecolor = '#DDEBF7' if row_idx % 2 == 0 else 'white'
            if row_idx == rows - 1:
                facecolor = '#D3D3D3'
                table.add_cell(row_idx + 1, col_idx, width, height, text=text, loc='center', facecolor=facecolor).set_text_props(weight='bold', fontsize=10)
            else:
                table.add_cell(row_idx + 1, col_idx, width, height, text=text, loc='center', facecolor=facecolor).set_text_props(fontsize=10)
    if ncols > 0:
        table[(0, 0)].width = 0.30  # Executive name column wider
    if ncols > 1:
        remaining_width = 0.70
        col_width = remaining_width / (ncols - 1)
        for i in range(1, ncols):
            table[(0, i)].width = col_width
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    ax.add_table(table)
    
    # CHANGE THIS LINE - Remove "FY" prefix from title
    clean_title = title.replace("FY ", "")
    plt.suptitle(clean_title, fontsize=14, weight='bold', color='#0070C0', y=1.02)
    
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
    plt.close()
    return img_buffer

def create_customer_ppt_slide(slide, df, title, sorted_months, is_last_slide=False, page_number=None):
    """Updated customer slide creation with page number support"""
    if df.empty or len(df.columns) < 2:
        st.warning(f"Skipping customer slide: DataFrame is empty or has insufficient columns {df.columns.tolist()}")
        return
    
    # Create title with consistent formatting
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title.upper()
    p = title_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    
    columns = list(df.columns)
    if 'EXECUTIVE NAME' not in columns:
        st.warning(f"Missing essential column in customer DataFrame: EXECUTIVE NAME")
        return
    
    # Calculate table dimensions to use full slide space
    num_rows = len(df) + 1  # +1 for header
    num_cols = len(columns)
    
    # Use more of the slide space
    table_width = Inches(12.0)
    table_height = Inches(5.8)
    left = Inches(0.67)
    top = Inches(1.2)
    
    # Calculate row height based on available space - Convert to int
    row_height = int(table_height / num_rows)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
    
    # Calculate column widths to fill the entire table width
    if num_cols == 2:
        exec_width = Inches(4.0)
        month_width = Inches(8.0)
    elif num_cols > 2:
        exec_width = Inches(4.0)
        remaining_width = table_width - exec_width
        month_width = remaining_width / (num_cols - 1)
    else:
        exec_width = Inches(12.0)
        month_width = Inches(1.0)
    
    # Set column widths
    table.columns[0].width = int(exec_width)
    for col_idx in range(1, num_cols):
        table.columns[col_idx].width = int(month_width)
    
    # Set row heights to fill the table
    for row_idx in range(num_rows):
        table.rows[row_idx].height = row_height
    
    # Format header row consistently
    for col_idx, col_name in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name.upper()
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
    
    # Format data rows consistently
    for row_idx, (index, row) in enumerate(df.iterrows(), start=1):
        executive_name = str(row.get('EXECUTIVE NAME', '')).upper()
        is_total_row = ('TOTAL' in executive_name) or ('GRAND TOTAL' in executive_name)
        
        for col_idx, col_name in enumerate(columns):
            cell = table.cell(row_idx, col_idx)
            try:
                value = row[col_name]
                cell.text = str(value).upper() if isinstance(value, str) else str(value)
            except (KeyError, ValueError) as e:
                cell.text = ""
                st.warning(f"Error accessing {col_name} at row {index} in customer slide: {e}")
            
            if is_total_row:
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
            else:
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.fill.solid()
                if (row_idx - 1) % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
    
    # Add page number if provided
    if page_number is not None:
        add_page_number(slide, page_number)

def create_customer_ppt_with_splitting(df, title_base, sorted_months, month_period):
    """
    Create multiple PowerPoint slides with table splitting for better readability
    """
    slides_data = []
    
    if df.empty:
        return slides_data
    
    # Remove GRAND TOTAL row temporarily for splitting
    total_row = None
    if not df.empty and df.iloc[-1].get('EXECUTIVE NAME') == 'GRAND TOTAL':
        total_row = df.iloc[-1:].copy()
        data_rows = df.iloc[:-1].copy()
    else:
        data_rows = df.copy()
    
    # Split threshold
    max_rows_per_slide = 15
    
    # If data fits in one slide
    if len(data_rows) <= max_rows_per_slide:
        # Add total row back if it exists
        if total_row is not None:
            final_df = pd.concat([data_rows, total_row], ignore_index=True)
        else:
            final_df = data_rows
        
        slides_data.append({
            'df': final_df,
            'title': f"NUMBER OF BILLED CUSTOMERS - {month_period}",
            'is_last': True
        })
    else:
        # Split into multiple slides
        num_parts = (len(data_rows) + max_rows_per_slide - 1) // max_rows_per_slide
        
        for i in range(num_parts):
            start_idx = i * max_rows_per_slide
            end_idx = min((i + 1) * max_rows_per_slide, len(data_rows))
            chunk = data_rows.iloc[start_idx:end_idx].copy()
            
            # Add total row only to the last slide
            if i == num_parts - 1 and total_row is not None:
                chunk = pd.concat([chunk, total_row], ignore_index=True)
            
            # Create slide title with slide number
            slide_title = f"NUMBER OF BILLED CUSTOMERS - {month_period})"
            
            slides_data.append({
                'df': chunk,
                'title': slide_title,
                'is_last': (i == num_parts - 1)
            })
    
    return slides_data

def extract_area_name(area):
   if pd.isna(area) or not str(area).strip():
       return None
   area = str(area).strip()
   area_upper = area.upper()
   if area_upper == 'HO' or area_upper.endswith('-HO'):
       return None
   branch_variations = {
       'PUDUCHERRY': ['PUDUCHERRY', 'PONDY', 'PUDUCHERRY - PONDY', 'PONDICHERRY', 'PUDUCHERI', 'aaaa - PUDUCHERRY', 'AAAA - PUDUCHERRY'],
       'COIMBATORE': ['COIMBATORE', 'CBE', 'COIMBATORE - CBE', 'COIMBATURE', 'aaaa - COIMBATORE', 'AAAA - COIMBATORE'],
       'KARUR': ['KARUR', 'KRR', 'KARUR - KRR', 'aaaa - KARUR', 'AAAA - KARUR'],
       'MADURAI': ['MADURAI', 'MDU', 'MADURAI - MDU', 'MADURA', 'aaaa - MADURAI', 'AAAA - MADURAI'],
       'CHENNAI': ['CHENNAI', 'CHN', 'aaaa - CHENNAI', 'AAAA - CHENNAI'],
   }
   for standard_name, variations in branch_variations.items():
       for variation in variations:
           if variation in area_upper:
               return standard_name
   prefixes = ['AAAA - ', 'aaaa - ', 'BBB - ', 'bbb - ', 'ASIA CRYSTAL COMMODITY LLP - ']
   for prefix in prefixes:
       if area_upper.startswith(prefix.upper()):
           return area[len(prefix):].strip().upper()
   separators = [' - ', '-', ':']
   for sep in separators:
       if sep in area_upper:
           return area_upper.split(sep)[-1].strip()
   return area_upper

def filter_os_qty(os_df, os_area_col, os_qty_col, os_due_date_col, os_exec_col, 
                 selected_branches=None, selected_years=None, till_month=None, selected_executives=None):
   required_columns = [os_area_col, os_qty_col, os_due_date_col, os_exec_col]
   for col in required_columns:
       if col not in os_df.columns:
           st.error(f"Column '{col}' not found in OS data.")
           return None, None, None
   
   os_df = os_df.copy()
   os_df[os_area_col] = os_df[os_area_col].apply(extract_area_name).astype(str).str.strip().str.upper()
   os_df[os_exec_col] = os_df[os_exec_col].apply(lambda x: 'BLANK' if pd.isna(x) or str(x).strip() == '' else str(x).strip().upper())
   
   try:
       os_df[os_due_date_col] = pd.to_datetime(os_df[os_due_date_col], errors='coerce')
   except Exception as e:
       st.error(f"Error converting '{os_due_date_col}' to datetime: {e}. Ensure dates are in 'YYYY-MM-DD' format.")
       return None, None, None
   
   # Convert negative values to 0 BEFORE division
   os_df[os_qty_col] = pd.to_numeric(os_df[os_qty_col], errors='coerce').fillna(0)
   os_df[os_qty_col] = os_df[os_qty_col].clip(lower=0)  # Convert negative values to 0
   os_df[os_qty_col] = os_df[os_qty_col] / 100000  # Then divide by 100000
   
   if os_df[os_qty_col].isna().any():
       st.warning(f"Non-numeric values in '{os_qty_col}' replaced with 0.")
   
   start_date, end_date = None, None
   if selected_years and till_month:
       month_map = {
           'January': 1, 'February': 2, 'March': 3, 'April': 4, 'May': 5, 'June': 6,
           'July': 7, 'August': 8, 'September': 9, 'October': 10, 'November': 11, 'December': 12
       }
       if till_month == "November Brodha":
           till_month = "November"
       till_month_num = month_map.get(till_month)
       if not till_month_num:
           st.error(f"Invalid month selected: {till_month}")
           return None, None, None
       selected_years = [int(year) for year in selected_years]
       earliest_year = min(selected_years)
       latest_year = max(selected_years)
       start_date = datetime(earliest_year, 1, 1)
       end_date = (datetime(latest_year, till_month_num, 1) + relativedelta(months=1) - relativedelta(days=1))
       os_df = os_df[
           (os_df[os_due_date_col].notna()) &
           (os_df[os_due_date_col] >= start_date) &
           (os_df[os_due_date_col] <= end_date)
       ]
       if os_df.empty:
           st.error(f"No data matches the period from Jan {earliest_year} to {end_date.strftime('%b %Y')}.")
           return None, None, None
   
   # Apply branch filter if provided
   if selected_branches:
       os_df = os_df[os_df[os_area_col].isin([b.upper() for b in selected_branches])]
       if os_df.empty:
           st.error("No data matches the selected branches.")
           return None, None, None
   
   # Executive filtering logic - only use selected_executives if provided
   if selected_executives:
       # Use only the explicitly selected executives
       selected_execs_upper = [str(e).upper() for e in selected_executives]
       executives_to_display = selected_execs_upper
   else:
       # If no executives selected, use all executives in the filtered data
       executives_to_display = sorted(os_df[os_exec_col].dropna().unique())
   
   # Filter data by selected executives
   if executives_to_display:
       os_df = os_df[os_df[os_exec_col].isin(executives_to_display)]
       if os_df.empty:
           st.error("No data matches the selected executives.")
           return None, None, None
   
   # Filter to only positive values (after negative conversion to 0)
   os_df = os_df[os_df[os_qty_col] > 0]
   if os_df.empty:
       st.error("No positive net values found in the filtered data.")
       return None, None, None
   
   # Group and aggregate data - Updated column name with /L unit
   os_grouped_qty = (os_df.groupby(os_exec_col)
                    .agg({os_qty_col: 'sum'})
                    .reset_index()
                    .rename(columns={os_exec_col: 'EXECUTIVE', os_qty_col: 'TARGET/L'}))

   # Ensure all executives_to_display are included (only selected ones)
   result_df = pd.DataFrame({'EXECUTIVE': executives_to_display})
   result_df = pd.merge(result_df, os_grouped_qty, on='EXECUTIVE', how='left').fillna({'TARGET/L': 0})
   
   # Add total row with consistent 2-decimal formatting
   total_row = pd.DataFrame([{
       'EXECUTIVE': 'TOTAL', 
       'TARGET/L': float(f"{result_df['TARGET/L'].sum():.2f}")
   }])
   result_df = pd.concat([result_df, total_row], ignore_index=True)
   
   # Apply consistent 2-decimal formatting
   result_df['TARGET/L'] = result_df['TARGET/L'].apply(lambda x: float(f"{x:.2f}") if pd.notna(x) else 0.00)
   
   return result_df, start_date, end_date

def create_od_table_image(df, title, columns_to_show=None):
   if columns_to_show is None:
       if 'EXECUTIVE' in df.columns:
           columns_to_show = ['EXECUTIVE', 'TARGET/L']
       else:
           columns_to_show = ['Area', 'TARGET/L']
   fig, ax = plt.subplots(figsize=(10, len(df) * 0.5))
   ax.axis('off')
   nrows, ncols = len(df), len(columns_to_show)
   table = Table(ax, bbox=[0, 0, 1, 1])
   for col_idx, col_name in enumerate(columns_to_show):
       table.add_cell(0, col_idx, 1.0/ncols, 1.0/nrows, text=col_name, loc='center', facecolor='#F2F2F2')
       table[0, col_idx].set_text_props(weight='bold', color='black', fontsize=12)
   column_mapping = {
       'EXECUTIVE': 'EXECUTIVE',
       'Area': 'Area',
       'TARGET/L': 'TARGET/L'
   }
   key_column = 'EXECUTIVE' if 'EXECUTIVE' in df.columns else 'Area'
   for row_idx in range(len(df)):
       for col_idx, display_col_name in enumerate(columns_to_show):
           actual_col_name = column_mapping.get(display_col_name, display_col_name)
           if actual_col_name not in df.columns:
               if display_col_name == 'TARGET/L' and 'TARGET/L' in df.columns:
                   actual_col_name = 'TARGET/L'
               else:
                   st.error(f"Column '{actual_col_name}' not found in DataFrame")
                   continue
           value = df.iloc[row_idx][actual_col_name]
           # Apply consistent 2-decimal formatting for numeric values
           text = str(value).upper() if actual_col_name == key_column else f"{float(value):.2f}"
           facecolor = '#DDEBF7' if row_idx % 2 == 0 else 'white'
           if row_idx == len(df) - 1 and df.iloc[row_idx][key_column] == 'TOTAL':
               facecolor = '#D3D3D3'
               table.add_cell(row_idx + 1, col_idx, 1.0/ncols, 1.0/nrows, 
                             text=text, loc='center', facecolor=facecolor).set_text_props(weight='bold', fontsize=12)
           else:
               table.add_cell(row_idx + 1, col_idx, 1.0/ncols, 1.0/nrows, 
                             text=text, loc='center', facecolor=facecolor).set_text_props(fontsize=10)
   table[(0, 0)].width = 0.6
   table[(0, 1)].width = 0.4
   table.auto_set_font_size(False)
   table.set_fontsize(10)
   ax.add_table(table)
   # Remove "(Value in Lakhs)" from title
   clean_title = title.replace(" (Value in Lakhs)", "").replace(" (Values in Lakhs)", "")
   plt.suptitle(clean_title, fontsize=16, weight='bold', color='black', y=1.05)
   img_buffer = BytesIO()
   plt.savefig(img_buffer, format='png', bbox_inches='tight', dpi=150)
   plt.close()
   return img_buffer

def create_od_ppt_slides(presentation, df, base_title):
    """Create multiple OD slides, splitting data if more than 15 rows per slide"""
    if df.empty or len(df.columns) < 2:
        st.warning(f"Skipping OD slides: DataFrame is empty or has insufficient columns {df.columns.tolist()}")
        return
    
    # Configuration
    max_rows_per_slide = 15
    total_rows = len(df)
    
    # Calculate number of slides needed
    num_slides = (total_rows + max_rows_per_slide - 1) // max_rows_per_slide
    
    for slide_num in range(num_slides):
        # Calculate row range for this slide
        start_idx = slide_num * max_rows_per_slide
        end_idx = min(start_idx + max_rows_per_slide, total_rows)
        
        # Get subset of data for this slide
        slide_df = df.iloc[start_idx:end_idx].copy()
        
        # Create slide title with part number if multiple slides
        if num_slides > 1:
            slide_title = f"{base_title}"
        else:
            slide_title = base_title
        
        # Add new slide
        slide_layout = presentation.slide_layouts[5]  # Blank layout
        slide = presentation.slides.add_slide(slide_layout)
        
        # Create the slide content
        create_single_od_slide(slide, slide_df, slide_title)

def create_single_od_slide(slide, df, title):
    """Updated OD slide creation with consistent formatting"""
    if df.empty or len(df.columns) < 2:
        st.warning(f"Skipping OD slide: DataFrame is empty or has insufficient columns {df.columns.tolist()}")
        return
    
    try:
        # Create title with consistent formatting
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        clean_title = title.replace(" (Value in Lakhs)", "").replace(" (Values in Lakhs)", "")
        title_frame.text = clean_title.upper()
        p = title_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.CENTER
        
        # Determine key column
        key_column = 'EXECUTIVE' if 'EXECUTIVE' in df.columns else 'Area'
        
        # Calculate table dimensions
        num_rows = len(df) + 1
        num_cols = 2
        
        table_width = Inches(12.0)
        table_height = Inches(5.8)
        left = Inches(0.67)
        top = Inches(1.2)
        
        row_height = int(table_height / num_rows)
        
        # Create table
        table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
        
        # Set column widths
        executive_width = Inches(6.0)
        target_width = Inches(6.0)
        
        table.columns[0].width = int(executive_width)
        table.columns[1].width = int(target_width)
        
        # Set row heights
        for row_idx in range(num_rows):
            table.rows[row_idx].height = row_height
        
        # Format header row
        headers = [key_column, "TARGET/L"]
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text.upper()
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
        
        # Format data rows with consistent formatting
        for row_idx in range(len(df)):
            executive_name = str(df.iloc[row_idx][key_column]).upper()
            is_total_row = ('TOTAL' in executive_name) or ('GRAND TOTAL' in executive_name)
            
            # Fill executive name
            exec_cell = table.cell(row_idx + 1, 0)
            exec_cell.text = str(df.iloc[row_idx][key_column]).upper()
            exec_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Fill target value with consistent 2-decimal formatting
            target_cell = table.cell(row_idx + 1, 1)
            target_value = df.iloc[row_idx]['TARGET/L']
            value_text = f"{float(target_value):.2f}"
            target_cell.text = value_text
            target_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Apply formatting based on row type
            for col_idx in range(2):
                cell = table.cell(row_idx + 1, col_idx)
                
                if is_total_row:
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(211, 211, 211)
                else:
                    cell.text_frame.paragraphs[0].font.size = Pt(11)
                    cell.fill.solid()
                    if row_idx % 2 == 0:
                        cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                    else:
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)
                
    except Exception as e:
        st.error(f"Error creating OD PPT slide: {e}")
        st.error(traceback.format_exc())

def create_thank_you_slide(prs, logo_file=None):
    """Create a thank you slide at the end of presentation"""
    try:
        blank_slide_layout = prs.slide_layouts[6]
        thank_you_slide = prs.slides.add_slide(blank_slide_layout)
        
        # Logo (if provided)
        if logo_file is not None:
            try:
                logo_buffer = BytesIO(logo_file.read())
                logo = thank_you_slide.shapes.add_picture(logo_buffer, Inches(5.665), Inches(2.0), width=Inches(2), height=Inches(2))
                logo_file.seek(0)
            except Exception as e:
                logger.error(f"Error adding logo to thank you slide: {e}")
        
        # Thank You message
        thank_you_box = thank_you_slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.33), Inches(1.5))
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.text = "THANK YOU"
        p = thank_you_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        

        
        # Add page number to thank you slide
        page_number = len(prs.slides)
        add_page_number(thank_you_slide, page_number)
        
        return thank_you_slide
        
    except Exception as e:
        logger.error(f"Error creating thank you slide: {e}")
        return None
    
# Legacy function for backward compatibility
def create_od_ppt_slide(slide, df, title):
    """Legacy function - use create_od_ppt_slides for automatic splitting"""
    create_single_od_slide(slide, df, title)

def calculate_od_values(os_jan, os_feb, total_sale, selected_month_str,
                       os_jan_due_date_col, os_jan_ref_date_col, os_jan_net_value_col, os_jan_exec_col, os_jan_sl_code_col, os_jan_area_col,
                       os_feb_due_date_col, os_feb_ref_date_col, os_feb_net_value_col, os_feb_exec_col, os_feb_sl_code_col, os_feb_area_col,
                       sale_bill_date_col, sale_due_date_col, sale_value_col, sale_exec_col, sale_sl_code_col, sale_area_col,
                       selected_executives, selected_branches=None):

   # Convert and validate numeric cols
   for df, col, file in [
       (os_jan, os_jan_net_value_col, "OS Jan"),
       (os_feb, os_feb_net_value_col, "OS Feb"),
       (total_sale, sale_value_col, "Total Sale")
   ]:
       try:
           df[col] = pd.to_numeric(df[col], errors='coerce')
           if df[col].isna().all():
               st.error(f"Column '{col}' in {file} contains no valid numeric data.")
               return None
       except Exception as e:
           st.error(f"Error processing column '{col}' in {file}: {e}")
           return None

   # Clip negative
   os_jan[os_jan_net_value_col] = os_jan[os_jan_net_value_col].clip(lower=0)
   os_feb[os_feb_net_value_col] = os_feb[os_feb_net_value_col].clip(lower=0)

   # Standardize branch
   os_jan[os_jan_area_col] = os_jan[os_jan_area_col].apply(extract_area_name).astype(str).str.strip().str.upper()
   os_feb[os_feb_area_col] = os_feb[os_feb_area_col].apply(extract_area_name).astype(str).str.strip().str.upper()
   total_sale[sale_area_col] = total_sale[sale_area_col].apply(extract_area_name).astype(str).str.strip().str.upper()

   # Branch filter
   if selected_branches:
       os_jan = os_jan[os_jan[os_jan_area_col].isin([b.upper() for b in selected_branches])]
       os_feb = os_feb[os_feb[os_feb_area_col].isin([b.upper() for b in selected_branches])]
       total_sale = total_sale[total_sale[sale_area_col].isin([b.upper() for b in selected_branches])]
       if os_jan.empty or os_feb.empty or total_sale.empty:
           st.error(f"No data found for selected branches: {', '.join(selected_branches)}")
           return None

   # Date conversion and exec normalization
   os_jan[os_jan_due_date_col] = pd.to_datetime(os_jan[os_jan_due_date_col], errors='coerce')
   os_jan[os_jan_ref_date_col] = pd.to_datetime(os_jan.get(os_jan_ref_date_col), errors='coerce')
   os_jan["SL Code"] = os_jan[os_jan_sl_code_col].astype(str)
   os_jan["EXECUTIVE"] = os_jan[os_jan_exec_col].astype(str).str.strip().str.upper()

   os_feb[os_feb_due_date_col] = pd.to_datetime(os_feb[os_feb_due_date_col], errors='coerce')
   os_feb[os_feb_ref_date_col] = pd.to_datetime(os_feb.get(os_feb_ref_date_col), errors='coerce')
   os_feb["SL Code"] = os_feb[os_feb_sl_code_col].astype(str)
   os_feb["EXECUTIVE"] = os_feb[os_feb_exec_col].astype(str).str.strip().str.upper()

   total_sale[sale_bill_date_col] = pd.to_datetime(total_sale[sale_bill_date_col], errors='coerce')
   total_sale[sale_due_date_col] = pd.to_datetime(total_sale[sale_due_date_col], errors='coerce')
   total_sale["SL Code"] = total_sale[sale_sl_code_col].astype(str)
   total_sale["EXECUTIVE"] = total_sale[sale_exec_col].astype(str).str.strip().str.upper()

   # Exec filter
   if selected_branches:
       branch_os_jan = os_jan[os_jan[os_jan_area_col].isin([b.upper() for b in selected_branches])]
       branch_os_feb = os_feb[os_feb[os_feb_area_col].isin([b.upper() for b in selected_branches])]
       branch_sale = total_sale[total_sale[sale_area_col].isin([b.upper() for b in selected_branches])]
       branch_execs = sorted(set(branch_os_jan["EXECUTIVE"].dropna().unique()) |
                             set(branch_os_feb["EXECUTIVE"].dropna().unique()) |
                             set(branch_sale["EXECUTIVE"].dropna().unique()))
       if selected_executives:
           sel_execs_upper = [str(e).strip().upper() for e in selected_executives]
           executives_to_display = [e for e in branch_execs if e in sel_execs_upper]
       else:
           executives_to_display = branch_execs
   else:
       executives_to_display = [str(e).strip().upper() for e in selected_executives] if selected_executives else \
                               sorted(set(os_jan["EXECUTIVE"].dropna().unique()) |
                                      set(os_feb["EXECUTIVE"].dropna().unique()) |
                                      set(total_sale["EXECUTIVE"].dropna().unique()))

   os_jan = os_jan[os_jan["EXECUTIVE"].isin(executives_to_display)]
   os_feb = os_feb[os_feb["EXECUTIVE"].isin(executives_to_display)]
   total_sale = total_sale[total_sale["EXECUTIVE"].isin(executives_to_display)]
   if os_jan.empty or os_feb.empty or total_sale.empty:
       st.error("No data after filtering.")
       return None

   specified_date = pd.to_datetime("01-" + selected_month_str, format="%d-%b-%y")
   specified_month_end = specified_date + pd.offsets.MonthEnd(0)

   # Due Target - Updated column name
   due_target_sum = os_jan[os_jan[os_jan_due_date_col] <= specified_month_end] \
       .groupby("EXECUTIVE")[os_jan_net_value_col].sum().reset_index()
   due_target_sum.columns = ["EXECUTIVE", "DUE TARGET/L"]

   # OS Jan Coll
   os_jan_coll_sum = os_jan[os_jan[os_jan_due_date_col] <= specified_month_end] \
       .groupby("EXECUTIVE")[os_jan_net_value_col].sum().reset_index()
   os_jan_coll_sum.columns = ["EXECUTIVE", "OS Jan Coll"]

   # OS Feb Coll
   os_feb_coll_sum = os_feb[(os_feb[os_feb_ref_date_col] < specified_date) &
                            (os_feb[os_feb_due_date_col] <= specified_month_end)] \
       .groupby("EXECUTIVE")[os_feb_net_value_col].sum().reset_index()
   os_feb_coll_sum.columns = ["EXECUTIVE", "OS Feb Coll"]

   # Collection + early Overall % - Ensuring exact same calculation logic
   collection = os_jan_coll_sum.merge(os_feb_coll_sum, on="EXECUTIVE", how="outer").fillna(0)
   collection["COLLECTION ACHIEVED/L"] = collection["OS Jan Coll"] - collection["OS Feb Coll"]
   collection = collection.merge(due_target_sum, on="EXECUTIVE", how="outer").fillna(0)
   collection["OVERALL % ACHIEVED"] = np.where(
       collection["DUE TARGET/L"] > 0,
       (collection["COLLECTION ACHIEVED/L"] / collection["DUE TARGET/L"]) * 100,
       0
   )
   # Overdue - Updated column name
   overdue_sum = total_sale[
       total_sale[sale_bill_date_col].between(specified_date, specified_month_end) &
       total_sale[sale_due_date_col].between(specified_date, specified_month_end)
   ].groupby("EXECUTIVE")[sale_value_col].sum().reset_index()
   overdue_sum.columns = ["EXECUTIVE", "FOR THE MONTH OVERDUE/L"]

   # Month collection + early Selected Month % - Ensuring exact same calculation logic
   sale_value_sum = overdue_sum.rename(columns={"FOR THE MONTH OVERDUE/L": "Sale Value"})
   os_feb_month_sum = os_feb[
       os_feb[os_feb_ref_date_col].between(specified_date, specified_month_end) &
       os_feb[os_feb_due_date_col].between(specified_date, specified_month_end)
   ].groupby("EXECUTIVE")[os_feb_net_value_col].sum().reset_index()
   os_feb_month_sum.columns = ["EXECUTIVE", "OS Month Collection"]

   month_collection = sale_value_sum.merge(os_feb_month_sum, on="EXECUTIVE", how="outer").fillna(0)
   month_collection["FOR THE MONTH COLLECTION/L"] = month_collection["Sale Value"] - month_collection["OS Month Collection"]
   month_collection["FOR THE MONTH % ACHIEVED"] = np.where(
       month_collection["Sale Value"] > 0,
       (month_collection["FOR THE MONTH COLLECTION/L"] / month_collection["Sale Value"]) * 100,
       0
   )

   # Merge all - Following original order and logic exactly
   final = collection.drop(columns=["OS Jan Coll", "OS Feb Coll"]) \
       .merge(overdue_sum, on="EXECUTIVE", how="outer") \
       .merge(month_collection[["EXECUTIVE", "FOR THE MONTH COLLECTION/L", "FOR THE MONTH % ACHIEVED"]],
              on="EXECUTIVE", how="outer").fillna(0)

   # Reorder columns to match original function's exact order
   final = final[["EXECUTIVE", "DUE TARGET/L", "COLLECTION ACHIEVED/L", "OVERALL % ACHIEVED", 
                  "FOR THE MONTH OVERDUE/L", "FOR THE MONTH COLLECTION/L", "FOR THE MONTH % ACHIEVED"]]

   # Preserve exec list
   final = pd.DataFrame({'EXECUTIVE': executives_to_display}).merge(final, on='EXECUTIVE', how='left').fillna(0)

   # Remove HO/HEAD OFFICE
   final = final[~final["EXECUTIVE"].str.upper().isin(["HO", "HEAD OFFICE"])]

   # Scale + rounding only after percentages done - Updated column names
   val_cols = ["DUE TARGET/L", "COLLECTION ACHIEVED/L", "FOR THE MONTH OVERDUE/L", "FOR THE MONTH COLLECTION/L"]
   final[val_cols] = final[val_cols].div(100000)
   
   # Apply consistent 2-decimal formatting
   for col in val_cols:
       final[col] = final[col].apply(lambda x: float(f"{x:.2f}") if pd.notna(x) else 0.00)
   
   for col in ["OVERALL % ACHIEVED", "FOR THE MONTH % ACHIEVED"]:
       final[col] = final[col].apply(lambda x: float(f"{x:.2f}") if pd.notna(x) else 0.00)

   # Sort + TOTAL
   final.sort_values("EXECUTIVE", inplace=True)
   total_row = {'EXECUTIVE': 'TOTAL'}
   for col in final.columns[1:]:
       if col in ["OVERALL % ACHIEVED", "FOR THE MONTH % ACHIEVED"]:
           total_row[col] = float(f"{np.average(final[col], weights=final['DUE TARGET/L'] if col == 'OVERALL % ACHIEVED' else final['FOR THE MONTH OVERDUE/L']):.2f}")
       else:
           total_row[col] = float(f"{final[col].sum():.2f}")
   final = pd.concat([final, pd.DataFrame([total_row])], ignore_index=True)

   return final
def get_available_months(os_jan, os_feb, total_sale,
                        os_jan_due_date_col, os_jan_ref_date_col,
                        os_feb_due_date_col, os_feb_ref_date_col,
                        sale_bill_date_col, sale_due_date_col):
   months = set()
   for df, date_cols in [
       (os_jan, [os_jan_due_date_col, os_jan_ref_date_col]),
       (os_feb, [os_feb_due_date_col, os_feb_ref_date_col]),
       (total_sale, [sale_bill_date_col, sale_due_date_col])
   ]:
       for col in date_cols:
           if col in df.columns:
               df[col] = pd.to_datetime(df[col], errors='coerce')
               valid_dates = df[col].dropna()
               month_years = valid_dates.dt.strftime('%b-%y').unique()
               months.update(month_years)
   months = sorted(list(months), key=lambda x: pd.to_datetime("01-" + x, format="%d-%b-%y"))
   return months

def standardize_name(name):
   if pd.isna(name) or not name:
       return ""
   name = str(name).strip().lower()
   name = ''.join(c for c in name if c.isalnum() or c.isspace())
   name = ' '.join(word.capitalize() for word in name.split())
   general_variants = ['general', 'gen', 'generals', 'general ', 'genral', 'generl']
   if any(variant in name.lower() for variant in general_variants):
       return 'General'
   return name

def create_sl_code_mapping(ly_df, cy_df, budget_df, ly_sl_code_col, cy_sl_code_col, budget_sl_code_col, 
                          ly_company_group_col, cy_company_group_col, budget_company_group_col):
   try:
       mappings = []
       for df, sl_code_col, company_group_col in [
           (ly_df, ly_sl_code_col, ly_company_group_col),
           (cy_df, cy_sl_code_col, cy_company_group_col),
           (budget_df, budget_sl_code_col, budget_company_group_col)
       ]:
           if sl_code_col in df.columns and company_group_col in df.columns:
               subset = df[[sl_code_col, company_group_col]].dropna()
               subset = subset[subset[sl_code_col] != ""]
               mappings.append(subset.rename(columns={sl_code_col: 'SL_CODE', company_group_col: 'COMPANY_GROUP'}))
       if not mappings:
           logger.warning("No valid SL Code mappings found in any dataset")
           return {}
       combined = pd.concat(mappings, ignore_index=True)
       combined['COMPANY_GROUP'] = combined['COMPANY_GROUP'].apply(standardize_name)
       mapping_df = combined.groupby('SL_CODE')['COMPANY_GROUP'].agg(lambda x: x.mode()[0] if not x.empty else "").reset_index()
       sl_code_map = dict(zip(mapping_df['SL_CODE'], mapping_df['COMPANY_GROUP']))
       return sl_code_map
   except Exception as e:
       logger.error(f"Error creating SL Code mapping: {e}")
       st.error(f"Error creating SL Code mapping: {e}")
       return {}

def apply_sl_code_mapping(df, sl_code_col, company_group_col, sl_code_map):
   if sl_code_col not in df.columns or not sl_code_map:
       return df[company_group_col].apply(standardize_name)
   try:
       def map_company(row):
           if pd.isna(row[sl_code_col]) or row[sl_code_col] == "":
               return standardize_name(row[company_group_col])
           sl_code = str(row[sl_code_col]).strip()
           return sl_code_map.get(sl_code, standardize_name(row[company_group_col]))
       return df.apply(map_company, axis=1)
   except Exception as e:
       logger.error(f"Error applying SL Code mapping: {e}")
       st.error(f"Error applying SL Code mapping: {e}")
       return df[company_group_col].apply(standardize_name)

def calculate_product_growth(ly_df, cy_df, budget_df, ly_month, cy_month, ly_date_col, cy_date_col, 
                           ly_qty_col, cy_qty_col, ly_value_col, cy_value_col, 
                           budget_qty_col, budget_value_col, ly_company_group_col, 
                           cy_company_group_col, budget_company_group_col, 
                           ly_product_group_col, cy_product_group_col, budget_product_group_col,
                           ly_sl_code_col, cy_sl_code_col, budget_sl_code_col,
                           ly_exec_col, cy_exec_col, budget_exec_col, 
                           selected_executives=None, selected_company_groups=None, selected_product_groups=None):
   ly_df = ly_df.copy()
   cy_df = cy_df.copy()
   budget_df = budget_df.copy()
   required_cols = [(ly_df, [ly_date_col, ly_qty_col, ly_value_col, ly_company_group_col, ly_product_group_col, ly_exec_col]),
                   (cy_df, [cy_date_col, cy_qty_col, cy_value_col, cy_company_group_col, cy_product_group_col, cy_exec_col]),
                   (budget_df, [budget_qty_col, budget_value_col, budget_company_group_col, budget_product_group_col, budget_exec_col])]
   for df, cols in required_cols:
       missing_cols = [col for col in cols if col not in df.columns]
       if missing_cols:
           logger.error(f"Missing columns in DataFrame: {missing_cols}")
           st.error(f"Missing columns: {missing_cols}")
           return None
   sl_code_map = create_sl_code_mapping(
       ly_df, cy_df, budget_df, 
       ly_sl_code_col, cy_sl_code_col, budget_sl_code_col,
       ly_company_group_col, cy_company_group_col, budget_company_group_col
   )
   ly_df[ly_company_group_col] = apply_sl_code_mapping(ly_df, ly_sl_code_col, ly_company_group_col, sl_code_map)
   cy_df[cy_company_group_col] = apply_sl_code_mapping(cy_df, cy_sl_code_col, cy_company_group_col, sl_code_map)
   budget_df[budget_company_group_col] = apply_sl_code_mapping(budget_df, budget_sl_code_col, budget_company_group_col, sl_code_map)
   ly_df[ly_product_group_col] = ly_df[ly_product_group_col].apply(standardize_name)
   cy_df[cy_product_group_col] = cy_df[cy_product_group_col].apply(standardize_name)
   budget_df[budget_product_group_col] = budget_df[budget_product_group_col].apply(standardize_name)
   if selected_executives:
       if ly_exec_col in ly_df.columns:
           ly_df = ly_df[ly_df[ly_exec_col].isin(selected_executives)]
       if cy_exec_col in cy_df.columns:
           cy_df = cy_df[cy_df[cy_exec_col].isin(selected_executives)]
       if budget_exec_col in budget_df.columns:
           budget_df = budget_df[budget_df[budget_exec_col].isin(selected_executives)]
   if ly_df.empty or cy_df.empty or budget_df.empty:
       st.warning("No data remains after executive filtering. Please check executive selections.")
       return None
   ly_df[ly_date_col] = pd.to_datetime(ly_df[ly_date_col], dayfirst=True, errors='coerce', format='mixed')
   cy_df[cy_date_col] = pd.to_datetime(cy_df[cy_date_col], dayfirst=True, errors='coerce', format='mixed')
   available_ly_months = ly_df[ly_date_col].dt.strftime('%b %y').dropna().unique().tolist()
   available_cy_months = cy_df[cy_date_col].dt.strftime('%b %y').dropna().unique().tolist()
   if not available_ly_months or not available_cy_months:
       st.error("No valid dates found in LY or CY data. Please check date columns.")
       return None
   if ly_month:
       ly_filtered_df = ly_df[ly_df[ly_date_col].dt.strftime('%b %y') == ly_month]
   else:
       latest_ly_month = max(available_ly_months, key=lambda x: pd.to_datetime(f"01 {x}", format="%d %b %y"))
       ly_filtered_df = ly_df[ly_df[ly_date_col].dt.strftime('%b %y') == latest_ly_month]
       ly_month = latest_ly_month
   if cy_month:
       cy_filtered_df = cy_df[cy_df[cy_date_col].dt.strftime('%b %y') == cy_month]
   else:
       latest_cy_month = max(available_cy_months, key=lambda x: pd.to_datetime(f"01 {x}", format="%d %b %y"))
       cy_filtered_df = cy_df[cy_df[cy_date_col].dt.strftime('%b %y') == latest_cy_month]
       cy_month = latest_cy_month
   if ly_filtered_df.empty or cy_filtered_df.empty:
       st.warning(f"No data for selected months (LY: {ly_month}, CY: {cy_month}). Please check month selections.")
       return None
   company_groups = pd.concat([ly_filtered_df[ly_company_group_col], cy_filtered_df[cy_company_group_col], budget_df[budget_company_group_col]]).dropna().unique().tolist()
   if selected_company_groups:
       selected_company_groups = [standardize_name(g) for g in selected_company_groups]
       valid_groups = set(company_groups)
       invalid_groups = [g for g in selected_company_groups if g not in valid_groups]
       if invalid_groups:
           st.warning(f"The following company groups are not found in the data: {invalid_groups}. Please verify selections.")
           selected_company_groups = [g for g in selected_company_groups if g in valid_groups]
           if not selected_company_groups:
               st.error("No valid company groups selected after validation. Please select valid company groups.")
               return None
       ly_filtered_df = ly_filtered_df[ly_filtered_df[ly_company_group_col].isin(selected_company_groups)]
       cy_filtered_df = cy_filtered_df[cy_filtered_df[cy_company_group_col].isin(selected_company_groups)]
       budget_df = budget_df[budget_df[budget_company_group_col].isin(selected_company_groups)]
       if ly_filtered_df.empty or cy_filtered_df.empty:
           st.warning(f"No data remains after filtering for company groups: {selected_company_groups}. Please check company group selections or data content.")
           return None
   product_groups = pd.concat([ly_filtered_df[ly_product_group_col], cy_filtered_df[cy_product_group_col], budget_df[budget_product_group_col]]).dropna().unique().tolist()
   if selected_product_groups:
       selected_product_groups = [standardize_name(g) for g in selected_product_groups]
       valid_product_groups = set(product_groups)
       invalid_product_groups = [g for g in selected_product_groups if g not in valid_product_groups]
       if invalid_product_groups:
           st.warning(f"The following product groups are not found in the data: {invalid_product_groups}. Please verify selections.")
           selected_product_groups = [g for g in selected_product_groups if g in valid_product_groups]
           if not selected_product_groups:
               st.error("No valid product groups selected after validation. Please select valid product groups.")
               return None
       ly_filtered_df = ly_filtered_df[ly_filtered_df[ly_product_group_col].isin(selected_product_groups)]
       cy_filtered_df = cy_filtered_df[cy_filtered_df[cy_product_group_col].isin(selected_product_groups)]
       budget_df = budget_df[budget_df[budget_product_group_col].isin(selected_product_groups)]
       if ly_filtered_df.empty or cy_filtered_df.empty:
           st.warning(f"No data remains after filtering for product groups: {selected_product_groups}. Please check product group selections or data content.")
           return None
   for df, qty_col, value_col in [(ly_filtered_df, ly_qty_col, ly_value_col), (cy_filtered_df, cy_qty_col, cy_value_col), (budget_df, budget_qty_col, budget_value_col)]:
       df[qty_col] = pd.to_numeric(df[qty_col], errors='coerce').fillna(0)
       df[value_col] = pd.to_numeric(df[value_col], errors='coerce').fillna(0)
   company_groups = selected_company_groups if selected_company_groups else sorted(set(company_groups))
   if not company_groups:
       st.warning("No valid company groups found in the data. Please check company group columns.")
       return None
   result = {}
   for company in company_groups:
       # Updated column names with QTY and VALUE in labels - ALL UPPERCASE
       qty_df = pd.DataFrame(columns=['PRODUCT GROUP', 'LAST YEAR QTY/MT', 'BUDGET QTY/MT', 'CURRENT YEAR QTY/MT'])
       value_df = pd.DataFrame(columns=['PRODUCT GROUP', 'LAST YEAR VALUE/L', 'BUDGET VALUE/L', 'CURRENT YEAR VALUE/L'])
       ly_company_df = ly_filtered_df[ly_filtered_df[ly_company_group_col] == company]
       cy_company_df = cy_filtered_df[cy_filtered_df[cy_company_group_col] == company]
       budget_company_df = budget_df[budget_df[budget_company_group_col] == company]
       if ly_company_df.empty and cy_company_df.empty and budget_company_df.empty:
           continue
       company_product_groups = pd.concat([
           ly_company_df[ly_product_group_col],
           cy_company_df[cy_product_group_col],
           budget_company_df[budget_product_group_col]
       ]).dropna().unique().tolist()
       if not company_product_groups:
           continue
       if selected_product_groups:
           company_product_groups = [pg for pg in company_product_groups if pg in selected_product_groups]
           if not company_product_groups:
               continue
       ly_qty = ly_company_df.groupby([ly_company_group_col, ly_product_group_col])[ly_qty_col].sum().reset_index()
       ly_qty = ly_qty.rename(columns={ly_product_group_col: 'PRODUCT GROUP', ly_qty_col: 'LAST YEAR QTY/MT'})
       cy_qty = cy_company_df.groupby([cy_company_group_col, cy_product_group_col])[cy_qty_col].sum().reset_index()
       cy_qty = cy_qty.rename(columns={cy_product_group_col: 'PRODUCT GROUP', cy_qty_col: 'CURRENT YEAR QTY/MT'})
       budget_qty = budget_company_df.groupby([budget_company_group_col, budget_product_group_col])[budget_qty_col].sum().reset_index()
       budget_qty = budget_qty.rename(columns={budget_product_group_col: 'PRODUCT GROUP', budget_qty_col: 'BUDGET QTY/MT'})
       ly_value = ly_company_df.groupby([ly_company_group_col, ly_product_group_col])[ly_value_col].sum().reset_index()
       ly_value = ly_value.rename(columns={ly_product_group_col: 'PRODUCT GROUP', ly_value_col: 'LAST YEAR VALUE/L'})
       cy_value = cy_company_df.groupby([cy_company_group_col, cy_product_group_col])[cy_value_col].sum().reset_index()
       cy_value = cy_value.rename(columns={cy_product_group_col: 'PRODUCT GROUP', cy_value_col: 'CURRENT YEAR VALUE/L'})
       budget_value = budget_company_df.groupby([budget_company_group_col, budget_product_group_col])[budget_value_col].sum().reset_index()
       budget_value = budget_value.rename(columns={budget_product_group_col: 'PRODUCT GROUP', budget_value_col: 'BUDGET VALUE/L'})
       product_qty_df = pd.DataFrame({'PRODUCT GROUP': company_product_groups})
       product_value_df = pd.DataFrame({'PRODUCT GROUP': company_product_groups})
       qty_df = product_qty_df.merge(ly_qty[['PRODUCT GROUP', 'LAST YEAR QTY/MT']], on='PRODUCT GROUP', how='left')\
                              .merge(budget_qty[['PRODUCT GROUP', 'BUDGET QTY/MT']], on='PRODUCT GROUP', how='left')\
                              .merge(cy_qty[['PRODUCT GROUP', 'CURRENT YEAR QTY/MT']], on='PRODUCT GROUP', how='left').fillna(0)
       value_df = product_value_df.merge(ly_value[['PRODUCT GROUP', 'LAST YEAR VALUE/L']], on='PRODUCT GROUP', how='left')\
                                  .merge(budget_value[['PRODUCT GROUP', 'BUDGET VALUE/L']], on='PRODUCT GROUP', how='left')\
                                  .merge(cy_value[['PRODUCT GROUP', 'CURRENT YEAR VALUE/L']], on='PRODUCT GROUP', how='left').fillna(0)
       
       # Round all numeric columns to 2 decimal places BEFORE any calculation
       numeric_cols_qty = ['LAST YEAR QTY/MT', 'BUDGET QTY/MT', 'CURRENT YEAR QTY/MT']
       numeric_cols_value = ['LAST YEAR VALUE/L', 'BUDGET VALUE/L', 'CURRENT YEAR VALUE/L']
       
       for col in numeric_cols_qty:
           qty_df[col] = qty_df[col].round(2)
       for col in numeric_cols_value:
           value_df[col] = value_df[col].round(2)
       
       # Calculate totals first
       total_ly_qty = qty_df['LAST YEAR QTY/MT'].sum()
       total_cy_qty = qty_df['CURRENT YEAR QTY/MT'].sum() 
       total_budget_qty = qty_df['BUDGET QTY/MT'].sum()
       
       total_ly_value = value_df['LAST YEAR VALUE/L'].sum()
       total_cy_value = value_df['CURRENT YEAR VALUE/L'].sum()
       total_budget_value = value_df['BUDGET VALUE/L'].sum()
       
       # Add total rows first
       qty_totals = pd.DataFrame({
           'PRODUCT GROUP': ['TOTAL'],
           'LAST YEAR QTY/MT': [round(total_ly_qty, 2)],
           'BUDGET QTY/MT': [round(total_budget_qty, 2)],
           'CURRENT YEAR QTY/MT': [round(total_cy_qty, 2)],
       })
       qty_df = pd.concat([qty_df, qty_totals], ignore_index=True)
       
       value_totals = pd.DataFrame({
           'PRODUCT GROUP': ['TOTAL'],
           'LAST YEAR VALUE/L': [round(total_ly_value, 2)],
           'BUDGET VALUE/L': [round(total_budget_value, 2)],
           'CURRENT YEAR VALUE/L': [round(total_cy_value, 2)],
       })
       value_df = pd.concat([value_df, value_totals], ignore_index=True)
       
       # NOW CALCULATE ALL GROWTH % USING IDENTICAL FORMULA
       # One function, one formula, applied to every single row identically
       qty_growth_list = []
       for _, row in qty_df.iterrows():
           cy = float(row['CURRENT YEAR QTY/MT'])
           ly = float(row['LAST YEAR QTY/MT'])
           if ly == 0:
               growth = 0.00 if cy == 0 else 100.00
           else:
               growth = round(((cy - ly) / ly) * 100, 2)
           qty_growth_list.append(growth)
       qty_df['GROWTH %'] = qty_growth_list
       
       value_growth_list = []
       for _, row in value_df.iterrows():
           cy = float(row['CURRENT YEAR VALUE/L'])
           ly = float(row['LAST YEAR VALUE/L'])
           if ly == 0:
               growth = 0.00 if cy == 0 else 100.00
           else:
               growth = round(((cy - ly) / ly) * 100, 2)
           value_growth_list.append(growth)
       value_df['GROWTH %'] = value_growth_list
       
       # Reorder columns
       qty_df = qty_df[['PRODUCT GROUP', 'LAST YEAR QTY/MT', 'BUDGET QTY/MT', 'CURRENT YEAR QTY/MT', 'GROWTH %']]
       value_df = value_df[['PRODUCT GROUP', 'LAST YEAR VALUE/L', 'BUDGET VALUE/L', 'CURRENT YEAR VALUE/L', 'GROWTH %']]
       
       result[company] = {'qty_df': qty_df, 'value_df': value_df}
   if not result:
       st.warning("No data available after filtering. Please review filters and data.")
       return None
   return result

def create_product_growth_ppt(group_results, month_title, logo_file=None):
   try:
       prs = Presentation()
       prs.slide_width = Inches(13.33)
       prs.slide_height = Inches(7.5)
       create_title_slide(prs, f"Product Growth – {month_title}", logo_file)
       for company, data in group_results.items():
           qty_df = data['qty_df']
           value_df = data['value_df']
           add_table_slide(prs, qty_df, f"{company} - Quantity Growth", percent_cols=[4])
           add_table_slide(prs, value_df, f"{company} - Value Growth", percent_cols=[4])
       ppt_buffer = BytesIO()
       prs.save(ppt_buffer)
       ppt_buffer.seek(0)
       return ppt_buffer
   except Exception as e:
       logger.error(f"Error creating Product Growth PPT: {e}")
       st.error(f"Error creating Product Growth PPT: {e}")
       return None

def sidebar_ui():
   with st.sidebar:
       st.title("Integrated Reports Dashboard")
       st.subheader("File Uploads")        
       sales_file = st.file_uploader("Current Month Sales Excel File", type=["xlsx"], key="upload_sales")
       if sales_file:
           st.session_state.sales_file = sales_file
           st.success("✅ Current Month Sales file uploaded")
       ly_sales_file = st.file_uploader("Last Year Sales Excel File", type=["xlsx"], key="upload_ly_sales")
       if ly_sales_file:
           st.session_state.ly_sales_file = ly_sales_file
           st.success("✅ Last Year Sales file uploaded")
       budget_file = st.file_uploader("Current Year Executive Target Excel File", type=["xlsx"], key="upload_budget")
       if budget_file:
           st.session_state.budget_file = budget_file
           st.success("✅ Executive Target file uploaded")
       os_jan_file = st.file_uploader("OS-Previous Month Excel File", type=["xlsx"], key="upload_os_jan")
       if os_jan_file:
           st.session_state.os_jan_file = os_jan_file
           st.success("✅ OS-Previous Month file uploaded")
       os_feb_file = st.file_uploader("OS-Current Month Excel File", type=["xlsx"], key="upload_os_feb")
       if os_feb_file:
           st.session_state.os_feb_file = os_feb_file
           st.success("✅ OS-Current Month file uploaded")
       logo_file = st.file_uploader("Upload Logo (Optional)", type=["png", "jpg", "jpeg"], key="upload_logo")
       if logo_file:
           st.session_state.logo_file = logo_file
           st.image(logo_file, width=100, caption="Logo Preview")
           st.success("✅ Logo uploaded")
       st.divider()
       st.write("📊 **Dashboard v1.0**")
       st.write("© 2025 Asia Crystal Commodity LLP")


def main():
    sidebar_ui()
    st.title("🔄 Integrated Reports Dashboard")

    required_files = {
        "Current Month Sales File": st.session_state.sales_file,
        "Last Year Sales File": st.session_state.ly_sales_file,
        "Current Year Executive Target File": st.session_state.budget_file,
        "OS-Previous Month Excel File": st.session_state.os_jan_file,
        "OS-Current Month Excel File": st.session_state.os_feb_file
    }

    missing_files = [name for name, file in required_files.items() if file is None]

    if missing_files:
        st.warning(f"Please upload the following files in the sidebar to access full functionality: {', '.join(missing_files)}")
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("#### Required Files:")
            st.markdown(f"- Current Month Sales: {'✅ Uploaded' if st.session_state.sales_file else '❌ Missing'}")
            st.markdown(f"- Last Year Sales: {'✅ Uploaded' if st.session_state.ly_sales_file else '❌ Missing'}")
            st.markdown(f"- Executive Target File: {'✅ Uploaded' if st.session_state.budget_file else '❌ Missing'}")
        with col2:
            st.markdown("####  ")
            st.markdown(f"- OS-Previous Month File: {'✅ Uploaded' if st.session_state.os_jan_file else '❌ Missing'}")
            st.markdown(f"- OS-Current Month File: {'✅ Uploaded' if st.session_state.os_feb_file else '❌ Missing'}")

    tabs = st.tabs([
        "📊 Target vs Billed",
        "💰 OD Target vs Collection",
        "📈 Product Growth",
        "👥 Number of Billed Customers & OD Target"
    ])

    with tabs[0]:
        st.header("Target vs Billed")
        if not st.session_state.sales_file or not st.session_state.budget_file:
            st.warning("⚠️ Please upload Sales and Target files to use this tab")
        else:
            try:
                # Initialize session state for results if not exists
                if 'budget_results' not in st.session_state:
                    st.session_state.budget_results = []
                    logger.info("Initialized st.session_state.budget_results")

                # Define find_column function within the tab's scope
                def find_column(columns, target_names, default_index=0):
                    for target in target_names:
                        for col in columns:
                            if col.lower() == target.lower():
                                return col
                    return columns[default_index] if columns else None

                sales_sheets = get_excel_sheets(st.session_state.sales_file)
                budget_sheets = get_excel_sheets(st.session_state.budget_file)
                st.subheader("Configure Files")
                col1, col2 = st.columns(2)
                with col1:
                    st.write("**Sales File**")
                    sales_sheet = st.selectbox("Sales Sheet", sales_sheets, key='sales_sheet')
                    sales_header_row = st.number_input("Sales Header Row (1-based)", min_value=1, max_value=10, value=1, key='sales_header_row') - 1
                with col2:
                    st.write("**Target File**")
                    budget_sheet = st.selectbox("Target Sheet", budget_sheets, key='budget_sheet')
                    budget_header_row = st.number_input("Target Header Row (1-based)", min_value=1, max_value=10, value=1, key='budget_header_row') - 1

                # Load DataFrames
                sales_df = pd.read_excel(st.session_state.sales_file, sheet_name=sales_sheet, header=sales_header_row)
                budget_df = pd.read_excel(st.session_state.budget_file, sheet_name=budget_sheet, header=budget_header_row)
                logger.info(f"Sales DataFrame shape: {sales_df.shape}, Budget DataFrame shape: {budget_df.shape}")

                # Column mappings
                column_mappings = {
                    'sales_date': ['Date'],
                    'sales_value': ['Value', 'Invoice Value'],
                    'sales_qty': ['Actual Quantity', 'Quantity'],
                    'sales_product_group': ['Type (Make)', 'Product Group'],
                    'sales_sl_code': ['Customer Code', 'SL Code'],
                    'sales_area': ['Branch', 'Area'],
                    'sales_exec': ['Executive Name', 'Executive'],
                    'budget_value': ['Budget Value', 'Value'],
                    'budget_qty': ['Budget Quantity', 'Quantity'],
                    'budget_product_group': ['Product Group', 'Type (Make)'],
                    'budget_sl_code': ['SL Code', 'Customer Code'],
                    'budget_area': ['Branch', 'Area'],
                    'budget_exec': ['Executive Name', 'Executive'],
                    'budget_exec_code': ['Executive Code', 'Exec Code']
                }
                default_columns = {}
                for key, targets in column_mappings.items():
                    if key.startswith('sales'):
                        default_columns[key] = find_column(sales_df.columns.tolist(), targets)
                    else:
                        default_columns[key] = find_column(budget_df.columns.tolist(), targets)

                with st.expander("Column Mappings"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.subheader("Sales Columns")
                        sales_date_col = st.selectbox("Date Column", sales_df.columns.tolist(), index=sales_df.columns.tolist().index(default_columns['sales_date']) if default_columns['sales_date'] in sales_df.columns else 0, key='sales_date')
                        sales_value_col = st.selectbox("Value Column", sales_df.columns.tolist(), index=sales_df.columns.tolist().index(default_columns['sales_value']) if default_columns['sales_value'] in sales_df.columns else 0, key='sales_value')
                        sales_qty_col = st.selectbox("Quantity Column", sales_df.columns.tolist(), index=sales_df.columns.tolist().index(default_columns['sales_qty']) if default_columns['sales_qty'] in sales_df.columns else 0, key='sales_qty')
                        sales_area_col = st.selectbox("Branch Column", sales_df.columns.tolist(), index=sales_df.columns.tolist().index(default_columns['sales_area']) if default_columns['sales_area'] in sales_df.columns else 0, key='sales_area')
                    with col2:
                        st.subheader("Target Columns")
                        budget_value_col = st.selectbox("Target Value Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_columns['budget_value']) if default_columns['budget_value'] in budget_df.columns else 0, key='budget_value')
                        budget_qty_col = st.selectbox("Target Quantity Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_columns['budget_qty']) if default_columns['budget_qty'] in budget_df.columns else 0, key='budget_qty')
                        budget_area_col = st.selectbox("Target Branch Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_columns['budget_area']) if default_columns['budget_area'] in budget_df.columns else 0, key='budget_area')
                        # NEW: Executive Code column
                        budget_exec_code_col = st.selectbox("Target Executive Code Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_columns['budget_exec_code']) if default_columns['budget_exec_code'] in budget_df.columns else 0, key='budget_exec_code')
                    col3, col4 = st.columns(2)
                    with col3:
                        sales_product_group_col = st.selectbox("Product Group Column", sales_df.columns.tolist(), index=sales_df.columns.tolist().index(default_columns['sales_product_group']) if default_columns['sales_product_group'] in sales_df.columns else 0, key='sales_product_group')
                        sales_sl_code_col = st.selectbox("SL Code Column", sales_df.columns.tolist(), index=sales_df.columns.tolist().index(default_columns['sales_sl_code']) if default_columns['sales_sl_code'] in sales_df.columns else 0, key='sales_sl_code')
                    with col4:
                        budget_product_group_col = st.selectbox("Target Product Group Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_columns['budget_product_group']) if default_columns['budget_product_group'] in budget_df.columns else 0, key='budget_product_group')
                        budget_sl_code_col = st.selectbox("Target SL Code Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_columns['budget_sl_code']) if default_columns['budget_sl_code'] in budget_df.columns else 0, key='budget_sl_code')
                    sales_exec_col = st.selectbox("Sales Executive Column", sales_df.columns.tolist(), index=sales_df.columns.tolist().index(default_columns['sales_exec']) if default_columns['sales_exec'] in sales_df.columns else 0, key='sales_exec')
                    budget_exec_col = st.selectbox("Target Executive Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_columns['budget_exec']) if default_columns['budget_exec'] in budget_df.columns else 0, key='budget_exec')

                # Get available months
                try:
                    sales_df[sales_date_col] = pd.to_datetime(sales_df[sales_date_col], dayfirst=True, errors='coerce')
                    available_months = sorted(sales_df[sales_date_col].dt.strftime('%b %y').dropna().unique().tolist())
                    logger.info(f"Available months: {available_months}")
                except Exception as e:
                    st.error(f"Error processing date column: {e}")
                    logger.error(f"Error processing date column: {e}", exc_info=True)
                    available_months = []

                if not available_months:
                    st.error("No valid months found in sales date column")
                else:
                    st.subheader("Select Sales Month")
                    selected_month = st.selectbox("Month", available_months, key='sales_month')

                    st.subheader("Filter Options")
                    filter_tab1, filter_tab2 = st.tabs(["Branches", "Executives"])
                    with filter_tab1:
                        raw_branches = set(sales_df[sales_area_col].dropna().astype(str).str.strip().str.upper().unique().tolist()) | \
                                    set(budget_df[budget_area_col].dropna().astype(str).str.strip().str.upper().unique().tolist())
                        all_branches = sorted(raw_branches)
                        branch_select_all = st.checkbox("Select All Branches", value=True, key='budget_branch_all')
                        selected_branches = []
                        if branch_select_all:
                            selected_branches = all_branches
                        else:
                            num_cols = 3
                            branch_cols = st.columns(num_cols)
                            for i, branch in enumerate(all_branches):
                                col_idx = i % num_cols
                                with branch_cols[col_idx]:
                                    if st.checkbox(branch, key=f'budget_branch_{branch}'):
                                        selected_branches.append(branch)
                    with filter_tab2:
                        all_executives = set(sales_df[sales_exec_col].dropna().astype(str).unique().tolist()) | set(budget_df[budget_exec_col].dropna().astype(str).unique().tolist())
                        all_executives = sorted(all_executives)
                        exec_select_all = st.checkbox("Select All Executives", value=True, key='exec_all')
                        selected_executives = []
                        if exec_select_all:
                            selected_executives = all_executives
                        else:
                            num_cols = 3
                            exec_cols = st.columns(num_cols)
                            for i, exec_name in enumerate(all_executives):
                                col_idx = i % num_cols
                                with exec_cols[col_idx]:
                                    if st.checkbox(exec_name, key=f'exec_{exec_name}'):
                                        selected_executives.append(exec_name)

                    if st.button("Generate Report", key='generate_report'):
                        if not selected_executives:
                            st.error("Please select at least one executive")
                        else:
                            with st.spinner("Generating report..."):
                                logger.info(f"Calling calculate_budget_values with selected_month: {selected_month}, "
                                            f"selected_branches: {selected_branches}, selected_executives: {selected_executives}")
                                budget_vs_billed_value_df, budget_vs_billed_qty_df, overall_sales_qty_df, overall_sales_value_df = calculate_budget_values(
                                    sales_df, budget_df, selected_month, selected_executives,
                                    sales_date_col, sales_area_col, sales_value_col, sales_qty_col,
                                    sales_product_group_col, sales_sl_code_col, sales_exec_col,
                                    budget_area_col, budget_value_col, budget_qty_col,
                                    budget_product_group_col, budget_sl_code_col, budget_exec_col,
                                    selected_branches=selected_branches
                                )
                                logger.info(f"calculate_budget_values returned: "
                                            f"budget_vs_billed_value_df: {budget_vs_billed_value_df is not None}, "
                                            f"budget_vs_billed_qty_df: {budget_vs_billed_qty_df is not None}, "
                                            f"overall_sales_qty_df: {overall_sales_qty_df is not None}, "
                                            f"overall_sales_value_df: {overall_sales_value_df is not None}")

                                # Store all non-None DataFrames
                                dfs_info = []
                                if budget_vs_billed_qty_df is not None:
                                    dfs_info.append({
                                        'df': budget_vs_billed_qty_df,
                                        'title': f"TARGET AGAINST BILLED - QUANTITY - {selected_month}",
                                        'percent_cols': [3]
                                    })
                                if budget_vs_billed_value_df is not None:
                                    dfs_info.append({
                                        'df': budget_vs_billed_value_df,
                                        'title': f"TARGET AGAINST BILLED - VALUE - {selected_month}",
                                        'percent_cols': [3]
                                    })
                                if overall_sales_qty_df is not None:
                                    dfs_info.append({
                                        'df': overall_sales_qty_df,
                                        'title': f"OVERALL SALES - QUANTITY - {selected_month}",
                                        'percent_cols': [3]
                                    })
                                if overall_sales_value_df is not None:
                                    dfs_info.append({
                                        'df': overall_sales_value_df,
                                        'title': f"OVERALL SALES - VALUE - {selected_month}",
                                        'percent_cols': [3]
                                    })

                                if dfs_info:
                                    st.success("Success!")
                                    if budget_vs_billed_qty_df is not None:
                                        st.subheader("Target vs Billed Quantity")
                                        st.dataframe(budget_vs_billed_qty_df, use_container_width=True)
                                        qty_image = create_table_image(budget_vs_billed_qty_df, f"TARGET AGAINST BILLED - QUANTITY - {selected_month}", percent_cols=[3])
                                        if qty_image:
                                            st.image(qty_image, use_column_width=True)
                                    if budget_vs_billed_value_df is not None:
                                        st.subheader("Target vs Billed Value")
                                        st.dataframe(budget_vs_billed_value_df, use_container_width=True)
                                        value_image = create_table_image(budget_vs_billed_value_df, f"TARGET AGAINST BILLED - VALUE - {selected_month}", percent_cols=[3])
                                        if value_image:
                                            st.image(value_image, use_column_width=True)
                                    if overall_sales_qty_df is not None:
                                        st.subheader("Overall Sales Quantity")
                                        st.dataframe(overall_sales_qty_df, use_container_width=True)
                                        overall_qty_image = create_table_image(overall_sales_qty_df, f"OVERALL SALES - QUANTITY - {selected_month}", percent_cols=[3])
                                        if overall_qty_image:
                                            st.image(overall_qty_image, use_column_width=True)
                                    if overall_sales_value_df is not None:
                                        st.subheader("Overall Sales Value")
                                        st.dataframe(overall_sales_value_df, use_container_width=True)
                                        overall_value_image = create_table_image(overall_sales_value_df, f"OVERALL SALES - VALUE - {selected_month}", percent_cols=[3])
                                        if overall_value_image:
                                            st.image(overall_value_image, use_column_width=True)

                                    # Store results in session state regardless of PPT success
                                    st.session_state.budget_results = dfs_info
                                    logger.info(f"Stored {len(dfs_info)} DataFrames in st.session_state.budget_results")

                                    # Generate PPT if at least one DataFrame is not None
                                    try:
                                        ppt_buffer = None
                                        if any(x is not None for x in [budget_vs_billed_value_df, budget_vs_billed_qty_df, overall_sales_qty_df, overall_sales_value_df]):
                                            logger.info("Attempting to create PPT with available DataFrames")
                                            ppt_buffer = create_budget_ppt(
                                                budget_vs_billed_value_df, budget_vs_billed_qty_df,
                                                overall_sales_qty_df, overall_sales_value_df,
                                                selected_month,
                                                st.session_state.logo_file
                                            )
                                        else:
                                            logger.warning("No valid DataFrames for PPT generation")
                                            st.warning("Cannot generate PPT: No valid DataFrames available. Check your data and selections.")

                                        if ppt_buffer:
                                            col1, col2 = st.columns(2)
                                            with col1:
                                                unique_id = str(uuid.uuid4())[:8]
                                                st.download_button(
                                                    label="📊 Download Budget vs Billed PPT",
                                                    data=ppt_buffer,
                                                    file_name=f"Target_vs_Billed_{selected_month}_{unique_id}.pptx",
                                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                                    key=f'budget_download_{unique_id}'
                                                )
                                            with col2:
                                                # NEW: Generate proof of calculation Excel
                                                try:
                                                    proof_excel = create_proof_of_calculation_excel(
                                                        budget_df, sales_df, selected_month,
                                                        budget_exec_col, budget_exec_code_col, budget_area_col,
                                                        budget_sl_code_col, budget_product_group_col,
                                                        budget_qty_col, budget_value_col,
                                                        sales_exec_col, sales_date_col, sales_area_col,
                                                        sales_sl_code_col, sales_product_group_col,
                                                        sales_qty_col, sales_value_col,
                                                        selected_executives, selected_branches
                                                    )
                                                    
                                                    if proof_excel:
                                                        unique_id_excel = str(uuid.uuid4())[:8]
                                                        st.download_button(
                                                            label="📋 Download Proof of Calculation (Excel)",
                                                            data=proof_excel,
                                                            file_name=f"Target_vs_Billed_Proof_{selected_month}_{unique_id_excel}.xlsx",
                                                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                                            key=f'proof_download_{unique_id_excel}'
                                                        )
                                                    else:
                                                        st.warning("Could not generate proof of calculation file.")
                                                except Exception as e:
                                                    st.error(f"Error generating proof file: {e}")
                                        else:
                                            logger.warning("PPT buffer is None, skipping download button")
                                            st.warning("Failed to generate PPT. Check your data and selections.")
                                    except Exception as e:
                                        error_msg = f"Error creating Budget PPT: {str(e)}"
                                        st.error(error_msg)
                                        logger.error(error_msg, exc_info=True)
                                        st.error(traceback.format_exc())
                                else:
                                    st.error("No valid DataFrames generated. Check data and selections.")
                                    logger.error("No valid DataFrames generated by calculate_budget_values")

            except Exception as e:
                error_msg = f"Error in Budget vs Billed tab: {str(e)}"
                st.error(error_msg)
                logger.error(error_msg, exc_info=True)
                st.error(traceback.format_exc())
    with tabs[1]:
        st.header("OD Target vs Collection Report")
        if not st.session_state.os_jan_file or not st.session_state.os_feb_file or not st.session_state.sales_file:
            st.warning("⚠️ Please upload OS-Previous Month, OS-Current Month and Sales files to use this tab")
        else:
            try:
                os_jan_sheets = get_excel_sheets(st.session_state.os_jan_file)
                os_feb_sheets = get_excel_sheets(st.session_state.os_feb_file)
                sales_sheets = get_excel_sheets(st.session_state.sales_file)
                st.subheader("Sheet Selection")
                col1, col2, col3 = st.columns(3)
                with col1:
                    os_jan_sheet = st.selectbox("OS-Previous Month Sheet", os_jan_sheets, key='od_os_jan_sheet')
                    os_jan_header = st.number_input("OS-Previous Month Header Row (1-based)", min_value=1, max_value=10, value=1, key='od_os_jan_header') - 1
                with col2:
                    os_feb_sheet = st.selectbox("OS-Current Month Sheet", os_feb_sheets, key='od_os_feb_sheet')
                    os_feb_header = st.number_input("OS-Current Month Header Row (1-based)", min_value=1, max_value=10, value=1, key='od_os_feb_header') - 1
                with col3:
                    sales_sheet = st.selectbox("Sales Sheet", sales_sheets, key='od_sales_sheet')
                    sales_header = st.number_input("Sales Header Row (1-based)", min_value=1, max_value=10, value=1, key='od_sales_header') - 1
                os_jan = pd.read_excel(st.session_state.os_jan_file, sheet_name=os_jan_sheet, header=os_jan_header)
                os_feb = pd.read_excel(st.session_state.os_feb_file, sheet_name=os_feb_sheet, header=os_feb_header)
                total_sale = pd.read_excel(st.session_state.sales_file, sheet_name=sales_sheet, header=sales_header)

                def find_column(columns, target_names, default_index=0):
                    for target in target_names:
                        for col in columns:
                            if col.lower() == target.lower():
                                return col
                    return columns[default_index] if columns else None

                # OS-First column auto-mapping
                os_jan_column_mappings = {
                    'due_date': ['Due Date'],
                    'ref_date': ['Ref. Date'],
                    'net_value': ['Net Value'],
                    'executive': ['Executive Name'],
                    'sl_code': ['Party Code'],
                    'area': ['Branch', 'Area']  # Added for branch selection
                }

                # OS-Second column auto-mapping
                os_feb_column_mappings = {
                    'due_date': ['Due Date'],
                    'ref_date': ['Ref. Date'],
                    'net_value': ['Net Value'],
                    'executive': ['Executive Name'],
                    'sl_code': ['Party Code'],
                    'area': ['Branch', 'Area']  # Added for branch selection
                }

                # Total Sale column auto-mapping
                sales_column_mappings = {
                    'bill_date': ['Date'],
                    'due_date': ['Due Date'],
                    'value': ['Invoice Value'],
                    'executive': ['Executive Name'],
                    'sl_code': ['Customer Code'],
                    'area': ['Branch', 'Area']  # Added for branch selection
                }

                # Initialize default selections
                default_os_jan_cols = {}
                for key, targets in os_jan_column_mappings.items():
                    default_os_jan_cols[key] = find_column(os_jan.columns.tolist(), targets)

                default_os_feb_cols = {}
                for key, targets in os_feb_column_mappings.items():
                    default_os_feb_cols[key] = find_column(os_feb.columns.tolist(), targets)

                default_sales_cols = {}
                for key, targets in sales_column_mappings.items():
                    default_sales_cols[key] = find_column(total_sale.columns.tolist(), targets)

                with st.expander("Column Mappings"):
                    st.subheader("OS-First Column Mapping")
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        os_jan_due_date_col = st.selectbox(
                            "Due Date Column",
                            os_jan.columns.tolist(),
                            index=os_jan.columns.tolist().index(default_os_jan_cols['due_date']) if default_os_jan_cols['due_date'] in os_jan.columns else 0,
                            key='od_os_jan_due_date'
                        )
                        os_jan_ref_date_col = st.selectbox(
                            "Reference Date Column",
                            os_jan.columns.tolist(),
                            index=os_jan.columns.tolist().index(default_os_jan_cols['ref_date']) if default_os_jan_cols['ref_date'] in os_jan.columns else 0,
                            key='od_os_jan_ref_date'
                        )
                    with col2:
                        os_jan_net_value_col = st.selectbox(
                            "Net Value Column",
                            os_jan.columns.tolist(),
                            index=os_jan.columns.tolist().index(default_os_jan_cols['net_value']) if default_os_jan_cols['net_value'] in os_jan.columns else 0,
                            key='od_os_jan_net_value'
                        )
                        os_jan_sl_code_col = st.selectbox(
                            "SL Code Column",
                            os_jan.columns.tolist(),
                            index=os_jan.columns.tolist().index(default_os_jan_cols['sl_code']) if default_os_jan_cols['sl_code'] in os_jan.columns else 0,
                            key='od_os_jan_sl_code'
                        )
                    with col3:
                        os_jan_exec_col = st.selectbox(
                            "Executive Column",
                            os_jan.columns.tolist(),
                            index=os_jan.columns.tolist().index(default_os_jan_cols['executive']) if default_os_jan_cols['executive'] in os_jan.columns else 0,
                            key='od_os_jan_exec'
                        )
                        os_jan_area_col = st.selectbox(
                            "Branch Column",
                            os_jan.columns.tolist(),
                            index=os_jan.columns.tolist().index(default_os_jan_cols['area']) if default_os_jan_cols['area'] in os_jan.columns else 0,
                            key='od_os_jan_area'
                        )

                    st.subheader("OS-Second Column Mapping")
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        os_feb_due_date_col = st.selectbox(
                            "Due Date Column",
                            os_feb.columns.tolist(),
                            index=os_feb.columns.tolist().index(default_os_feb_cols['due_date']) if default_os_feb_cols['due_date'] in os_feb.columns else 0,
                            key='od_os_feb_due_date'
                        )
                        os_feb_ref_date_col = st.selectbox(
                            "Reference Date Column",
                            os_feb.columns.tolist(),
                            index=os_feb.columns.tolist().index(default_os_feb_cols['ref_date']) if default_os_feb_cols['ref_date'] in os_feb.columns else 0,
                            key='od_os_feb_ref_date'
                        )
                    with col2:
                        os_feb_net_value_col = st.selectbox(
                            "Net Value Column",
                            os_feb.columns.tolist(),
                            index=os_feb.columns.tolist().index(default_os_feb_cols['net_value']) if default_os_feb_cols['net_value'] in os_feb.columns else 0,
                            key='od_os_feb_net_value'
                        )
                        os_feb_sl_code_col = st.selectbox(
                            "SL Code Column",
                            os_feb.columns.tolist(),
                            index=os_feb.columns.tolist().index(default_os_feb_cols['sl_code']) if default_os_feb_cols['sl_code'] in os_feb.columns else 0,
                            key='od_os_feb_sl_code'
                        )
                    with col3:
                        os_feb_exec_col = st.selectbox(
                            "Executive Column",
                            os_feb.columns.tolist(),
                            index=os_feb.columns.tolist().index(default_os_feb_cols['executive']) if default_os_feb_cols['executive'] in os_feb.columns else 0,
                            key='od_os_feb'
                        )
                        os_feb_area_col = st.selectbox(
                            "Branch Column",
                            os_feb.columns.tolist(),
                            index=os_feb.columns.tolist().index(default_os_feb_cols['area']) if default_os_feb_cols['area'] in os_feb.columns else 0,
                            key='od_os_feb_area'
                        )

                    st.subheader("Total Sale Column Mapping")
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        sale_bill_date_col = st.selectbox(
                            "Sales Bill Date Column",
                            total_sale.columns.tolist(),
                            index=total_sale.columns.tolist().index(default_sales_cols['bill_date']) if default_sales_cols['bill_date'] in total_sale.columns else 0,
                            key='od_sale_bill_date'
                        )
                        sale_due_date_col = st.selectbox(
                            "Sales Due Date Column",
                            total_sale.columns.tolist(),
                            index=total_sale.columns.tolist().index(default_sales_cols['due_date']) if default_sales_cols['due_date'] in total_sale.columns else 0,
                            key='od_sale_due_date'
                        )
                    with col2:
                        sale_value_col = st.selectbox(
                            "Sales Value Column",
                            total_sale.columns.tolist(),
                            index=total_sale.columns.tolist().index(default_sales_cols['value']) if default_sales_cols['value'] in total_sale.columns else 0,
                            key='od_sale_value'
                        )
                        sale_sl_code_col = st.selectbox(
                            "Sales SL Code Column",
                            total_sale.columns.tolist(),
                            index=total_sale.columns.tolist().index(default_sales_cols['sl_code']) if default_sales_cols['sl_code'] in total_sale.columns else 0,
                            key='od_sale_sl_code'
                        )
                    with col3:
                        sale_exec_col = st.selectbox(
                            "Sales Executive Column",
                            total_sale.columns.tolist(),
                            index=total_sale.columns.tolist().index(default_sales_cols['executive']) if default_sales_cols['executive'] in total_sale.columns else 0,
                            key='od_sale_exec'
                        )
                        sale_area_col = st.selectbox(
                            "Sales Branch Column",
                            total_sale.columns.tolist(),
                            index=total_sale.columns.tolist().index(default_sales_cols['area']) if default_sales_cols['area'] in total_sale.columns else 0,
                            key='od_sale_area'
                        )
                    available_months = get_available_months(
                        os_jan, os_feb, total_sale,
                        os_jan_due_date_col, os_jan_ref_date_col,
                        os_feb_due_date_col, os_feb_ref_date_col,
                        sale_bill_date_col, sale_due_date_col
                    )
                    if not available_months:
                        st.error("No valid months found in the date columns.")
                    else:
                        st.subheader("Select Month")
                        selected_month_str = st.selectbox("Month", available_months, key='od_month')
                        st.subheader("Filter Options")
                        filter_tab1, filter_tab2 = st.tabs(["Branches", "Executives"])
                        with filter_tab1:
                            raw_branches = set(os_jan[os_jan_area_col].apply(extract_area_name).dropna().astype(str).str.upper().unique().tolist()) | \
                                        set(os_feb[os_feb_area_col].apply(extract_area_name).dropna().astype(str).str.upper().unique().tolist()) | \
                                        set(total_sale[sale_area_col].apply(extract_area_name).dropna().astype(str).str.upper().unique().tolist())
                            all_branches = sorted(raw_branches)
                            branch_select_all = st.checkbox("Select All Branches", value=True, key='od_vs_branch_all')
                            selected_branches = []
                            if branch_select_all:
                                selected_branches = all_branches
                            else:
                                num_cols = 3
                                branch_cols = st.columns(num_cols)
                                for i, branch in enumerate(all_branches):
                                    col_idx = i % num_cols
                                    with branch_cols[col_idx]:
                                        if st.checkbox(branch, key=f'od_vs_branch_{branch}'):
                                            selected_branches.append(branch)
                        with filter_tab2:
                            all_executives = set()
                            for df, exec_col in [
                                (os_jan, os_jan_exec_col),
                                (os_feb, os_feb_exec_col),
                                (total_sale, sale_exec_col)
                            ]:
                                if exec_col in df.columns:
                                    execs = df[exec_col].dropna().astype(str).unique().tolist()
                                    all_executives.update(execs)
                            all_executives = sorted(list(all_executives))

                            exec_select_all = st.checkbox("Select All Executives", value=True, key='od_vs_exec_all')
                            selected_od_executives = []
                            if exec_select_all:
                                selected_od_executives = all_executives
                            else:
                                num_cols = 3
                                exec_cols = st.columns(num_cols)
                                for i, exec_name in enumerate(all_executives):
                                    col_idx = i % num_cols
                                    with exec_cols[col_idx]:
                                        if st.checkbox(exec_name, key=f'od_vs_exec_{exec_name}'):
                                            selected_od_executives.append(exec_name)
                        if st.button("Generate OD Target vs Collection Report", key='od_vs_generate'):
                            if not selected_od_executives:
                                st.error("Please select at least one executive.")
                            else:
                                with st.spinner("Generating report..."):
                                    final_df = calculate_od_values(
                                        os_jan, os_feb, total_sale, selected_month_str,
                                        os_jan_due_date_col, os_jan_ref_date_col, os_jan_net_value_col, os_jan_exec_col, os_jan_sl_code_col, os_jan_area_col,
                                        os_feb_due_date_col, os_feb_ref_date_col, os_feb_net_value_col, os_feb_exec_col, os_feb_sl_code_col, os_feb_area_col,
                                        sale_bill_date_col, sale_due_date_col, sale_value_col, sale_exec_col, sale_sl_code_col, sale_area_col,
                                        selected_od_executives, selected_branches
                                    )
                                    if final_df is not None and not final_df.empty:
                                        st.success("Success!")
                                        st.subheader("OD Target vs Collection Results")
                                        st.dataframe(final_df)
                                        img_buffer = create_table_image(final_df, f"OD TARGET vs VALUE - {selected_month_str}", percent_cols=[3, 6])
                                        if img_buffer:
                                            st.image(img_buffer, use_column_width=True)
                                        prs = Presentation()
                                        prs.slide_width = Inches(13.33)
                                        prs.slide_height = Inches(7.5)

                                        create_title_slide(prs, f"OD Target vs Collection - {selected_month_str}", st.session_state.logo_file)

                                        add_table_slide(prs, final_df, f"OD Target vs Collection - {selected_month_str}", percent_cols=[3, 6])

                                        ppt_buffer = BytesIO()
                                        prs.save(ppt_buffer)
                                        ppt_buffer.seek(0)

                                        unique_id = str(uuid.uuid4())[:8]
                                        st.download_button(
                                            label="Download OD Target vs Collection PPT",
                                            data=ppt_buffer,
                                            file_name=f"OD_Target_vs_Collection_{selected_month_str}_{unique_id}.pptx",
                                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                            key=f"od_vs_download_{unique_id}"
                                        )
                                        st.session_state.od_vs_results = [
                                            {'df': final_df, 'title': f"OD TARGET vs COLLECTION - {selected_month_str}", 'percent_cols': [3, 6]}
                                        ]
                                    else:
                                        st.error("Failed to generate report. Check your data and selections.")
            except Exception as e:
                st.error(f"Error in OD Target vs Collection tab: {e}")
                st.error(traceback.format_exc())

    with tabs[2]:
       st.header("Product Growth Dashboard")
       if not st.session_state.ly_sales_file or not st.session_state.sales_file or not st.session_state.budget_file:
         st.warning("⚠️ Please upload Last Year Sales, Current Month Sales, and Executive Target files to use this tab.")
       else:
         try:
            # Get sheets from all three files
            ly_sales_sheets = get_excel_sheets(st.session_state.ly_sales_file)
            cy_sales_sheets = get_excel_sheets(st.session_state.sales_file)
            budget_sheets = get_excel_sheets(st.session_state.budget_file)

            st.subheader("Configure Files")
            col1, col2, col3 = st.columns(3)

            with col1:
                st.write("**Last Year Sales File**")
                ly_sales = st.selectbox("Last Year Sales Sheet", ly_sales_sheets, key='pg_ly_sales')
                ly_header_row = st.number_input("Last Year Header Row (1-based)", min_value=1, max_value=10, value=1, key='pg_ly_header_row') - 1
            with col2:
                st.write("**Current Month Sales File**")
                cy_sales = st.selectbox("Current Month Sales Sheet", cy_sales_sheets, key='pg_sales_cy_sales')
                cy_header_row = st.number_input("Current Month Header Row (1-based)", min_value=1, max_value=10, value=1, key='pg_cy_header_row') - 1
            with col3:
                st.write("**Target File**")
                budget_product_sheets = st.selectbox("Target Sheet", budget_sheets, key='pg_budget_sheet')  # Changed key to 'pg_budget_sheet'
                budget_header_row = st.number_input("Target Header Row (1-based)", min_value=1, value=1, key='pg_budget_header_row') - 1

            # Load data from separate files
            ly_df = pd.read_excel(st.session_state.ly_sales_file, sheet_name=ly_sales, header=ly_header_row)
            cy_df = pd.read_excel(st.session_state.sales_file, sheet_name=cy_sales, header=cy_header_row)
            budget_df = pd.read_excel(st.session_state.budget_file, sheet_name=budget_product_sheets, header=budget_header_row)

            def find_column(columns, target_names, default=None):
                """Helper function to find a column exactly matching any of the target names (case-insensitive)."""
                for target in target_names:
                    for col in columns:
                        if col.lower() == target.lower():
                            return col
                return default

            # Column mappings
            ly_column_mappings = {
                'date': ['Date'],
                'value': ['Value'],
                'product_group': ['Type (Make)'],
                'quantity': ['Actual Quantity'],
                'company_group': ['Company Group'],
                'executive': ['Executive Name'],
                'sl_code': ['Customer Code']
            }

            cy_column_mappings = {
                'date': ['Date'],
                'value': ['Value'],
                'product_group': ['Type (Make)'],
                'quantity': ['Actual Quantity'],
                'company_group': ['Company Group'],
                'executive': ['Executive Name'],
                'sl_code': ['Customer Code']
            }

            budget_column_mappings = {
                'company_group': ['Company Group'],
                'product_group': ['Product Group'],
                'quantity': ['Budget Quantity', 'Quantity'],
                'value': ['Budget Value', 'Value'],
                'executive': ['Executive Name'],
                'sl_code': ['SL Code']
            }

            default_ly_cols = {}
            for key, targets in ly_column_mappings.items():
                default_ly_cols[key] = find_column(ly_df.columns.tolist(), targets)

            default_cy_cols = {}
            for key, targets in cy_column_mappings.items():
                default_cy_cols[key] = find_column(cy_df.columns.tolist(), targets)

            default_budget_cols = {}
            for key, targets in budget_column_mappings.items():
                default_budget_cols[key] = find_column(budget_df.columns.tolist(), targets)

            with st.expander("Column Mappings"):
                st.subheader("Last Year Columns")
                col1, col2, col3 = st.columns(3)
                with col1:
                    ly_date_col = st.selectbox("LY Date Column", ly_df.columns.tolist(), index=ly_df.columns.tolist().index(default_ly_cols['date']) if default_ly_cols['date'] else 0, key='pg_ly_date')
                    ly_qty_col = st.selectbox("LY Quantity Column", ly_df.columns.tolist(), index=ly_df.columns.tolist().index(default_ly_cols['quantity']) if default_ly_cols['quantity'] else 0, key='pg_ly_qty')
                with col2:
                    ly_value_col = st.selectbox("LY Value Column", ly_df.columns.tolist(), index=ly_df.columns.tolist().index(default_ly_cols['value']) if default_ly_cols['value'] else 0, key='pg_ly_value')
                    ly_company_group_col = st.selectbox("LY Company Group Column", ly_df.columns.tolist(), index=ly_df.columns.tolist().index(default_ly_cols['company_group']) if default_ly_cols['company_group'] else 0, key='pg_ly_company_group')
                with col3:
                    ly_product_group_col = st.selectbox("LY Product Group Column", ly_df.columns.tolist(), index=ly_df.columns.tolist().index(default_ly_cols['product_group']) if default_ly_cols['product_group'] else 0, key='pg_ly_product')
                    ly_sl_code_col = st.selectbox("LY SL Code Column", ly_df.columns.tolist(), index=ly_df.columns.tolist().index(default_ly_cols['sl_code']) if default_ly_cols['sl_code'] else 0, key='pg_ly_sl')
                    ly_exec_col = st.selectbox("LY Executive Column", ly_df.columns.tolist(), index=ly_df.columns.tolist().index(default_ly_cols['executive']) if default_ly_cols['executive'] else 0, key='pg_ly_exec')

                st.subheader("Current Month Columns")
                col1, col2, col3 = st.columns(3)
                with col1:
                    cy_date_col = st.selectbox("CY Date Column", cy_df.columns.tolist(), index=cy_df.columns.tolist().index(default_cy_cols['date']) if default_cy_cols['date'] else 0, key='pg_cy_date')
                    cy_qty_col = st.selectbox("CY Quantity Column", cy_df.columns.tolist(), index=cy_df.columns.tolist().index(default_cy_cols['quantity']) if default_cy_cols['quantity'] else 0, key='pg_cy_qty')
                with col2:
                    cy_value_col = st.selectbox("CY Value Column", cy_df.columns.tolist(), index=cy_df.columns.tolist().index(default_cy_cols['value']) if default_cy_cols['value'] else 0, key='pg_cy_value')
                    cy_company_group_col = st.selectbox("CY Company Group Column", cy_df.columns.tolist(), index=cy_df.columns.tolist().index(default_cy_cols['company_group']) if default_cy_cols['company_group'] else 0, key='pg_cy_company_group')
                with col3:
                    cy_product_group_col = st.selectbox("CY Product Group Column", cy_df.columns.tolist(), index=cy_df.columns.tolist().index(default_cy_cols['product_group']) if default_cy_cols['product_group'] else 0, key='pg_cy_product')
                    cy_sl_code_col = st.selectbox("CY SL Code Column", cy_df.columns.tolist(), index=cy_df.columns.tolist().index(default_cy_cols['sl_code']) if default_cy_cols['sl_code'] else 0, key='pg_cy_sl')
                    cy_exec_col = st.selectbox("CY Executive Column", cy_df.columns.tolist(), index=cy_df.columns.tolist().index(default_cy_cols['executive']) if default_cy_cols['executive'] else 0, key='pg_cy_exec')

                st.subheader("Target Columns")
                col1, col2, col3 = st.columns(3)
                with col1:
                    budget_qty_col = st.selectbox("Target Quantity Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_budget_cols['quantity']) if default_budget_cols['quantity'] else 0, key='pg_budget_qty')
                    budget_value_col = st.selectbox("Target Value Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_budget_cols['value']) if default_budget_cols['value'] else 0, key='pg_budget_value')
                with col2:
                    budget_company_group_col = st.selectbox("Target Company Group Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_budget_cols['company_group']) if default_budget_cols['company_group'] else 0, key='pg_budget_company_group')
                    budget_sl_code_col = st.selectbox("Target SL Code Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_budget_cols['sl_code']) if default_budget_cols['sl_code'] else 0, key='pg_budget_sl')
                with col3:
                    budget_product_group_col = st.selectbox("Target Product Group Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_budget_cols['product_group']) if default_budget_cols['product_group'] else 0, key='pg_budget_product')
                    budget_exec_col = st.selectbox("Target Executive Column", budget_df.columns.tolist(), index=budget_df.columns.tolist().index(default_budget_cols['executive']) if default_budget_cols['executive'] else 0, key='pg_budget_exec')

            # Replace empty strings with NaN and standardize
            for df, col in [
                (ly_df, ly_company_group_col), (cy_df, cy_company_group_col), (budget_df, budget_company_group_col),
                (ly_df, ly_product_group_col), (cy_df, cy_product_group_col), (budget_df, budget_product_group_col),
            ]:
                df[col] = df[col].replace("", np.nan)
                df[col] = df[col].fillna("")
            
            ly_groups = sorted(set(ly_df[ly_company_group_col].apply(standardize_name).unique()))
            cy_groups = sorted(set(cy_df[cy_company_group_col].apply(standardize_name).unique()))
            budget_groups = sorted(set(budget_df[budget_company_group_col].apply(standardize_name).unique()))
            all_company_groups = sorted(set([g for g in ly_groups + cy_groups + budget_groups if g]))
            ly_product_groups = sorted(set(ly_df[ly_product_group_col].apply(standardize_name).unique()))
            cy_product_groups = sorted(set(cy_df[cy_product_group_col].apply(standardize_name).unique()))
            budget_product_groups = sorted(set(budget_df[budget_product_group_col].apply(standardize_name).unique()))
            all_product_groups = sorted(set([g for g in ly_product_groups + cy_product_groups + budget_product_groups if g]))

            st.subheader("Select Sales Months")
            ly_dates = pd.to_datetime(ly_df[ly_date_col], dayfirst=True, errors='coerce')
            cy_dates = pd.to_datetime(cy_df[cy_date_col], dayfirst=True, errors='coerce')
            ly_months = sorted(ly_dates.dt.strftime('%b %y').dropna().unique().tolist())
            cy_months = sorted(cy_dates.dt.strftime('%b %y').dropna().unique().tolist())
            col1, col2 = st.columns(2)
            with col1:
                selected_ly_month = st.selectbox("Last Year Month", options=ly_months, index=0, key='pg_ly_month')
            with col2:
                selected_cy_month = st.selectbox("Current Year Month", options=cy_months, index=0, key='pg_cy_month')
            st.subheader("Filter Options")
            filter_tabs = st.tabs(["Executives", "Company Groups", "Product Groups"])
            with filter_tabs[0]:
                all_execs = set()
                for df, col in [(ly_df, ly_exec_col), (cy_df, cy_exec_col), (budget_df, budget_exec_col)]:
                    execs = df[col].dropna().astype(str).unique().tolist()
                    all_execs.update(execs)
                all_execs = sorted(all_execs)
                pg_exec_select_all = st.checkbox("Select All Executives", value=True, key='pg_exec_all')
                selected_executives = []
                if pg_exec_select_all:
                    selected_executives = all_execs
                else:
                    num_cols = 3
                    exec_cols = st.columns(num_cols)
                    for i, exec_name in enumerate(all_execs):
                        col_idx = i % num_cols
                        with exec_cols[col_idx]:
                            if st.checkbox(exec_name, key=f'pg_exec_{exec_name}'):
                                selected_executives.append(exec_name)
            with filter_tabs[1]:
                pg_company_select_all = st.checkbox("Select All Company Groups", value=True, key='pg_company_all')
                selected_companies = []
                if pg_company_select_all:
                    selected_companies = all_company_groups
                else:
                    num_cols = 3
                    company_cols = st.columns(num_cols)
                    for i, group in enumerate(all_company_groups):
                        col_idx = i % num_cols
                        with company_cols[col_idx]:
                            if st.checkbox(group, key=f'pg_company_{group}'):
                                selected_companies.append(group)
            with filter_tabs[2]:
                pg_product_select_all = st.checkbox("Select All Product Groups", value=True, key='pg_product_all')
                selected_products = []
                if pg_product_select_all:
                    selected_products = all_product_groups
                else:
                    num_cols = 3
                    product_cols = st.columns(num_cols)
                    for i, group in enumerate(all_product_groups):
                        col_idx = i % num_cols
                        with product_cols[col_idx]:
                            if st.checkbox(group, key=f'pg_product_{group}'):
                                selected_products.append(group)
            if st.button("Generate Product Growth Report", key='pg_generate'):
                with st.spinner("Generating report..."):
                    month_title = f"LY: {selected_ly_month} vs CY: {selected_cy_month}"
                    group_results = calculate_product_growth(
                        ly_df, cy_df, budget_df, selected_ly_month, selected_cy_month,
                        ly_date_col, cy_date_col, ly_qty_col, cy_qty_col, ly_value_col, cy_value_col,
                        budget_qty_col, budget_value_col, ly_company_group_col,
                        cy_company_group_col, budget_company_group_col,
                        ly_product_group_col, cy_product_group_col, budget_product_group_col,
                        ly_sl_code_col, cy_sl_code_col, budget_sl_code_col,
                        ly_exec_col, cy_exec_col, budget_exec_col,
                        selected_executives,
                        selected_companies,
                        selected_products
                    )
                    if group_results:
                        st.success("Success!")
                        dfs_info = []
                        for company, data in group_results.items():
                            st.subheader(f"Company: {company}")
                            numeric_cols_qty = ['LAST_YEAR_QTY/MT', 'BUDGET_QTY', 'CURRENT_YEAR_QTY/MT', 'GROWTH%']
                            for col in numeric_cols_qty:
                                if col in data['qty_df'].columns:
                                    data['qty_df'][col] = data['qty_df'][col].round(2)
                            numeric_cols_value = ['LAST_YEAR_VALUE/L', 'BUDGET_VALUE', 'CURRENT_YEAR_VALUE/L', 'GROWTH %']
                            for col in numeric_cols_value:
                                if col in data['value_df'].columns:
                                    data['value_df'][col] = data['value_df'][col].round(2)
                            st.write(f"**{company} - Quantity Growth**")
                            st.dataframe(data['qty_df'])
                            st.write(f"**{company} - Value Growth**")
                            st.dataframe(data['value_df'])
                            dfs_info.append({'df': data['qty_df'], 'title': f"{company} - Quantity Growth", 'percent_cols': [4]})
                            dfs_info.append({'df': data['value_df'], 'title': f"{company} - Value Growth", 'percent_cols': [4]})
                        st.session_state.product_results = dfs_info
                        ppt_buffer = create_product_growth_ppt(
                            group_results,
                            month_title,
                            st.session_state.logo_file
                        )
                        if ppt_buffer:
                            unique_id = str(uuid.uuid4())[:8]
                            st.download_button(
                                label="Download PPT",
                                data=ppt_buffer,
                                file_name=f"Product_Growth_{month_title.replace(', ', '_')}_{unique_id}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key=f"pg_download_{unique_id}"
                            )
                    else:
                        st.error("Failed to generate report. Check your data.")
         except Exception as e:
            st.error(f"Error in Product Growth: {e}")
            st.error(traceback.format_exc())
    with tabs[3]:
        st.header("Customer & OD Analysis")
        nbc_tab, od_tab = st.tabs(["Billed Customers", "OD Target"])

        with nbc_tab:
            if not st.session_state.sales_file:
                st.warning("No sales data uploaded.")
            else:
                try:
                    sales_sheets = get_excel_sheets(st.session_state.sales_file)
                    st.subheader("Sales Sheet Selection")
                    sheet_name = st.selectbox("Select Sales Sheet", sales_sheets, key='sales_sheet_nbc')
                    sales_header_row = st.number_input("Sales Header (1-based)", min_value=1, max_value=10, value=1, step=1, key='sales_header_nbc') - 1
                    sales_df = pd.read_excel(st.session_state.sales_file, sheet_name=sheet_name, header=sales_header_row)
                    columns = sales_df.columns.tolist()
                    st.subheader("Column Mapping")
                    col1, col2 = st.columns(2)
                    with col1:
                        date_col = st.selectbox(
                            "Date",
                            columns,
                            index=columns.index('Date') if 'Date' in columns else 0,
                            help="Choose the column containing dates.",
                            key='nbc_date'
                        )
                        branch_col = st.selectbox(
                            "Branch",
                            columns,
                            index=columns.index('Branch') if 'Branch' in columns else 0,
                            help="Choose the column containing branch names.",
                            key='nbc_branch'
                        )
                    with col2:
                        customer_id_col = st.selectbox(
                            "SL Code",
                            columns,
                            index=columns.index('Customer Code') if 'Customer Code' in columns else 0,
                            help="Choose the column containing customer IDs.",
                            key='nbc_customer_id'
                        )
                        executive_col = st.selectbox(
                            "Executive",
                            columns,
                            index=columns.index('Executive Name') if 'Executive Name' in columns else 0,
                            help="Choose the column containing executive names.",
                            key='nbc_executive'
                        )
                    
                    # Get available months dynamically
                    try:
                        sales_df[date_col] = pd.to_datetime(sales_df[date_col], errors='coerce', dayfirst=True, format='mixed')
                        available_months = sorted(sales_df[date_col].dt.strftime('%b %Y').dropna().unique().tolist())
                    except Exception as e:
                        st.error(f"Error processing date column: {e}")
                        available_months = []
                    
                    # Month filter
                    st.subheader("Select Sales Month")
                    if not available_months:
                        st.warning("No valid months found in the date column.")
                        selected_months = []
                    else:
                        selected_months = st.multiselect(
                            "Select Months",
                            options=available_months,
                            default=available_months,
                            key='nbc_months_filter'
                        )

                    st.subheader("Filter Options")
                    filter_tab1, filter_tab2 = st.tabs(["Branches", "Executives"])
                    
                    with filter_tab1:
                        # Extract unique branches directly from the branch column
                        raw_branches = sales_df[branch_col].dropna().astype(str).str.strip().str.upper().unique().tolist()
                        all_branches = sorted(raw_branches)
                        branch_select_all = st.checkbox("Select All Branches", value=True, key='nbc_branch_all')
                        if branch_select_all:
                            selected_branches = all_branches
                        else:
                            num_cols = 3
                            branch_cols = st.columns(num_cols)
                            selected_branches = []
                            for i, branch in enumerate(all_branches):
                                col_idx = i % num_cols
                                with branch_cols[col_idx]:
                                    if st.checkbox(branch, key=f'nbc_branch_{branch}'):
                                        selected_branches.append(branch)
                    
                    with filter_tab2:
                        all_executives = sorted(sales_df[executive_col].dropna().astype(str).str.strip().str.upper().unique().tolist())
                        exec_select_all = st.checkbox("Select All Executives", value=True, key='nbc_exec_all')
                        if exec_select_all:
                            selected_executives = all_executives
                        else:
                            num_cols = 3
                            exec_cols = st.columns(num_cols)
                            selected_executives = []
                            for i, exec_name in enumerate(all_executives):
                                col_idx = i % num_cols
                                with exec_cols[col_idx]:
                                    if st.checkbox(exec_name, key=f'nbc_exec_{exec_name}'):
                                        selected_executives.append(exec_name)
                    
                    if st.button("Generate Report", key='nbc_generate'):
                        with st.spinner("Generating report..."):
                            if not selected_months:
                                st.error("Please select at least one month.")
                            else:
                                results = create_customer_table(
                                    sales_df,
                                    date_col,
                                    branch_col,
                                    customer_id_col,
                                    executive_col,
                                    selected_months=selected_months,
                                    selected_branches=selected_branches,
                                    selected_executives=selected_executives
                                )
                                if results:
                                    st.success("Report generated successfully!")
                                    st.subheader("Results")
                                    for fy, (result_df, sorted_months) in results.items():
                                        st.write(f"**Financial Year: {fy}**")
                                        st.dataframe(result_df, use_container_width=True)
                                        title = f"NUMBER OF BILLED CUSTOMERS - {fy}"
                                        img_buffer = create_customer_table_image(result_df, title, sorted_months, fy)
                                        if img_buffer:
                                            st.image(img_buffer, use_column_width=True)
                                        
                                        # Initialize PowerPoint presentation
                                        prs = Presentation()
                                        prs.slide_width = Inches(13.33)
                                        prs.slide_height = Inches(7.5)
                                        create_title_slide(prs, title, st.session_state.logo_file)
                                        
                                        # Use the new splitting function to create multiple slides if needed
                                        slides_data = create_customer_ppt_with_splitting(result_df, title, sorted_months, fy)
                                        
                                        for slide_info in slides_data:
                                            slide_layout = prs.slide_layouts[6]
                                            slide = prs.slides.add_slide(slide_layout)
                                            create_customer_ppt_slide(slide, slide_info['df'], slide_info['title'], sorted_months, slide_info['is_last'])
                                        
                                        ppt_buffer = BytesIO()
                                        prs.save(ppt_buffer)
                                        ppt_buffer.seek(0)
                                        st.download_button(
                                            label="Download Billed Customers PPT",
                                            data=ppt_buffer,
                                            file_name=f"Billed_customers_{fy}.pptx",
                                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                            key=f'nbc_download_{fy}'
                                        )
                                        st.session_state.customers_results = [
                                            {'df': result_df, 'title': f"NUMBER OF BILLED CUSTOMERS - {fy}"}
                                        ]
                                else:
                                    st.error("Failed to generate report. Check your data.")
                except Exception as e:
                    st.error(f"Error in tab: {e}")
                    st.error(traceback.format_exc())
        with od_tab:
            os_file_choice = st.radio(
                "Choose OS file for OD Target calculation",
                ["OS-Previous Month", "OS-Current Month"],
                key="od_file_choice"
            )
            chosen_os_file = st.session_state.os_jan_file if os_file_choice == "OS-Previous Month" else st.session_state.os_feb_file
            if not chosen_os_file:
                st.warning("⚠️ No OS file selected.")
            else:
                try:
                    os_sheets = get_excel_sheets(chosen_os_file)
                    st.subheader("Select Sheet")
                    os_sheet = st.selectbox("Select OS Sheet", os_sheets, key='od_sheet')
                    header_row = st.number_input("Header Row (1-based)", min_value=1, max_value=10, value=1, step=1, key='od_header_row') - 1
                    os_df = pd.read_excel(chosen_os_file, sheet_name=os_sheet, header=header_row)
                    if st.checkbox("Preview Raw OS Data"):
                        st.write("Raw OS Data (first 20):")
                        st.dataframe(os_df.head(20))
                    columns = os_df.columns.tolist()
                    st.subheader("OS Column Mapping")
                    col1, col2 = st.columns(2)
                    with col1:
                        os_area_col = st.selectbox(
                            "Select Area",
                            columns,
                            index=columns.index('Branch') if 'Branch' in columns else 0,
                            help="Contains branch names",
                            key='os_area_col'
                        )
                        os_qty_col = st.selectbox(
                            "Select Net Value",
                            columns,
                            index=columns.index('Net Value') if 'Net Value' in columns else 0,
                            help="Contains net values",
                            key='os_qty_col'
                        )
                    with col2:
                        os_due_date_col = st.selectbox(
                            "Select Due Date",
                            columns,
                            index=columns.index('Due Date') if 'Due Date' in columns else 0,
                            help="Contains due dates",
                            key='os_due_date_col'
                        )
                        os_exec_col = st.selectbox(
                            "Select Executive Column",
                            columns,
                            index=columns.index('Executive Name') if 'Executive Name' in columns else 0,
                            key='os_exec'
                        )
                    st.subheader("Due Date Filter")
                    try:
                        os_df[os_due_date_col] = pd.to_datetime(os_df[os_due_date_col], errors='coerce')
                        years = sorted(os_df[os_due_date_col].dt.year.dropna().astype(int).unique().tolist())
                    except Exception as e:
                        st.error(f"Error processing due dates: {e}. Ensure valid date format.")
                        years = []
                    if not years:
                        st.warning("No valid due dates found in the dataset.")
                    else:
                        selected_years = st.multiselect(
                            "Select years for filtering",
                            options=[str(year) for year in years],
                            default=[str(year) for year in years],
                            key='year_multiselect'
                        )
                        if not selected_years:
                            st.error("Please select at least one year.")
                        else:
                            month_options = ['January', 'February', 'March', 'April', 'May', 'June',
                                            'July', 'August', 'September', 'October', 'November', 'December']
                            till_month = st.selectbox("Select Month", month_options, key='till_month')

                    st.subheader("Filter Options")
                    filter_tabs = st.tabs(["Branches", "Executives"])
                    with filter_tabs[0]:
                        os_branches = sorted(set([b for b in os_df[os_area_col].apply(extract_area_name).dropna().unique() if b]))
                        if not os_branches:
                            st.error("No valid branches found in OS data.")
                        else:
                            os_branch_select_all = st.checkbox("Select All OS Branches", value=True, key='od_branch_all')
                            selected_os_branches = []
                            if os_branch_select_all:
                                selected_os_branches = os_branches
                            else:
                                num_cols = 3
                                branch_cols = st.columns(num_cols)
                                for i, branch in enumerate(os_branches):
                                    col_idx = i % num_cols
                                    with branch_cols[col_idx]:
                                        if st.checkbox(branch, key=f'od_branch_{branch}'):
                                            selected_os_branches.append(branch)
                    with filter_tabs[1]:
                        os_executives = sorted(os_df[os_exec_col].dropna().astype(str).unique().tolist())
                        os_exec_select_all = st.checkbox("Select All Executives", value=True, key='od_exec_all')
                        selected_os_executives = []
                        if os_exec_select_all:
                            selected_os_executives = os_executives
                        else:
                            num_cols = 3
                            exec_cols = st.columns(num_cols)
                            for i, exec_name in enumerate(os_executives):
                                col_idx = i % num_cols
                                with exec_cols[col_idx]:
                                    if st.checkbox(exec_name, key=f'od_exec_{exec_name}'):
                                        selected_os_executives.append(exec_name)
                    if st.button("Generate Report", key='od_generate'):
                        if not selected_years or not till_month:
                            st.error("Please select at least one year and one month.")
                        else:

                            if 'od_results' in st.session_state:
                                st.session_state.od_results = []
                            
                            with st.spinner("Generating report..."):
                                od_target_df, start_date, end_date = filter_os_qty(
                                    os_df, os_area_col, os_qty_col, os_due_date_col, os_exec_col,
                                    selected_branches=selected_os_branches,
                                    selected_years=selected_years,
                                    till_month=till_month,
                                    selected_executives=selected_os_executives
                                )
                                if od_target_df is not None:
                                    start_str = start_date.strftime('%b %Y') if start_date else 'All Periods'
                                    end_str = end_date.strftime('%b %Y') if end_date else 'All Periods'
                                    od_title = f"OD Target - {end_str}"
                                    st.subheader(od_title)
                                    st.dataframe(od_target_df)
                                    img_buffer = create_od_table_image(od_target_df, od_title)
                                    if img_buffer:
                                        st.image(img_buffer, use_column_width=True)
                                    
                                    # ✅ UPDATED POWERPOINT CREATION SECTION
                                    prs = Presentation()
                                    prs.slide_width = Inches(13.33)
                                    prs.slide_height = Inches(7.5)
                                    create_title_slide(prs, od_title, st.session_state.logo_file)
                                    
                                    # 🔄 NEW: Use the splitting function instead of single slide
                                    create_od_ppt_slides(prs, od_target_df, od_title)
                                    
                                    # Continue with the rest of PowerPoint processing
                                    ppt_buffer = BytesIO()
                                    prs.save(ppt_buffer)
                                    ppt_buffer.seek(0)
                                    st.download_button(
                                        label="Download OD Target PPT",
                                        data=ppt_buffer,
                                        file_name=f"OD_target_by_executive_{end_str}.pptx",
                                        mime="application",
                                        key='od_target_download'
                                    )
                                    if 'od_results' not in st.session_state:
                                        st.session_state['od_results'] = []
                                    st.session_state.od_results.append({'df': od_target_df, 'title': od_title})
                                else:
                                    st.error("Failed to generate report. Please check your data and selections.")
                except Exception as e:
                    st.error(f"Error in tab: {e}")
                    st.error(traceback.format_exc())
   # Summary Report Generator
    st.divider()
    st.header("Summary Report Generator")
    all_dfs_info = []
    collected_sections = []
    if hasattr(st.session_state, 'budget_results') and st.session_state.budget_results:
        all_dfs_info.extend(st.session_state.budget_results)
        collected_sections.append(f"Budget: {len(st.session_state.budget_results)} reports")
    if hasattr(st.session_state, 'od_vs_results') and st.session_state.od_vs_results:
        all_dfs_info.extend(st.session_state.od_vs_results)
        collected_sections.append(f"OD vs Collection: {len(st.session_state.od_vs_results)} reports")
    if hasattr(st.session_state, 'product_results') and st.session_state.product_results:
        all_dfs_info.extend(st.session_state.product_results)
        collected_sections.append(f"Product: {len(st.session_state.product_results)} reports")
    if hasattr(st.session_state, 'customers_results') and st.session_state.customers_results:
        all_dfs_info.extend(st.session_state.customers_results)
        collected_sections.append(f"Customers: {len(st.session_state.customers_results)} reports")
    if hasattr(st.session_state, 'od_results') and st.session_state.od_results:
        all_dfs_info.extend(st.session_state.od_results)
        collected_sections.append(f"OD Target: {len(st.session_state.od_results)} reports")
    if all_dfs_info:
        st.info(f"Reports collected: {', '.join(collected_sections)}")
        title = st.text_input("Enter Consolidated Report Title", "ACCLP Integrated Report")
        if st.button("Generate Consolidated PPT"):
            with st.spinner("Creating consolidated PowerPoint..."):
                consolidated_ppt = create_consolidated_ppt(
                    all_dfs_info,
                    st.session_state.logo_file,
                    title
                )
                if consolidated_ppt:
                    unique_id = str(uuid.uuid4())[:8]
                    st.success(f"PowerPoint created successfully with {len(all_dfs_info)} slides!")
                    st.download_button(
                        label="Download Consolidated PPT",
                        data=consolidated_ppt,
                        file_name=f"ACCLP_Consolidated_Report_{unique_id}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"consolidated_download_{unique_id}"
                    )
                else:
                    st.error("Failed to create consolidated PowerPoint. Please check the reports data.")
    else:
       st.info("No reports generated yet.")

if __name__ == "__main__":
   main()
