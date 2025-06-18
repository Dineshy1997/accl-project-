import pandas as pd
import streamlit as st
import re
import os
import gc
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import numpy as np
from datetime import datetime

# Import streamlit-plotly-events for interactivity
from streamlit_plotly_events import plotly_events

# Force matplotlib to use Agg backend for cloud deployment
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# Set environment variables for better rendering
os.environ['MPLBACKEND'] = 'Agg'

# Configure matplotlib for high-quality output on cloud
plt.rcParams.update({
    'figure.dpi': 200,
    'savefig.dpi': 300,
    'savefig.bbox': 'tight',
    'savefig.facecolor': 'white',
    'figure.facecolor': 'white',
    'axes.facecolor': 'white',
    'font.size': 14,
    'axes.labelsize': 16,
    'axes.titlesize': 18,
    'xtick.labelsize': 14,
    'ytick.labelsize': 14,
    'legend.fontsize': 14,
    'figure.figsize': (12, 8),
    'savefig.format': 'png',
    'image.interpolation': 'bilinear'
})

# Set Streamlit page configuration for wide layout with cloud optimization
st.set_page_config(
    layout='wide',
    page_title="📊 Excel Dashboard - Data Table & Visualizations",
    page_icon="📊",
    initial_sidebar_state="expanded"
)

# Add cloud-optimized CSS
def add_cloud_css():
    st.markdown("""
    <style>
        /* Improve overall appearance */
        .main .block-container {
            padding-top: 1rem;
            padding-bottom: 1rem;
            max-width: 95%;
        }
        
        /* Force high-quality image rendering */
        img {
            image-rendering: -webkit-optimize-contrast;
            image-rendering: optimize-contrast;
            image-rendering: crisp-edges;
            -ms-interpolation-mode: nearest-neighbor;
        }
        
        /* Ensure charts scale properly */
        .js-plotly-plot, .plotly {
            width: 100% !important;
            height: auto !important;
        }
        
        /* Better chart containers */
        .element-container {
            margin: 0.5rem 0;
        }
        
        /* Improve dataframe styling */
        .dataframe {
            font-size: 12px;
        }
        
        /* Better metric styling */
        [data-testid="metric-container"] {
            background-color: #f0f2f6;
            border: 1px solid #e0e0e0;
            padding: 0.5rem;
            border-radius: 0.5rem;
            margin: 0.25rem;
        }
        
        /* Improve tabs */
        .stTabs [data-baseweb="tab-list"] {
            gap: 4px;
        }
        
        .stTabs [data-baseweb="tab"] {
            height: 40px;
            padding-left: 15px;
            padding-right: 15px;
            background-color: #f0f2f6;
            border-radius: 8px 8px 0 0;
            font-size: 14px;
        }
        
        /* Better expander styling */
        .streamlit-expanderHeader {
            background-color: #f0f2f6;
            border-radius: 5px;
            font-size: 14px;
        }
        
        /* Interactive chart styling */
        .interactive-chart {
            border: 2px solid #e0e0e0;
            border-radius: 10px;
            padding: 10px;
            margin: 10px 0;
        }
        
        /* Selected data styling */
        .selected-data {
            background-color: #f0f8ff;
            border: 1px solid #4CAF50;
            border-radius: 5px;
            padding: 10px;
            margin: 10px 0;
        }
        
        /* Better responsive design */
        @media (max-width: 768px) {
            .main .block-container {
                padding: 0.5rem;
            }
            .stTabs [data-baseweb="tab"] {
                height: 35px;
                padding: 0 10px;
                font-size: 12px;
            }
        }
    </style>
    """, unsafe_allow_html=True)

add_cloud_css()

st.title("📊 Excel Dashboard - Data Table & Visualizations")

# Define exclusion terms for branches
BRANCH_EXCLUDE_TERMS = ['Total', 'TOTAL', 'Grand', 'GRAND', 'CHN Total', 'ERD SALES', 'North Total', 'WEST SALES', 'GROUP COMPANIES']

# Memory management for cloud deployment
def optimize_memory():
    """Clear memory periodically to prevent issues on cloud."""
    gc.collect()
    plt.close('all')

# Utility function to safely convert values to JSON-serializable types
def safe_convert_value(x):
    """Ultra-safe value conversion that handles all pandas types."""
    try:
        if x is None or (hasattr(x, 'isna') and pd.isna(x)) or pd.isna(x):
            return None
        str_val = str(x)
        if str_val.lower() in ['nan', 'none', 'nat', '', 'null']:
            return None
        return str_val.strip()
    except:
        return None

# Convert DataFrame to JSON-serializable types while preserving numerics
def make_jsonly_serializable(df):
    """Convert DataFrame columns to JSON-serializable types while preserving numerics."""
    if df.empty:
        return df
    df = df.copy()
    for col in df.columns:
        try:
            if pd.api.types.is_numeric_dtype(df[col]):
                if pd.api.types.is_integer_dtype(df[col]):
                    df[col] = df[col].astype('Int64')
                else:
                    df[col] = df[col].astype(float)
            else:
                df[col] = [safe_convert_value(val) for val in df[col]]
        except Exception as e:
            st.warning(f"Error processing column '{col}': {e}")
            df[col] = [str(val) if val is not None else None for val in df[col]]
    return df.reset_index(drop=True)

# Find table end by detecting TOTAL SALES or GRAND TOTAL rows
def find_table_end(df, start_idx):
    """Find where table ends by looking for TOTAL SALES or GRAND TOTAL rows."""
    for i in range(start_idx, len(df)):
        row_text = ' '.join(str(cell) for cell in df.iloc[i].values if pd.notna(cell)).upper()
        if any(term in row_text for term in ['TOTAL SALES', 'GRAND TOTAL', 'OVERALL TOTAL']):
            return i + 1
    return len(df)

# Create charts optimized for Streamlit Cloud
def create_cloud_optimized_chart(data, x_col, y_col, chart_type='bar', title="Chart", color_override=None):
    """Create charts optimized for Streamlit Cloud rendering with wider bars."""
    
    # Create figure with optimal settings for cloud - wider figure
    fig, ax = plt.subplots(figsize=(16, 8), dpi=200)  # Increased width
    
    # Set high-quality rendering
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')
    
    # Determine color - use override if provided, otherwise default colors
    if color_override:
        bar_color = color_override
    else:
        bar_color = '#2E86AB'
    
    if chart_type == 'bar':
        # Wider bars with more spacing
        bars = ax.bar(data[x_col], data[y_col], 
                     color=bar_color, alpha=0.8, 
                     edgecolor='#1B4965', linewidth=1,
                     width=0.6)  # Increased bar width
    
    elif chart_type == 'line':
        line_color = color_override if color_override else '#2E86AB'
        ax.plot(data[x_col], data[y_col], 
               marker='o', linewidth=3, markersize=8,
               color=line_color, markerfacecolor='#F18F01', 
               markeredgecolor='#1B4965', markeredgewidth=1)
        ax.grid(True, alpha=0.3, linestyle='--', linewidth=0.5)
    
    elif chart_type == 'pie':
        colors = plt.cm.Set3(np.linspace(0, 1, len(data)))
        wedges, texts, autotexts = ax.pie(data[y_col], labels=data[x_col], 
                                         autopct='%1.1f%%', startangle=90, 
                                         colors=colors, textprops={'fontsize': 14, 'weight': 'bold'})
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
            autotext.set_fontsize(14)
    
    # Enhanced styling with original sizes
    ax.set_title(title, fontsize=20, fontweight='bold', pad=15)
    if chart_type != 'pie':
        ax.set_xlabel(x_col, fontsize=16, fontweight='bold', labelpad=15)
        ax.set_ylabel(y_col, fontsize=16, fontweight='bold', labelpad=15)
        
        # Better tick formatting
        ax.tick_params(axis='both', which='major', labelsize=14, width=1.5, length=6)
        
        # Format large numbers on y-axis
        if data[y_col].max() > 1000:
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'{x/1000:.0f}K' if x >= 1000 else f'{x:.0f}'))
        
        # Rotate x-axis labels if needed
        max_label_length = max(len(str(label)) for label in data[x_col])
        if max_label_length > 10:
            plt.setp(ax.get_xticklabels(), rotation=45, ha='right')
        
        # Remove spines for cleaner look
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_linewidth(1)
        ax.spines['bottom'].set_linewidth(1)
    
    # Ensure tight layout
    plt.tight_layout(pad=3.0)
    
    return fig
# Create Plotly charts optimized for cloud with interactivity
def create_plotly_chart_cloud_optimized(data, x_col, y_col, chart_type, title, color_override=None):
    """Create Plotly charts optimized for cloud deployment with interactivity and wider bars."""
    try:
        import plotly.express as px
        import plotly.graph_objects as go
        
        # Determine color - use override if provided
        default_color = color_override if color_override else '#2E86AB'
        
        # Common layout settings optimized for cloud with wider charts
        layout_config = {
            'title': {
                'text': title,
                'x': 0.5,
                'xanchor': 'center',
                'font': {'size': 20, 'family': 'Arial, sans-serif'}
            },
            'font': {'size': 14, 'family': 'Arial, sans-serif'},
            'plot_bgcolor': 'white',
            'paper_bgcolor': 'white',
            'height': 600,
            'width': 1200,  # Explicit wider width
            'margin': {'l': 80, 'r': 80, 't': 100, 'b': 120},  # Increased bottom margin
            'showlegend': True if chart_type == 'pie' else False,
            'bargap': 0.2,  # Space between bars
            'bargroupgap': 0.1  # Space between bar groups
        }
        
        if chart_type == 'bar':
            fig = px.bar(data, x=x_col, y=y_col, 
                        color_discrete_sequence=[default_color],
                        width=1200)  # Explicit width setting
            
            # Make bars wider
            fig.update_traces(width=0.6)  # Adjust bar width (0-1)
            
            # Remove value annotations for cleaner look
            fig.update_traces(texttemplate=None, textposition=None)
            
            # Update axes
            fig.update_xaxes(
                title_font={'size': 16, 'family': 'Arial, sans-serif'},
                tickfont={'size': 14},
                tickangle=45 if max(len(str(x)) for x in data[x_col]) > 10 else 0
            )
            fig.update_yaxes(
                title_font={'size': 16, 'family': 'Arial, sans-serif'},
                tickfont={'size': 14}
            )
            
        elif chart_type == 'line':
            fig = px.line(data, x=x_col, y=y_col, 
                         markers=True, color_discrete_sequence=[default_color],
                         width=1200)  # Explicit width setting
            
            fig.update_traces(
                line={'width': 4}, 
                marker={'size': 10}
            )
            
            fig.update_xaxes(
                title_font={'size': 16, 'family': 'Arial, sans-serif'},
                tickfont={'size': 14},
                tickangle=45 if max(len(str(x)) for x in data[x_col]) > 10 else 0
            )
            fig.update_yaxes(
                title_font={'size': 16, 'family': 'Arial, sans-serif'},
                tickfont={'size': 14}
            )
            
        elif chart_type == 'pie':
            fig = px.pie(data, values=y_col, names=x_col, width=1200)
            fig.update_traces(
                textposition='inside', 
                textinfo='percent+label',
                textfont={'size': 14, 'family': 'Arial, sans-serif'}
            )
        
        # Apply layout
        fig.update_layout(**layout_config)
        
        # Configure for better cloud rendering
        config = {
            'displayModeBar': True,
            'displaylogo': False,
            'modeBarButtonsToRemove': ['pan2d', 'lasso2d', 'select2d'],
            'toImageButtonOptions': {
                'format': 'png',
                'filename': 'chart',
                'height': 800,
                'width': 1200,
                'scale': 2
            }
        }
        
        return fig, config
    except ImportError:
        return None, None
# Handle click events for interactive charts
def handle_chart_click(selected_points, data, x_col, y_col, chart_title):
    """Handle click events on charts and display detailed information."""
    if selected_points and len(selected_points) > 0:
        point = selected_points[0]
        
        # Get the clicked point information
        if 'pointIndex' in point:
            idx = point['pointIndex']
            if idx < len(data):
                clicked_row = data.iloc[idx]
                
                st.markdown(f"""
                <div class="selected-data">
                <h4>📌 Selected Data Point</h4>
                <p><strong>{x_col}:</strong> {clicked_row[x_col]}</p>
                <p><strong>{y_col}:</strong> {clicked_row[y_col]:,.2f}</p>
                </div>
                """, unsafe_allow_html=True)
                
                # Show detailed data for the selected point
                with st.expander("🔍 Detailed Information"):
                    st.json({
                        "Chart": chart_title,
                        "Selected Item": str(clicked_row[x_col]),
                        "Value": float(clicked_row[y_col]),
                        "All Data": clicked_row.to_dict()
                    })

# Create PowerPoint slide with chart image (optimized)
def create_ppt_with_chart(title, chart_data, x_col, y_col, chart_type='bar', color_override=None):
    """Creates PowerPoint slide with chart image."""
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    
    # Add title to slide as plain text (not clickable)
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
    else:
        # Fallback: create a text box for title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = title
    
    # Check if chart_data is valid
    if chart_data is None or chart_data.empty:
        st.error(f"Error: No data provided for {title}.")
        error_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        ef = error_box.text_frame
        ef.text = "Error: No data available"
        ppt_bytes = BytesIO()
        ppt.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes
    
    # Check if y_col exists and contains numeric data
    if y_col not in chart_data.columns:
        st.error(f"Error: Column {y_col} not found in data for {title}.")
        error_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        ef = error_box.text_frame
        ef.text = f"Error: Column {y_col} not found"
        ppt_bytes = BytesIO()
        ppt.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes
    
    if not pd.api.types.is_numeric_dtype(chart_data[y_col]):
        st.error(f"Error: Column {y_col} is not numeric for {title}. Cannot create chart.")
        error_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        ef = error_box.text_frame
        ef.text = f"Error: No numeric data available for {y_col}"
        ppt_bytes = BytesIO()
        ppt.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes
    
    # For pie charts, filter out non-positive values
    if chart_type == 'pie':
        chart_data = chart_data[chart_data[y_col] > 0].copy()
        if chart_data.empty:
            st.warning(f"No positive values available for pie chart in {title}. Skipping chart creation.")
            error_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
            ef = error_box.text_frame
            ef.text = f"No positive values available for {y_col} pie chart"
            ppt_bytes = BytesIO()
            ppt.save(ppt_bytes)
            ppt_bytes.seek(0)
            return ppt_bytes
    
    # Create chart with optimized settings
    fig = create_cloud_optimized_chart(chart_data, x_col, y_col, chart_type, title, color_override)
    
    # Save chart to buffer
    img_buffer = BytesIO()
    fig.savefig(img_buffer, format='png', dpi=200, bbox_inches='tight')  # Reduced dpi
    plt.close(fig)
    img_buffer.seek(0)
    
    # Add chart to slide
    slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
    
    ppt_bytes = BytesIO()
    ppt.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

# Ensure y-axis column contains numeric data
def ensure_numeric_data(data, y_col):
    """Ensure the y-axis column contains numeric data."""
    if y_col not in data.columns:
        return False
    try:
        data[y_col] = pd.to_numeric(data[y_col].astype(str).str.replace(',', ''), errors='coerce')
        data.dropna(subset=[y_col], inplace=True)
    except Exception as e:
        st.warning(f"Failed to convert {y_col} to numeric: {e}")
        return False
    return not data.empty

# Helper function to safely get chart data for master PPT
def get_chart_data_for_ppt(data, label, first_col, visual_type):
    """Safely extract chart data for PPT generation with proper column handling."""
    if data is None or (isinstance(data, pd.DataFrame) and data.empty):
        return None, None, None
    
    try:
        if label == "Budget vs Actual":
            if 'Month' in data.columns and 'Value' in data.columns and 'Metric' in data.columns:
                return data, "Month", "Value"
            elif 'Metric' in data.columns and 'Value' in data.columns:
                return data, "Metric", "Value"
            else:
                return None, None, None
                
        elif label in ["Branch Performance", "Product Performance"]:
            if len(data.columns) >= 2:
                x_col = data.columns[0]
                y_col = data.columns[1]
                return data, x_col, y_col
            else:
                return None, None, None
                
        elif label in ["Branch Monthwise", "Product Monthwise"]:
            if 'Month' in data.columns and 'Value' in data.columns:
                return data, "Month", "Value"
            else:
                return None, None, None
                
        elif "YTD" in label:
            # Handle YTD data specifically
            if 'Period' in data.columns:
                return data, "Period", label.replace("YTD ", "")
            elif len(data.columns) >= 2:
                return data, data.columns[0], data.columns[1]
            else:
                return None, None, None
                
        else:
            if "Month" in data.columns:
                label_clean = label.replace(",", "").replace(" ", "")
                if label_clean in data.columns:
                    return data, "Month", label_clean
                elif "Value" in data.columns:
                    return data, "Month", "Value"
            elif "Period" in data.columns:
                label_clean = label.replace(",", "").replace(" ", "")
                if label_clean in data.columns:
                    return data, "Period", label_clean
                elif "Value" in data.columns:
                    return data, "Period", "Value"
            
        return None, None, None
        
    except Exception as e:
        st.warning(f"Error processing chart data for {label}: {e}")
        return None, None, None

# Helper function to extract clean month-year from column names
def extract_month_year(col_name):
    """Extract clean month-year format from column names, removing Gr/Ach prefixes."""
    col_str = str(col_name).strip()
    
    # Remove common prefixes that we don't want in the x-axis labels
    col_str = re.sub(r'^(Gr[-\s]*|Ach[-\s]*|Act[-\s]*|Budget[-\s]*|LY[-\s]*)', '', col_str, flags=re.IGNORECASE)
    
    # Extract month-year pattern
    month_year_match = re.search(r'(\w{3,})[-–\s]*(\d{2})', col_str, re.IGNORECASE)
    if month_year_match:
        month, year = month_year_match.groups()
        return f"{month.capitalize()}-{year}"
    
    return col_str

# Cloud-optimized visualization display function with interactivity
def display_visualization_cloud_optimized(tab, label, data, x_col, y_col, visual_type, color_override=None):
    """Display visualization optimized for Streamlit Cloud with interactivity."""
    with tab:
        if data is None or data.empty:
            st.warning(f"No data available for {label}")
            return
        
        if not ensure_numeric_data(data, y_col):
            st.warning(f"No numeric data available for {label}")
            return None
        
        # Filter out non-positive values for pie charts
        if visual_type == "Pie Chart":
            data = data[data[y_col] > 0]
            if data.empty:
                st.warning(f"No positive values available for {label} pie chart")
                return None
        
        st.markdown(f"### {label} - {table_name}")
        st.markdown("💡 **Click on chart elements to see detailed information**")
        
        # Try Plotly first (recommended for cloud) with interactivity
        try:
            import plotly.express as px
            
            chart_type_map = {
                "Bar Chart": "bar",
                "Line Chart": "line", 
                "Pie Chart": "pie"
            }
            
            fig, config = create_plotly_chart_cloud_optimized(
                data, x_col, y_col, 
                chart_type_map[visual_type], 
                f"{label} - {table_name}",
                color_override
            )
            
            if fig is not None:
                # Display with cloud-optimized config and interactivity
                st.markdown('<div class="interactive-chart">', unsafe_allow_html=True)
                
                # Use plotly_events for interactivity
                selected_points = plotly_events(
                    fig, 
                    click_event=True, 
                    hover_event=False,
                    select_event=False,
                    override_height=700,
                    key=f"{label}_{visual_type}_chart"
                )
                
                # Handle click events
                handle_chart_click(selected_points, data, x_col, y_col, f"{label} - {table_name}")
                
                st.markdown('</div>', unsafe_allow_html=True)
            else:
                raise ImportError("Plotly not available, using matplotlib")
            
        except ImportError:
            # Fallback to matplotlib with cloud optimization (no interactivity)
            st.info("Interactive features require Plotly. Displaying static chart.")
            chart_type = visual_type.lower().replace(" chart", "")
            fig = create_cloud_optimized_chart(
                data, x_col, y_col, chart_type, f"{label} - {table_name}", color_override
            )
            
            # Use higher DPI for cloud display
            st.pyplot(fig, use_container_width=True)
            plt.close(fig)  # Important: free memory
        
        # Display data table
        with st.expander("📊 View Data Table"):
            st.dataframe(data, use_container_width=True)
        
        # Memory cleanup
        optimize_memory()

# File uploader
uploaded_file = st.sidebar.file_uploader("📂 Upload Excel File", type=["xlsx"])

if uploaded_file:
    try:
        # Clear memory before processing new file
        optimize_memory()
        
        # Read the Excel file in chunks if it's large
        file_size = uploaded_file.size
        chunk_size = 1024 * 1024  # 1MB chunks
        
        if file_size > 10 * 1024 * 1024:  # If file is larger than 10MB
            st.warning("Large file detected. Processing in chunks for better performance...")
            
            # Create a temporary file
            with open("temp_upload.xlsx", "wb") as f:
                f.write(uploaded_file.getbuffer())
            
            # Process the file in chunks
            xls = pd.ExcelFile("temp_upload.xlsx")
        else:
            xls = pd.ExcelFile(uploaded_file)
            
        sheet_names = xls.sheet_names
        selected_sheet = st.sidebar.selectbox("📄 Select a Sheet", sheet_names)
        
        # Read only the first 1000 rows initially for faster loading
        df_sheet = pd.read_excel(xls, sheet_name=selected_sheet, header=None, nrows=1000)
        
        # Try alternative reading method if data structure is suboptimal
        if df_sheet.shape[1] < 10 and df_sheet.iloc[:, 0].astype(str).str.len().max() > 200:
            try:
                df_sheet_alt = pd.read_excel(xls, sheet_name=selected_sheet, header=None, engine='openpyxl', nrows=1000)
                if df_sheet_alt.shape[1] > df_sheet.shape[1]:
                    df_sheet = df_sheet_alt
                    st.info("✅ Improved data structure using alternative reading method")
            except:
                pass
                
        
        # Dynamic processing of long header rows
        if df_sheet.shape[1] < 20:
            new_data = []
            months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
            metrics = ['Budget', 'LY', 'Act', 'Gr', 'Ach', 'YTD']
            year_pattern = r'\d{2,4}(?:[-–]\d{2,4})?'
            
            for idx, row in df_sheet.iterrows():
                if pd.notna(row.iloc[0]):
                    row_text = str(row.iloc[0]).strip()
                    if any(metric in row_text for metric in metrics) or re.search(r'SALES\s*(in\s*(MT|Value|Ton[n]?age))?', row_text, re.IGNORECASE):
                        patterns = []
                        patterns.append(r'SALES\s*in\s*(MT|Value|Ton[n]?age)', re.IGNORECASE)
                        for metric in metrics:
                            for month in months:
                                patterns.append(rf'{metric}[-–\s]*{month}[-–\s]*{year_pattern}', re.IGNORECASE)
                            patterns.append(rf'{metric}[-–\s]*YTD[-–\s]*{year_pattern}\s*\([^)]*\)', re.IGNORECASE)
                            patterns.append(rf'YTD[-–\s]*{year_pattern}\s*\([^)]*\)\s*{metric}', re.IGNORECASE)
                        
                        positions = []
                        for pattern in patterns:
                            for match in re.finditer(pattern, row_text):
                                positions.append((match.start(), match.group()))
                        positions.sort()
                        parts = [item[1].strip() for item in positions]
                        
                        if len(parts) < 5:
                            parts = [part.strip() for part in row_text.split() if part.strip()]
                        
                        new_data.append(parts)
                    else:
                        new_data.append(row_text.split())
                else:
                    new_data.append([])
            
            if new_data:
                max_cols = max(len(row) for row in new_data)
                for row in new_data:
                    while len(row) < max_cols:
                        row.append(None)
                df_sheet = pd.DataFrame(new_data)
        
    except Exception as e:
        st.error(f"Error reading Excel file: {e}")
        st.stop()

    # Determine sheet type and index
    sheet_index = sheet_names.index(selected_sheet)
    is_first_sheet = sheet_index == 0
    is_sales_monthwise = 'sales analysis month wise' in selected_sheet.lower() or ('sales' in selected_sheet.lower() and 'month' in selected_sheet.lower())

    # Initialize table_name variable
    table_name = ""

    if is_first_sheet:
        st.subheader("📋 First Sheet - Table Detection")
        
        table1_start = None
        table2_start = None
        for i in range(len(df_sheet)):
            row_text = ' '.join(str(cell) for cell in df_sheet.iloc[i].values if pd.notna(cell))
            if re.search(r'\bsales\s*in\s*mt\b', row_text, re.IGNORECASE) and table1_start is None:
                table1_start = i
            elif re.search(r'\bsales\s*in\s*(value|tonnage|tonage)\b', row_text, re.IGNORECASE) and table1_start is not None and table2_start is None:
                table2_start = i
        
        table_options = []
        if table1_start is not None:
            table_options.append("Table 1: SALES IN MT")
        if table2_start is not None:
            table_options.append("Table 2: SALES IN VALUE")
        
        if table_options:
            table_choice = st.sidebar.radio("📌 Select Table", table_options, key="first_sheet_table_select")
            table_name = table_choice
            
            if table_choice == "Table 1: SALES IN MT" and table1_start is not None:
                st.write("### Table 1: SALES IN MT")
                table1_end = table2_start if table2_start is not None else len(df_sheet)
                table1 = df_sheet.iloc[table1_start:table1_end].dropna(how='all').reset_index(drop=True)
                
                if not table1.empty:
                    header_row_idx = None
                    for i in range(min(3, len(table1))):
                        row_text = ' '.join(str(cell) for cell in table1.iloc[i].values if pd.notna(cell))
                        if re.search(r'\b(budget|ly|act|gr|ach)\b', row_text, re.IGNORECASE):
                            header_row_idx = i
                            break
                    
                    if header_row_idx is not None:
                        header_row = table1.iloc[header_row_idx]
                        new_columns = []
                        for i, val in enumerate(header_row):
                            if pd.notna(val) and str(val).strip():
                                col_name = str(val).strip()
                                if len(col_name) > 100:
                                    parts = re.split(r'\s+(?=Budget|LY|Act|Gr|Ach|YTD)', col_name)
                                    if len(parts) > 1:
                                        col_name = parts[0]
                                new_columns.append(col_name)
                            else:
                                new_columns.append(f'Unnamed_{i}')
                        
                        if len(new_columns) < len(table1.columns):
                            while len(new_columns) < len(table1.columns):
                                new_columns.append(f'Unnamed_{len(new_columns)}')
                        elif len(new_columns) > len(table1.columns):
                            new_columns = new_columns[:len(table1.columns)]
                        
                        if pd.notna(header_row.iloc[0]) and len(str(header_row.iloc[0])) > 100:
                            header_text = str(header_row.iloc[0])
                            split_headers = re.split(r'\s+(?=Budget-|LY-|Act-|Gr-|Ach-|YTD-)', header_text)
                            split_headers = [h.strip() for h in split_headers if h.strip()]
                            if len(split_headers) > 5:
                                new_columns = split_headers
                                while len(new_columns) < len(table1.columns):
                                    new_columns.append(f'Unnamed_{len(new_columns)}')
                                new_columns = new_columns[:len(table1.columns)]
                        
                        table1.columns = new_columns
                        table1 = table1.iloc[header_row_idx + 1:].reset_index(drop=True)
                        
                        if not table1.empty:
                            table1 = make_jsonly_serializable(table1)
                            st.dataframe(table1, use_container_width=True)
                            csv1 = table1.to_csv(index=False).encode('utf-8')
                            st.download_button(
                                "⬇️ Download Table 1 as CSV", 
                                csv1, 
                                "sales_in_mt.csv", 
                                "text/csv",
                                key="download_table1_csv"
                            )
                        else:
                            st.warning("No data available for Table 1 after processing.")
                    else:
                        st.error("Could not find column headers for Table 1.")
                        st.dataframe(table1)
            
            elif table_choice == "Table 2: SALES IN VALUE" and table2_start is not None:
                st.write("### Table 2: SALES IN VALUE")
                table2 = df_sheet.iloc[table2_start:].dropna(how='all').reset_index(drop=True)
                
                if not table2.empty:
                    header_row_idx = None
                    for i in range(min(5, len(table2))):
                        row_text = ' '.join(str(cell) for cell in table2.iloc[i].values if pd.notna(cell))
                        if re.search(r'\b(budget|ly|act|gr|ach)\b', row_text, re.IGNORECASE):
                            header_row_idx = i
                            break
                    
                    if header_row_idx is not None:
                        header_row = table2.iloc[header_row_idx]
                        new_columns = []
                        for i, val in enumerate(header_row):
                            if pd.notna(val) and str(val).strip():
                                col_name = str(val).strip()
                                if len(col_name) > 100:
                                    parts = re.split(r'\s+(?=Budget|LY|Act|Gr|Ach|YTD)', col_name)
                                    if len(parts) > 1:
                                        col_name = parts[0]
                                new_columns.append(col_name)
                            else:
                                new_columns.append(f'Unnamed_{i}')
                        
                        if len(new_columns) < len(table2.columns):
                            while len(new_columns) < len(table2.columns):
                                new_columns.append(f'Unnamed_{len(new_columns)}')
                        elif len(new_columns) > len(table2.columns):
                            new_columns = new_columns[:len(table2.columns)]
                        
                        if pd.notna(header_row.iloc[0]) and len(str(header_row.iloc[0])) > 100:
                            header_text = str(header_row.iloc[0])
                            split_headers = re.split(r'\s+(?=Budget-|LY-|Act-|Gr-|Ach-|YTD-)', header_text)
                            split_headers = [h.strip() for h in split_headers if h.strip()]
                            if len(split_headers) > 5:
                                new_columns = split_headers
                                while len(new_columns) < len(table2.columns):
                                    new_columns.append(f'Unnamed_{len(new_columns)}')
                                new_columns = new_columns[:len(table2.columns)]
                        
                        table2.columns = new_columns
                        table2 = table2.iloc[header_row_idx + 1:].reset_index(drop=True)
                        
                        if not table2.empty:
                            table2 = make_jsonly_serializable(table2)
                            st.dataframe(table2, use_container_width=True)
                            csv2 = table2.to_csv(index=False).encode('utf-8')
                            st.download_button(
                                "⬇️ Download Table 2 as CSV", 
                                csv2, 
                                "sales_in_value_tonage.csv", 
                                "text/csv",
                                key="download_table2_csv"
                            )
                        else:
                            st.warning("No data available for Table 2 after processing.")
                    else:
                        st.error("Could not find column headers for Table 2.")
                        st.dataframe(table2)
                else:
                    st.warning("Table 2 is empty or contains no valid data.")
        else:
            st.warning("No tables ('SALES IN MT' or 'SALES IN VALUE/TONAGE') found in the first sheet.")
            df_sheet_clean = make_jsonly_serializable(df_sheet)
            st.dataframe(df_sheet_clean, use_container_width=True)
            csv = df_sheet_clean.to_csv(index=False).encode('utf-8')
            st.download_button(
                "⬇️ Download Raw Data as CSV", 
                csv, 
                "raw_data.csv", 
                "text/csv",
                key="download_raw_data_csv"
            )
    
    elif is_sales_monthwise:
        st.subheader(f"📋 {selected_sheet} - Sales Month Wise Analysis")
        
        table1_start = None
        table2_start = None
        for i in range(len(df_sheet)):
            row_text = ' '.join(str(cell) for cell in df_sheet.iloc[i].values if pd.notna(cell))
            if re.search(r'\bsales\s*in\s*mt\b', row_text, re.IGNORECASE) and table1_start is None:
                table1_start = i
            elif re.search(r'\bsales\s*in\s*(tonnage|tonage)\b', row_text, re.IGNORECASE) and table1_start is not None and table2_start is None:
                table2_start = i
        
        table_options = []
        if table1_start is not None:
            table_options.append("Table 1: SALES IN MT")
        if table2_start is not None:
            table_options.append("Table 2: SALES IN TONAGE")
        
        if table_options:
            table_choice = st.sidebar.radio("📌 Select Table", table_options, key="sales_monthwise_table_select")
            table_name = table_choice
            
            if table_choice == "Table 1: SALES IN MT" and table1_start is not None:
                if sheet_index >= 1 and sheet_index <= 4:
                    table1_end = find_table_end(df_sheet, table1_start)
                else:
                    table1_end = table2_start if table2_start is not None else len(df_sheet)
                
                table1 = df_sheet.iloc[table1_start:table1_end].dropna(how='all').reset_index(drop=True)
                
                if not table1.empty:
                    header_row_idx = None
                    for i in range(min(5, len(table1))):
                        row_text = ' '.join(str(cell) for cell in table1.iloc[i].values if pd.notna(cell))
                        if re.search(r'\b(budget|ly|act|gr|ach)\b', row_text, re.IGNORECASE):
                            header_row_idx = i
                            break
                    
                    if header_row_idx is not None:
                        header_row = table1.iloc[header_row_idx]
                        new_columns = [str(val).strip() if pd.notna(val) else f'Unnamed_{i}' 
                                      for i, val in enumerate(header_row)]
                        table1.columns = new_columns
                        table1 = table1.iloc[header_row_idx + 1:].reset_index(drop=True)
                        
                        if 2 <= sheet_index <= 4:
                            if not table1.empty:
                                table1 = table1.drop(index=0).reset_index(drop=True)
                            else:
                                st.warning("Table 1 is empty after processing, cannot delete first row.")
                        
                        if sheet_index == 1 and not table1.empty:
                            table1 = table1[~table1[table1.columns[0]].str.contains('REGIONS', case=False, na=False)].reset_index(drop=True)
                        
                        if not table1.empty:
                            table1 = make_jsonly_serializable(table1)
                            st.dataframe(table1, use_container_width=True)
                            csv1 = table1.to_csv(index=False).encode('utf-8')
                            st.download_button(
                                "⬇️ Download Table 1 as CSV", 
                                csv1, 
                                "sales_in_mt.csv", 
                                "text/csv",
                                key="download_table1_csv"
                            )
                        else:
                            st.warning("No data available for Table 1 after processing.")
                    else:
                        st.error("Could not find column headers for Table 1.")
                        st.dataframe(table1)
            
            elif table_choice == "Table 2: SALES IN TONAGE" and table2_start is not None:
                if sheet_index >= 1 and sheet_index <= 4:
                    table2_end = find_table_end(df_sheet, table2_start)
                    table2 = df_sheet.iloc[table2_start:table2_end].dropna(how='all').reset_index(drop=True)
                else:
                    table2 = df_sheet.iloc[table2_start:].dropna(how='all').reset_index(drop=True)
                
                if not table2.empty:
                    header_row_idx = None
                    for i in range(min(5, len(table2))):
                        row_text = ' '.join(str(cell) for cell in table2.iloc[i].values if pd.notna(cell))
                        if re.search(r'\b(budget|ly|act|gr|ach)\b', row_text, re.IGNORECASE):
                            header_row_idx = i
                            break
                    
                    if header_row_idx is not None:
                        header_row = table2.iloc[header_row_idx]
                        new_columns = [str(val).strip() if pd.notna(val) else f'Unnamed_{i}' 
                                      for i, val in enumerate(header_row)]
                        table2.columns = new_columns
                        table2 = table2.iloc[header_row_idx + 1:].reset_index(drop=True)
                        
                        if 2 <= sheet_index <= 4:
                            if not table2.empty:
                                table2 = table2.drop(index=0).reset_index(drop=True)
                            else:
                                st.warning("Table 2 is empty after processing, cannot delete first row.")
                        
                        if sheet_index == 1 and not table2.empty:
                            table2 = table2[~table2[table2.columns[0]].str.contains('REGIONS', case=False, na=False)].reset_index(drop=True)
                        
                        if not table2.empty:
                            table2 = make_jsonly_serializable(table2)
                            st.dataframe(table2, use_container_width=True)
                            csv2 = table2.to_csv(index=False).encode('utf-8')
                            st.download_button(
                                "⬇️ Download Table 2 as CSV", 
                                csv2, 
                                "sales_in_tonage.csv", 
                                "text/csv",
                                key="download_table2_csv"
                            )
                        else:
                            st.warning("No data available for Table 2 after processing.")
                    else:
                        st.error("Could not find column headers for Table 2.")
                        st.dataframe(table2)
                else:
                    st.warning("Table 2 is empty or contains no valid data.")
        else:
            st.warning("No tables ('SALES IN MT' or 'SALES IN TONAGE') found in the sheet.")
            df_sheet_clean = make_jsonly_serializable(df_sheet)
            st.dataframe(df_sheet_clean, use_container_width=True)
            csv = df_sheet_clean.to_csv(index=False).encode('utf-8')
            st.download_button(
                "⬇️ Download Raw Data as CSV", 
                csv, 
                "raw_data.csv", 
                "text/csv",
                key="download_raw_data_csv"
            )
    
    else:
        is_product_analysis = ('product' in selected_sheet.lower() or 
                             'ts-pw' in selected_sheet.lower() or 
                             'ero-pw' in selected_sheet.lower())
        is_branch_analysis = 'region wise analysis' in selected_sheet.lower()

        if is_branch_analysis:
            table1_header = "Sales in MT"
            table2_header = "Sales in Value"
        else:
            table1_header = "Sales in Tonage"
            table2_header = "Sales in Value"

        def extract_tables(df, table1_header, table2_header):
            table1_idx, table2_idx = None, None
            for i in range(len(df)):
                row_text = ' '.join(df.iloc[i].astype(str).str.lower().tolist())
                if table1_idx is None and table1_header.lower() in row_text:
                    table1_idx = i
                elif table2_idx is None and table2_header.lower() in row_text and i > (table1_idx or 0):
                    table2_idx = i
            return table1_idx, table2_idx

        idx1, idx2 = extract_tables(df_sheet, table1_header, table2_header)

        if idx1 is None:
            st.error(f"❌ Could not locate '{table1_header}' header in the sheet.")
        else:
            if sheet_index >= 1 and sheet_index <= 4:
                table1_end = find_table_end(df_sheet, idx1 + 1)
                if idx2 is not None:
                    table1_end = min(table1_end, idx2)
            else:
                table1_end = idx2 if idx2 is not None else len(df_sheet)
            
            table1 = df_sheet.iloc[idx1+1:table1_end].dropna(how='all').reset_index(drop=True)
            table1.columns = df_sheet.iloc[idx1].apply(lambda x: str(x).strip() if pd.notna(x) else '')
            
            if not table1.empty:
                first_row = table1.iloc[0]
                first_cell = str(first_row[0]).strip().lower() if pd.notna(first_row[0]) else ""
                is_subheader = (first_cell == "" or 
                               first_cell.startswith("unnamed") or 
                               any(term.lower() in first_cell for term in BRANCH_EXCLUDE_TERMS))
                if is_subheader:
                    table1 = table1.drop(index=0).reset_index(drop=True)

            if idx2 is not None:
                if sheet_index >= 1 and sheet_index <= 4:
                    table2_end = find_table_end(df_sheet, idx2 + 1)
                    table2 = df_sheet.iloc[idx2+1:table2_end].dropna(how='all').reset_index(drop=True)
                else:
                    table2 = df_sheet.iloc[idx2+1:].dropna(how='all').reset_index(drop=True)
                
                table2.columns = df_sheet.iloc[idx2].apply(lambda x: str(x).strip() if pd.notna(x) else '')
                
                if not table2.empty:
                    first_row = table2.iloc[0]
                    first_cell = str(first_row[0]).strip().lower() if pd.notna(first_row[0]) else ""
                    is_subheader = (first_cell == "" or 
                                   first_cell.startswith("unnamed") or 
                                   any(term.lower() in first_cell for term in BRANCH_EXCLUDE_TERMS))
                    if is_subheader:
                        table2 = table2.drop(index=0).reset_index(drop=True)
            else:
                table2 = None

            table_options = [f"Table 1: {table1_header.upper()}"]
            if table2 is not None:
                table_options.append(f"Table 2: {table2_header.upper()}")
            table_choice = st.sidebar.radio("📌 Select Table", table_options)
            table_name = table_choice
            table_df = table1 if table_choice == table_options[0] else table2

            table_df = make_jsonly_serializable(table_df)
            table_df.columns = table_df.columns.map(str)

            def rename_columns(columns):
                renamed = []
                ytd_base = None
                prev_month = None
                prev_year = None
                for col in columns:
                    col_clean = col.strip().replace(",", "").replace("–", "-")
                    ytd_act_match = re.search(r'(YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))\s*Act', col_clean, re.IGNORECASE)
                    if ytd_act_match:
                        ytd_base = ytd_act_match.group(1).replace("–", "-")
                        renamed.append(f"Act-{ytd_base}")
                        continue
                    ytd_gr_match = re.search(r'(Gr-YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))', col_clean, re.IGNORECASE)
                    if ytd_gr_match:
                        gr_ytd = ytd_gr_match.group(1).replace("–", "-")
                        renamed.append(f"Gr-{gr_ytd.split('Gr-')[1]}")
                        continue
                    ytd_ach_match = re.search(r'(Ach-YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))', col_clean, re.IGNORECASE)
                    ytd_ach_alt_match = re.search(r'(YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))\s*Ach', col_clean, re.IGNORECASE)
                    if ytd_ach_match:
                        ach_ytd = ytd_ach_match.group(1).replace("–", "-")
                        renamed.append(ach_ytd)
                        continue
                    elif ytd_ach_alt_match:
                        ytd_part = ytd_ach_alt_match.group(1).replace("–", "-")
                        renamed.append(f"Ach-{ytd_part}")
                        continue
                    monthly_match = re.search(r'(\b\w{3,})[\s-]*(\d{2})', col_clean)
                    if monthly_match:
                        prev_month, prev_year = monthly_match.groups()
                        prev_month = prev_month.capitalize()
                    if col_clean.lower().startswith("gr") and prev_month and prev_year:
                        renamed.append(f"Gr - {prev_month}-{prev_year}")
                    elif col_clean.lower().startswith("ach") and prev_month and prev_year:
                        renamed.append(f"Ach - {prev_month}-{prev_year}")
                    else:
                        renamed.append(col)
                return renamed

            table_df.columns = rename_columns(table_df.columns)

            if table_df.columns.duplicated().any():
                table_df = table_df.loc[:, ~table_df.columns.duplicated()]
                st.warning("⚠️ Duplicate columns detected and removed.")

            branch_list = []
            product_list = []

            def extract_unique_values(df, first_col, exclude_terms=None):
                if exclude_terms is None:
                    exclude_terms = ['Total', 'TOTAL', 'Grand', 'GRAND', 'CHN Total', 'ERD SALES']
                
                valid_rows = df[~df[first_col].str.contains('|'.join(exclude_terms), na=False, case=False)]
                valid_rows = valid_rows[valid_rows[first_col].notna()]
                
                unique_values = valid_rows[first_col].astype(str).str.strip().unique()
                return sorted(unique_values)

            if is_branch_analysis:
                branch_list = extract_unique_values(table_df, table_df.columns[0])
            elif is_product_analysis:
                product_list = extract_unique_values(table_df, table_df.columns[0])

            months = sorted(set(re.findall(r'\b(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\b', 
                       ' '.join(table_df.columns), flags=re.IGNORECASE)))
            years = sorted(set(re.findall(r'[-–](\d{2})\b', ' '.join(table_df.columns))))

            # Add "Select All" option to filters
            months_options = ["Select All"] + months
            years_options = ["Select All"] + years
            branches_options = ["Select All"] + branch_list if is_branch_analysis else []
            products_options = ["Select All"] + product_list if is_product_analysis else []

            selected_month = st.sidebar.selectbox("📅 Filter by Month", months_options, index=0)
            selected_year = st.sidebar.selectbox("📆 Filter by Year", years_options, index=0)
            selected_branch = st.sidebar.selectbox("🌍 Filter by Branch", branches_options, index=0) if is_branch_analysis else None
            selected_product = st.sidebar.selectbox("📦 Filter by Product", products_options, index=0) if is_product_analysis else None

            # Handle "Select All" logic
            selected_months = months if selected_month == "Select All" else [selected_month] if selected_month else []
            selected_years = years if selected_year == "Select All" else [selected_year] if selected_year else []
            selected_branches = branch_list if selected_branch == "Select All" else [selected_branch] if selected_branch else []
            selected_products = product_list if selected_product == "Select All" else [selected_product] if selected_product else []

            filtered_df = table_df.copy()
            first_col = filtered_df.columns[0]

            if selected_branches and is_branch_analysis:
                filtered_df = filtered_df[filtered_df[first_col].astype(str).isin(selected_branches)]
            if selected_products and is_product_analysis:
                filtered_df = filtered_df[filtered_df[first_col].astype(str).isin(selected_products)]

            def column_filter(col):
                col_str = str(col).lower().replace(",", "").replace("–", "-")
                if "ytd" in col_str:
                    return any(f"-{y}" in col_str for y in selected_years) if selected_years else True
                month_match = any(m.lower() in col_str for m in selected_months)
                year_match = any(f"-{y}" in col_str for y in selected_years) if selected_years else True
                return month_match and year_match

            visual_cols = [col for col in table_df.columns if column_filter(col)]
            display_df = filtered_df[[first_col] + visual_cols] if visual_cols else filtered_df[[first_col]]

            display_df = make_jsonly_serializable(display_df)

            def convert_to_numeric(series):
                try:
                    return pd.to_numeric(series.astype(str).str.replace(',', ''), errors='coerce')
                except:
                    return series

            formatted_df = display_df.copy()
            numeric_cols = []
            for col in formatted_df.columns:
                if col == formatted_df.columns[0]:
                    continue
                formatted_df[col] = convert_to_numeric(formatted_df[col])
                if pd.api.types.is_numeric_dtype(formatted_df[col]):
                    numeric_cols.append(col)
                    formatted_df[col] = formatted_df[col].round(2)

            style_dict = {col: "{:.2f}" for col in numeric_cols}
            st.subheader("📋 Filtered Table View")
            st.dataframe(formatted_df.style.format(style_dict, na_rep="-"), use_container_width=True)

            csv = display_df.to_csv(index=False).encode('utf-8')
            st.download_button(
                "⬇️ Download Filtered Data as CSV", 
                csv, 
                "filtered_data.csv",
                "text/csv",
                key="download_filtered_data_csv"
            )

            st.sidebar.markdown("---")
            st.sidebar.subheader("📊 Visualization Options")
            
            visual_type = st.sidebar.selectbox(
                "Select Visualization Type",
                ["Bar Chart", "Pie Chart", "Line Chart"],
                index=0,
                key="visualization_type_select"
            )  

            tabs = st.tabs([
                "📊 Budget vs Actual", "📊 Budget", "📊 LY", "📊 Act", "📊 Gr", "📊 Ach", 
                "📈 YTD Budget", "📈 YTD LY", "📈 YTD Act", "📈 YTD Gr", "📈 YTD Ach", 
                "🌍 Branch Performance", "🌍 Branch Monthwise", 
                "📦 Product Performance", "📦 Product Monthwise"
            ])

            tab_names = [
                "Budget vs Actual", "Budget", "LY", "Act", "Gr", "Ach",
                "YTD Budget", "YTD LY", "YTD Act", "YTD Gr", "YTD Ach",
                "Branch Performance", "Branch Monthwise",
                "Product Performance", "Product Monthwise"
            ]
            tabs_dict = dict(zip(tab_names, tabs))

            def plot_budget_vs_actual(tab, visual_type):
                with tab:
                    budget_cols = [col for col in table_df.columns 
                                  if str(col).lower().startswith('budget') and 'ytd' not in str(col).lower()
                                  and column_filter(col)]
                    act_cols = [col for col in table_df.columns 
                               if str(col).lower().startswith('act') and 'ytd' not in str(col).lower()
                               and column_filter(col)]
                
                    if not (budget_cols and act_cols):
                        st.info("No matching Budget or Act columns found for comparison")
                        return None
                
                    budget_months = [re.search(r'(\w{3,})[-–](\d{2})', str(col), re.IGNORECASE) for col in budget_cols]
                    act_months = [re.search(r'(\w{3,})[-–](\d{2})', str(col), re.IGNORECASE) for col in act_cols]
                    common_months = set((m.group(1), m.group(2)) for m in budget_months if m) & \
                                    set((m.group(1), m.group(2)) for m in act_months if m)
                
                    if not common_months:
                        st.info("No common months found for Budget vs Actual comparison")
                        return None
                
                    selected_budget_cols = []
                    selected_act_cols = []
                    for month, year in common_months:
                        for col in budget_cols:
                            if re.search(rf'\b{month}[-–]{year}\b', str(col), re.IGNORECASE):
                                selected_budget_cols.append(col)
                        for col in act_cols:
                            if re.search(rf'\b{month}[-–]{year}\b', str(col), re.IGNORECASE):
                                selected_act_cols.append(col)
                
                    chart_data = filtered_df[[first_col] + selected_budget_cols + selected_act_cols].copy()
                
                    for col in selected_budget_cols + selected_act_cols:
                        chart_data[col] = pd.to_numeric(chart_data[col].astype(str).str.replace(',', ''), 
                                                       errors='coerce')
                
                    chart_data = chart_data.dropna()
                
                    if chart_data.empty:
                        st.warning("No valid numeric data available for Budget vs Act comparison")
                        return None
                
                    # Prepare data for pie chart if selected
                    if visual_type == "Pie Chart":
                        # Aggregate totals for Budget and Actual
                        budget_total = chart_data[selected_budget_cols].sum().sum()
                        act_total = chart_data[selected_act_cols].sum().sum()
                        pie_data = pd.DataFrame({
                            "Metric": ["Budget", "Act"],
                            "Value": [budget_total, act_total]
                        })
                        pie_data = pie_data[pie_data["Value"] > 0]  # Remove zero or negative values
                        if pie_data.empty:
                            st.warning("No valid data for Budget vs Actual pie chart after aggregation")
                            return None
                        display_visualization_cloud_optimized(tab, "Budget vs Actual", pie_data, "Metric", "Value", visual_type)
                        ppt_type = 'pie'
                        chart_data_for_ppt = pie_data
                        x_col_for_ppt = "Metric"
                    else:
                        # Existing grouped bar chart or line chart logic
                        chart_data_melt = chart_data.melt(id_vars=first_col, 
                                                         var_name="Month_Metric", 
                                                         value_name="Value")
                        chart_data_melt['Metric'] = chart_data_melt['Month_Metric'].apply(
                            lambda x: 'Budget' if 'budget' in str(x).lower() else 'Act'
                        )
                        chart_data_melt['Month'] = chart_data_melt['Month_Metric'].apply(
                            lambda x: re.search(r'(\w{3,})[-–](\d{2})', str(x), re.IGNORECASE).group(0) 
                                      if re.search(r'(\w{3,})[-–](\d{2})', str(x), re.IGNORECASE) else x
                        )
                        
                        chart_data_melt = make_jsonly_serializable(chart_data_melt)
                        chart_data_agg = chart_data_melt.groupby(['Month', 'Metric'])['Value'].sum().reset_index()
                        
                        if chart_data_agg.empty or 'Value' not in chart_data_agg.columns:
                            st.warning("Aggregation failed: No valid data for Budget vs Actual comparison")
                            return None
                        
                        chart_data_agg['Value'] = pd.to_numeric(chart_data_agg['Value'], errors='coerce')
                        if chart_data_agg['Value'].isna().all():
                            st.warning("No numeric values available in aggregated data for Budget vs Actual")
                            return None
                        
                        if not ensure_numeric_data(chart_data_agg, 'Value'):
                            st.warning("No numeric data available for Budget vs Actual comparison")
                            return None
                        
                        st.markdown(f"### Budget vs Actual Comparison - {table_name}")
                        st.markdown("💡 **Click on chart elements to see detailed information**")
                        
                        try:
                            import plotly.express as px
                            import plotly.graph_objects as go
                            
                            fig = go.Figure()
                            budget_data = chart_data_agg[chart_data_agg['Metric'] == 'Budget']
                            act_data = chart_data_agg[chart_data_agg['Metric'] == 'Act']
                            
                            if not budget_data.empty:
                                fig.add_trace(go.Bar(
                                    x=budget_data['Month'],
                                    y=budget_data['Value'],
                                    name='Budget',
                                    marker_color='#2E86AB'
                                ))
                            
                            if not act_data.empty:
                                fig.add_trace(go.Bar(
                                    x=act_data['Month'],
                                    y=act_data['Value'],
                                    name='Act',
                                    marker_color='#FF8C00'
                                ))
                            
                            fig.update_layout(
                                title={'text': f"Budget vs Actual Comparison - {table_name}", 'x': 0.5, 'xanchor': 'center', 'font': {'size': 16}},
                                xaxis_title="Month",
                                yaxis_title="Value",
                                font={'size': 12},
                                plot_bgcolor='white',
                                paper_bgcolor='white',
                                height=600,
                                margin={'l': 60, 'r': 60, 't': 80, 'b': 60},
                                showlegend=True,
                                barmode='group'
                            )
                            fig.update_xaxes(title_font={'size': 14}, tickfont={'size': 12})
                            fig.update_yaxes(title_font={'size': 14}, tickfont={'size': 12})
                            
                            # Add interactivity
                            st.markdown('<div class="interactive-chart">', unsafe_allow_html=True)
                            selected_points = plotly_events(
                                fig, 
                                click_event=True, 
                                hover_event=False,
                                select_event=False,
                                override_height=600,
                                key=f"budget_vs_actual_{visual_type}_chart"
                            )
                            
                            # Handle click events for grouped data
                            if selected_points and len(selected_points) > 0:
                                point = selected_points[0]
                                if 'pointIndex' in point and 'curveNumber' in point:
                                    curve_num = point['curveNumber']
                                    point_idx = point['pointIndex']
                                    
                                    if curve_num == 0 and not budget_data.empty and point_idx < len(budget_data):
                                        clicked_row = budget_data.iloc[point_idx]
                                        metric_type = "Budget"
                                    elif curve_num == 1 and not act_data.empty and point_idx < len(act_data):
                                        clicked_row = act_data.iloc[point_idx]
                                        metric_type = "Act"
                                    else:
                                        clicked_row = None
                                        metric_type = "Unknown"
                                    
                                    if clicked_row is not None:
                                        st.markdown(f"""
                                        <div class="selected-data">
                                        <h4>📌 Selected Data Point</h4>
                                        <p><strong>Metric:</strong> {metric_type}</p>
                                        <p><strong>Month:</strong> {clicked_row['Month']}</p>
                                        <p><strong>Value:</strong> {clicked_row['Value']:,.2f}</p>
                                        </div>
                                        """, unsafe_allow_html=True)
                            
                            st.markdown('</div>', unsafe_allow_html=True)
                            
                        except ImportError:
                            fig, ax = plt.subplots(figsize=(10, 6), dpi=150)
                            fig.patch.set_facecolor('white')
                            ax.set_facecolor('white')
                            budget_data = chart_data_agg[chart_data_agg['Metric'] == 'Budget']
                            act_data = chart_data_agg[chart_data_agg['Metric'] == 'Act']
                            if not budget_data.empty and not act_data.empty:
                                bar_width = 0.35
                                months = sorted(set(chart_data_agg['Month']))
                                x_pos = np.arange(len(months))
                                budget_values = [budget_data[budget_data['Month'] == month]['Value'].iloc[0] 
                                               if len(budget_data[budget_data['Month'] == month]) > 0 else 0 for month in months]
                                act_values = [act_data[act_data['Month'] == month]['Value'].iloc[0] 
                                            if len(act_data[act_data['Month'] == month]) > 0 else 0 for month in months]
                                ax.bar(x_pos - bar_width/2, budget_values, bar_width, label='Budget', color='#2E86AB')
                                ax.bar(x_pos + bar_width/2, act_values, bar_width, label='Act', color='#FF8C00')
                                ax.set_xlabel('Month', fontsize=14, fontweight='bold', labelpad=10)
                                ax.set_ylabel('Value', fontsize=14, fontweight='bold', labelpad=10)
                                ax.set_title(f"Budget vs Actual Comparison - {table_name}", fontsize=16, fontweight='bold', pad=15)
                                ax.set_xticks(x_pos)
                                ax.set_xticklabels(months, rotation=0, ha='center')
                                ax.legend()
                                ax.tick_params(axis='both', which='major', labelsize=12, width=1, length=4)
                                ax.spines['top'].set_visible(False)
                                ax.spines['right'].set_visible(False)
                                ax.spines['left'].set_linewidth(1)
                                ax.spines['bottom'].set_linewidth(1)
                            
                            plt.tight_layout(pad=2.0)
                            st.pyplot(fig, use_container_width=True)
                            plt.close(fig)
                        
                        ppt_type = 'bar' if visual_type == 'Bar Chart' else 'line'
                        chart_data_for_ppt = chart_data_agg
                        x_col_for_ppt = "Month"
                    
                    # Display data table
                    with st.expander("📊 View Data Table"):
                        st.dataframe(chart_data_for_ppt, use_container_width=True)
                
                    ppt_bytes = create_ppt_with_chart(
                        title=f"Budget vs Actual - {table_name} - {selected_sheet}",
                        chart_data=chart_data_for_ppt,
                        x_col=x_col_for_ppt,
                        y_col="Value",
                        chart_type=ppt_type
                    )
                
                    st.download_button(
                        "⬇️ Download Budget vs Actual PPT",
                        ppt_bytes,
                        "budget_vs_actual.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_budget_vs_actual_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data_for_ppt

            def plot_monthly_comparison(tab, label, visual_type):
                with tab:
                    normalized_label = label.replace(",", "")
                    plot_cols = [col for col in table_df.columns 
                               if str(col).lower().replace(",", "").startswith(normalized_label.lower()) 
                               and 'ytd' not in str(col).lower()
                               and column_filter(col)]
                
                    if not plot_cols:
                        st.info(f"No matching columns found for '{label}'")
                        return None
                
                    chart_data = filtered_df[[first_col] + plot_cols].copy()
                
                    for col in plot_cols:
                        chart_data[col] = pd.to_numeric(chart_data[col].astype(str).str.replace(',', ''), 
                                                       errors='coerce')
                
                    chart_data = chart_data.melt(id_vars=first_col, 
                                              var_name="Month", 
                                              value_name=label)
                    chart_data = chart_data.dropna()
                
                    chart_data[label] = pd.to_numeric(chart_data[label], errors='coerce')
                
                    if chart_data.empty or not ensure_numeric_data(chart_data, label):
                        st.warning(f"No valid numeric data available for '{label}' after conversion.")
                        return None
                
                    # Clean month names using helper function
                    chart_data['Month'] = chart_data['Month'].apply(extract_month_year)
                
                    month_order = {'Apr': 1, 'May': 2, 'Jun': 3, 'Jul': 4, 'Aug': 5, 'Sep': 6,
                                   'Oct': 7, 'Nov': 8, 'Dec': 9, 'Jan': 10, 'Feb': 11, 'Mar': 12}
                
                    def get_sort_key(month_str):
                        month_match = re.search(r'(\w{3,})[-–](\d{2})', month_str, re.IGNORECASE)
                        if month_match:
                            month, year = month_match.groups()
                            month_idx = month_order.get(month.capitalize(), 99)
                            year_int = int(year)
                            if month_idx >= 10:
                                fiscal_year = year_int - 1
                            else:
                                fiscal_year = year_int
                            return (fiscal_year, month_idx)
                        return (0, 99)
                
                    chart_data = chart_data.sort_values(by='Month', key=lambda x: x.map(get_sort_key))
                    chart_data = make_jsonly_serializable(chart_data)
                
                    # Set orange color for Act tab
                    color_override = '#FF8C00' if label == "Act" else None
                
                    display_visualization_cloud_optimized(tab, f"{label} by Month", chart_data, "Month", label, visual_type, color_override)
                    
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'
                    ppt_bytes = create_ppt_with_chart(
                        f"{label} Analysis - {table_name} - {selected_sheet}",
                        chart_data,
                        "Month",
                        label,
                        ppt_type,
                        color_override
                    )
                    
                    st.download_button(
                        f"⬇️ Download {label} PPT",
                        ppt_bytes,
                        f"{label.lower().replace(' ', '_')}_analysis.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_{label.lower().replace(' ', '_')}_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data

            def plot_ytd_comparison(tab, pattern, label, visual_type):
                with tab:
                    ytd_cols = []
                    normalized_label = label.replace(",", "").lower()
                    
                    for col in table_df.columns:
                        col_str = str(col).lower().replace(",", "").replace("–", "-")
                        if normalized_label == 'gr':
                            if (re.search(r'gr-ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)', col_str, re.IGNORECASE) and
                                column_filter(col)):
                                ytd_cols.append(col)
                        elif normalized_label == 'ach':
                            if (re.search(r'ach-ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)', col_str, re.IGNORECASE) and
                                column_filter(col)):
                                ytd_cols.append(col)
                            elif (re.search(r'(ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\).*ach|ach.*ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\))', col_str, re.IGNORECASE) and
                                  column_filter(col)):
                                ytd_cols.append(col)
                        else:
                            if (re.search(r'ytd.*\b' + normalized_label + r'\b|' + normalized_label + r'.*ytd', col_str, re.IGNORECASE) or
                                re.search(r'ytd[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)\s*' + normalized_label, col_str, re.IGNORECASE)) and \
                               column_filter(col):
                                ytd_cols.append(col)
                    
                    if not ytd_cols:
                        st.warning(f"No YTD {label} columns found. Expected columns like '{label}-YTD-25-26-(Apr to Jun)'.")
                        return None
                    
                    clean_labels = []
                    for col in ytd_cols:
                        col_str = str(col)
                        year_match = re.search(r'(\d{2,4})\s*[-–]\s*(\d{2,4})\s*\((.*?)\)', col_str, re.IGNORECASE)
                        if year_match:
                            start_year, end_year, month_range = year_match.groups()
                            start_year = start_year[-2:] if len(start_year) > 2 else start_year
                            end_year = end_year[-2:] if len(end_year) > 2 else end_year
                            fiscal_year = f"{start_year}-{end_year}"
                            month_range_clean = re.sub(r'\s*to\s*', ' - ', month_range, flags=re.IGNORECASE)
                            if label.lower() == 'budget':
                                clean_label = f"Budget {fiscal_year} ({month_range_clean})"
                            elif label.lower() == 'ly':
                                clean_label = f"LY {fiscal_year} ({month_range_clean})"
                            else:
                                clean_label = f"{label} {fiscal_year} ({month_range_clean})"
                        else:
                            fiscal_year = "Unknown"
                            month_range_clean = "Apr - Jun"
                            if label.lower() == 'budget':
                                clean_label = f"Budget {fiscal_year} ({month_range_clean})"
                            elif label.lower() == 'ly':
                                clean_label = f"LY {fiscal_year} ({month_range_clean})"
                            else:
                                clean_label = f"{label} {fiscal_year} ({month_range_clean})"
                            st.warning(f"Could not parse year or month range in column '{col}'. Using default '{clean_label}'.")
                        clean_labels.append(clean_label)
                    
                    month_order = {'Apr':1, 'May':2, 'Jun':3, 'Jul':4, 'Aug':5, 'Sep':6,
                                   'Oct':7, 'Nov':8, 'Dec':9, 'Jan':10, 'Feb':11, 'Mar':12}
                    
                    def get_sort_key(col_name):
                        month_match = re.search(r'\((\w{3})', col_name, re.IGNORECASE)
                        return month_order.get(month_match.group(1).capitalize(), 0) if month_match else 0
                    
                    sorted_cols = [first_col] + sorted(clean_labels, key=get_sort_key)
                    comparison_data = filtered_df[[first_col] + ytd_cols].copy()
                    comparison_data.columns = [first_col] + clean_labels
                    comparison_data = comparison_data[sorted_cols]
                    
                    for col in clean_labels:
                        comparison_data[col] = pd.to_numeric(comparison_data[col].astype(str).str.replace(',', ''), errors='coerce')
                    
                    chart_data = comparison_data.melt(id_vars=first_col, 
                                                     var_name="Period", 
                                                     value_name=label)
                    chart_data = chart_data.dropna()
                    
                    if not ensure_numeric_data(chart_data, label):
                        st.warning(f"No numeric data available for YTD {label} comparisons")
                        return None
                    
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    # Set orange color for Act YTD tab
                    color_override = '#FF8C00' if label == "Act" else None
                    
                    # Custom visualization for YTD with straight x-axis labels and interactivity
                    st.markdown(f"### {label} YTD Comparisons - {table_name}")
                    st.markdown("💡 **Click on chart elements to see detailed information**")
                    
                    # Try Plotly first with straight x-axis labels
                    try:
                        import plotly.express as px
                        
                        # Determine color
                        default_color = color_override if color_override else '#2E86AB'
                        
                        if visual_type == "Bar Chart":
                            fig = px.bar(chart_data, x="Period", y=label, 
                                        title=f"{label} YTD Comparisons - {table_name}",
                                        color_discrete_sequence=[default_color])
                            
                            # Force straight x-axis labels for YTD
                            fig.update_xaxes(
                                title_font={'size': 14, 'family': 'Arial, sans-serif'},
                                tickfont={'size': 12},
                                tickangle=0
                            )
                            fig.update_yaxes(
                                title_font={'size': 14, 'family': 'Arial, sans-serif'},
                                tickfont={'size': 12}
                            )
                            
                        elif visual_type == "Line Chart":
                            fig = px.line(chart_data, x="Period", y=label, 
                                         title=f"{label} YTD Comparisons - {table_name}",
                                         markers=True, color_discrete_sequence=[default_color])
                            
                            fig.update_traces(line={'width': 3}, marker={'size': 8})
                            fig.update_xaxes(
                                title_font={'size': 14, 'family': 'Arial, sans-serif'},
                                tickfont={'size': 12},
                                tickangle=0
                            )
                            fig.update_yaxes(
                                title_font={'size': 14, 'family': 'Arial, sans-serif'},
                                tickfont={'size': 12}
                            )
                            
                        elif visual_type == "Pie Chart":
                            fig = px.pie(chart_data, values=label, names="Period", 
                                        title=f"{label} YTD Distribution - {table_name}")
                            fig.update_traces(
                                textposition='inside', 
                                textinfo='percent+label',
                                textfont={'size': 12, 'family': 'Arial, sans-serif'}
                            )
                        
                        # Apply common layout
                        fig.update_layout(
                            title={'x': 0.5, 'xanchor': 'center', 'font': {'size': 16, 'family': 'Arial, sans-serif'}},
                            font={'size': 12, 'family': 'Arial, sans-serif'},
                            plot_bgcolor='white',
                            paper_bgcolor='white',
                            height=600,
                            margin={'l': 60, 'r': 60, 't': 80, 'b': 60}
                        )
                        
                        # Add interactivity
                        st.markdown('<div class="interactive-chart">', unsafe_allow_html=True)
                        selected_points = plotly_events(
                            fig, 
                            click_event=True, 
                            hover_event=False,
                            select_event=False,
                            override_height=600,
                            key=f"ytd_{label}_{visual_type}_chart"
                        )
                        
                        # Handle click events
                        handle_chart_click(selected_points, chart_data, "Period", label, f"{label} YTD Comparisons - {table_name}")
                        st.markdown('</div>', unsafe_allow_html=True)
                        
                    except ImportError:
                        # Fallback to matplotlib with straight x-axis labels
                        st.info("Interactive features require Plotly. Displaying static chart.")
                        fig, ax = plt.subplots(figsize=(10, 6), dpi=150)
                        fig.patch.set_facecolor('white')
                        ax.set_facecolor('white')
                        
                        # Determine color for matplotlib
                        bar_color = color_override if color_override else '#2E86AB'
                        
                        if visual_type == "Bar Chart":
                            ax.bar(chart_data['Period'], chart_data[label], 
                                  color=bar_color, alpha=0.8, edgecolor='#1B4965', linewidth=1)
                        
                        elif visual_type == "Line Chart":
                            ax.plot(chart_data['Period'], chart_data[label], 
                                   marker='o', linewidth=3, markersize=8, 
                                   color=bar_color, markerfacecolor='#F18F01', 
                                   markeredgecolor='#1B4965', markeredgewidth=1)
                            ax.grid(True, alpha=0.3, linestyle='--', linewidth=0.5)
                        
                        elif visual_type == "Pie Chart":
                            # Data already filtered for positive values
                            colors = plt.cm.Set3(np.linspace(0, 1, len(chart_data)))
                            wedges, texts, autotexts = ax.pie(chart_data[label], labels=chart_data['Period'], 
                                                             autopct='%1.1f%%', startangle=90, 
                                                             colors=colors, textprops={'fontsize': 10, 'weight': 'bold'})
                            for autotext in autotexts:
                                autotext.set_color('white')
                                autotext.set_fontweight('bold')
                                autotext.set_fontsize(10)
                        
                        ax.set_title(f"{label} YTD Comparisons - {table_name}", fontsize=16, fontweight='bold', pad=15)
                        if visual_type != "Pie Chart":
                            ax.set_xlabel("Period", fontsize=14, fontweight='bold', labelpad=10)
                            ax.set_ylabel(label, fontsize=14, fontweight='bold', labelpad=10)
                            ax.tick_params(axis='both', which='major', labelsize=12, width=1, length=4)
                            
                            # Format large numbers on y-axis
                            if chart_data[label].max() > 1000:
                                ax.yaxis.set_major_formatter(plt.FuncFormatter(
                                    lambda x, p: f'{x/1000:.0f}K' if x >= 1000 else f'{x:.0f}'))
                            
                            # Straight x-axis labels for YTD
                            plt.setp(ax.get_xticklabels(), rotation=0, ha='center')
                            
                            ax.spines['top'].set_visible(False)
                            ax.spines['right'].set_visible(False)
                            ax.spines['left'].set_linewidth(1)
                            ax.spines['bottom'].set_linewidth(1)                            
                            ax.spines['top'].set_visible(False)
                            ax.spines['right'].set_visible(False)
                            ax.spines['left'].set_linewidth(1)
                            ax.spines['bottom'].set_linewidth(1)
                        
                        plt.tight_layout(pad=2.0)
                        st.pyplot(fig, use_container_width=True)
                        plt.close(fig)
                    
                    # Display data table
                    with st.expander("📊 View Data Table"):
                        st.dataframe(chart_data, use_container_width=True)
                    
                    # Generate PPT for YTD chart
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'
                    ppt_bytes = create_ppt_with_chart(
                        f"YTD {label} Analysis - {table_name} - {selected_sheet}",
                        chart_data,
                        "Period",
                        label,
                        ppt_type,
                        color_override
                    )
                    
                    st.download_button(
                        f"⬇️ Download YTD {label} PPT",
                        ppt_bytes,
                        f"ytd_{label.lower().replace(' ', '_')}_analysis.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_ytd_{label.lower().replace(' ', '_')}_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data

            def plot_branch_performance(tab, visual_type):
                with tab:
                    if not is_branch_analysis:
                        st.info("This visualization is only available for region analysis sheets")
                        return
                
                    ytd_act_col = None
                    for col in table_df.columns:
                        col_str = str(col).strip()
                        if col_str == "Act-YTD-25-26 (Apr to Mar)" or \
                           re.search(r'YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)\s*Act', col_str, re.IGNORECASE):
                            ytd_act_col = col
                            break
                
                    if ytd_act_col is None:
                        st.warning("Could not find YTD Act column for region performance analysis")
                        return
                
                    first_col = table_df.columns[0]
                    regions_df = table_df[~table_df[first_col].str.contains('|'.join(BRANCH_EXCLUDE_TERMS), na=False, case=False)].copy()
                    regions_df = regions_df.dropna(subset=[first_col, ytd_act_col])
                
                    if regions_df.empty:
                        st.warning("No branch data available after filtering")
                        return
                
                    regions_df[ytd_act_col] = pd.to_numeric(regions_df[ytd_act_col].astype(str).str.replace(',', ''), errors='coerce')
                    regions_df = regions_df.dropna(subset=[ytd_act_col])
                
                    if not ensure_numeric_data(regions_df, ytd_act_col):
                        st.warning("No numeric data available for region performance")
                        return
                
                    regions_df = regions_df.sort_values(by=ytd_act_col, ascending=False)
                    
                    st.markdown(f"### Branch Performance Analysis - {table_name}")
                    
                    # Use cloud-optimized visualization with interactivity
                    display_visualization_cloud_optimized(tab, "Branch Performance", regions_df, first_col, ytd_act_col, visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        top_region = regions_df.iloc[0]
                        st.metric("Top Performer", top_region[first_col], f"{top_region[ytd_act_col]:,.0f}")
                    with col2:
                        total_performance = regions_df[ytd_act_col].sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                    with col3:
                        avg_performance = regions_df[ytd_act_col].mean()
                        st.metric("Average Performance", f"{avg_performance:,.0f}")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.markdown("#### 🏆 Top 5 Regions")
                        top_5 = regions_df.head(5)[[first_col, ytd_act_col]]
                        st.dataframe(top_5, use_container_width=True, hide_index=True)
                    
                    with col2:
                        st.markdown("#### 📉 Bottom 5 Regions")
                        bottom_5 = regions_df.tail(5)[[first_col, ytd_act_col]]
                        st.dataframe(bottom_5, use_container_width=True, hide_index=True)
                
                    regions_df = make_jsonly_serializable(regions_df)
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'
                    ppt_bytes = create_ppt_with_chart(
                        f"Branch Performance - {table_name} - {selected_sheet}",
                        regions_df,
                        first_col,
                        ytd_act_col,
                        ppt_type
                    )
                
                    st.download_button(
                        "⬇️ Download Region Performance PPT",
                        ppt_bytes,
                        "region_performance.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_region_performance_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return regions_df

            def plot_branch_monthwise(tab, visual_type):
                with tab:
                    if not is_branch_analysis:
                        st.info("This visualization is only available for region analysis sheets")
                        return
                
                    regions_df = filtered_df[~filtered_df[first_col].str.contains('|'.join(BRANCH_EXCLUDE_TERMS), na=False, case=False)].copy()
                
                    act_cols = []
                    for col in table_df.columns:
                        col_str = str(col).lower()
                        for year in years:
                            if (re.search(rf'\bact\b.*(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)[-\s]*{year}', col_str, re.IGNORECASE) and 'ytd' not in col_str):
                                act_cols.append(col)
                
                    if not act_cols:
                        st.warning(f"Could not find monthly Act columns for the selected years ({', '.join(years)})")
                        return
                
                    month_order = ['Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar']
                
                    def get_sort_key(col_name):
                        col_name = str(col_name).lower()
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', col_name, re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', col_name)
                        month_idx = month_order.index(month_match.group(1).capitalize()) if month_match and month_match.group(1).capitalize() in month_order else 99
                        year = int(year_match.group(1)) if year_match else 0
                        return (year, month_idx)
                
                    act_cols_sorted = sorted(act_cols, key=get_sort_key)
                
                    monthwise_data = regions_df[[first_col] + act_cols_sorted].copy()
                
                    clean_col_names = []
                    for col in act_cols_sorted:
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', str(col), re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', str(col))
                        if month_match and year_match:
                            clean_col_names.append(f"{month_match.group(1).capitalize()}-{year_match.group(1)}")
                        else:
                            clean_col_names.append(str(col))
                
                    monthwise_data.columns = [first_col] + clean_col_names
                
                    for col in clean_col_names:
                        monthwise_data[col] = pd.to_numeric(monthwise_data[col].astype(str).str.replace(',', ''), errors='coerce')
                    
                    monthwise_data = monthwise_data.dropna()
                    
                    if monthwise_data.empty:
                        st.warning("No numeric data available for region monthwise performance after filtering")
                        return
                
                    st.markdown(f"### Branch Monthwise Performance ({', '.join(selected_years if selected_years else years)})")
                
                    chart_data = monthwise_data.melt(id_vars=first_col, 
                                                  var_name="Month", 
                                                  value_name="Value")
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    display_visualization_cloud_optimized(tab, "Branch Monthwise", chart_data, "Month", "Value", visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        total_by_month = monthwise_data[clean_col_names].sum()
                        best_month = total_by_month.idxmax()
                        st.metric("Best Month", best_month, f"{total_by_month[best_month]:,.0f}")
                    with col2:
                        avg_monthly = total_by_month.mean()
                        st.metric("Monthly Average", f"{avg_monthly:,.0f}")
                    with col3:
                        total_performance = total_by_month.sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line'
                    ppt_bytes = create_ppt_with_chart(
                        f"Region Monthwise Performance - {selected_sheet}",
                        chart_data,
                        "Month",
                        "Value",
                        ppt_type
                    )
                
                    st.download_button(
                        "⬇️ Download Region Monthwise PPT",
                        ppt_bytes,
                        "region_monthwise.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_region_monthwise_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data
            def plot_product_performance(tab, visual_type):
                with tab:
                    if not is_product_analysis:
                        st.info("This visualization is only available for product analysis sheets")
                        return
                
                    ytd_act_col = None
                    for col in table_df.columns:
                        col_str = str(col).strip()
                        if col_str == "Act-YTD-25-26 (Apr to Mar)" or \
                           re.search(r'YTD[-–\s]*\d{2}[-–\s]*\d{2}\s*\([^)]*\)\s*Act', col_str, re.IGNORECASE):
                            ytd_act_col = col
                            break
                
                    if ytd_act_col is None:
                        st.warning("Could not find YTD Act column for product performance analysis")
                        return
                
                    first_col = table_df.columns[0]
                    exclude_terms = ['Total', 'TOTAL', 'Grand', 'GRAND', 'Total Sales']
                    products_df = table_df[~table_df[first_col].str.contains('|'.join(exclude_terms), na=False, case=False)].copy()
                    products_df = products_df.dropna(subset=[first_col, ytd_act_col])
                
                    if products_df.empty:
                        st.warning("No product data available after filtering")
                        return
                
                    products_df[ytd_act_col] = pd.to_numeric(products_df[ytd_act_col].astype(str).str.replace(',', ''), errors='coerce')
                    products_df = products_df.dropna(subset=[ytd_act_col])
                
                    if not ensure_numeric_data(products_df, ytd_act_col):
                        st.warning("No numeric data available for product performance")
                        return
                
                    products_df = products_df.sort_values(by=ytd_act_col, ascending=False)
                    
                    st.markdown("### Product Performance Analysis")
                    
                    display_visualization_cloud_optimized(tab, "Product Performance", products_df, first_col, ytd_act_col, visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        top_product = products_df.iloc[0]
                        st.metric("Top Performer", top_product[first_col], f"{top_product[ytd_act_col]:,.0f}")
                    with col2:
                        total_performance = products_df[ytd_act_col].sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                    with col3:
                        avg_performance = products_df[ytd_act_col].mean()
                        st.metric("Average Performance", f"{avg_performance:,.0f}")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.markdown("#### 🏆 Top 5 Products")
                        top_5 = products_df.head(5)[[first_col, ytd_act_col]]
                        st.dataframe(top_5, use_container_width=True, hide_index=True)
                    
                    with col2:
                        st.markdown("#### 📉 Bottom 5 Products")
                        bottom_5 = products_df.tail(5)[[first_col, ytd_act_col]]
                        st.dataframe(bottom_5, use_container_width=True, hide_index=True)
                
                    products_df = make_jsonly_serializable(products_df)
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line' if visual_type == "Line Chart" else 'pie'
                    ppt_bytes = create_ppt_with_chart(
                        f"Product Performance - {selected_sheet}",
                        products_df,
                        first_col,
                        ytd_act_col,
                        ppt_type
                    )
                
                    st.download_button(
                        "⬇️ Download Product Performance PPT",
                        ppt_bytes,
                        "product_performance.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_product_performance_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return products_df

            def plot_product_monthwise(tab, visual_type):
                with tab:
                    if not is_product_analysis:
                        st.info("This visualization is only available for product analysis sheets")
                        return
                
                    act_cols = []
                    for col in table_df.columns:
                        col_str = str(col).lower()
                        for year in years:
                            if (re.search(rf'\bact\b.*(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)[-\s]*{year}', col_str, re.IGNORECASE) and 'ytd' not in col_str):
                                act_cols.append(col)
                
                    if not act_cols:
                        st.warning(f"Could not find monthly Act columns for the selected years ({', '.join(years)})")
                        return
                
                    month_order = ['Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar']
                
                    def get_sort_key(col_name):
                        col_name = str(col_name).lower()
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', col_name, re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', col_name)
                        month_idx = month_order.index(month_match.group(1).capitalize()) if month_match and month_match.group(1).capitalize() in month_order else 99
                        year = int(year_match.group(1)) if year_match else 0
                        return (year, month_idx)
                
                    act_cols_sorted = sorted(act_cols, key=get_sort_key)
                
                    first_col = table_df.columns[0]
                    exclude_terms = ['Total', 'TOTAL', 'Grand', 'GRAND', 'Total Sales']
                    products_df = table_df[~table_df[first_col].str.contains('|'.join(exclude_terms), na=False, case=False)].copy()
                    monthwise_data = products_df[[first_col] + act_cols_sorted].copy()
                    
                    clean_col_names = []
                    for col in act_cols_sorted:
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', str(col), re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', str(col))
                        if month_match and year_match:
                            clean_col_names.append(f"{month_match.group(1).capitalize()}-{year_match.group(1)}")
                        else:
                            clean_col_names.append(str(col))
                    
                    monthwise_data.columns = [first_col] + clean_col_names
                    
                    for col in clean_col_names:
                        monthwise_data[col] = pd.to_numeric(monthwise_data[col].astype(str).str.replace(',', ''), errors='coerce')
                    
                    monthwise_data = monthwise_data.dropna()
                    
                    if monthwise_data.empty:
                        st.warning("No numeric data available for product monthwise performance after filtering")
                        return
                
                    st.write(f"### Product Monthwise Performance ({', '.join(selected_years if selected_years else years)})")
                
                    chart_data = monthwise_data.melt(id_vars=first_col, 
                                                  var_name="Month", 
                                                  value_name="Value")
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    display_visualization_cloud_optimized(tab, "Product Monthwise", chart_data, "Month", "Value", visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        total_by_month = monthwise_data[clean_col_names].sum()
                        best_month = total_by_month.idxmax()
                        st.metric("Best Month", best_month, f"{total_by_month[best_month]:,.0f}")
                    with col2:
                        avg_monthly = total_by_month.mean()
                        st.metric("Monthly Average", f"{avg_monthly:,.0f}")
                    with col3:
                        total_performance = total_by_month.sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line'
                    ppt_bytes = create_ppt_with_chart(
                        f"Product Monthwise Performance - {selected_sheet}",
                        chart_data,
                        "Month",
                        "Value",
                        ppt_type
                    )
                
                    st.download_button(
                        "⬇️ Download Product Monthwise PPT",
                        ppt_bytes,
                        "product_monthwise.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_product_monthwise_ppt_{selected_sheet}_{sheet_index}"
                    )


            def plot_product_monthwise(tab, visual_type):
                with tab:
                    if not is_product_analysis:
                        st.info("This visualization is only available for product analysis sheets")
                        return
                
                    act_cols = []
                    for col in table_df.columns:
                        col_str = str(col).lower()
                        for year in years:
                            if (re.search(rf'\bact\b.*(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)[-\s]*{year}', col_str, re.IGNORECASE) and 'ytd' not in col_str):
                                act_cols.append(col)
                
                    if not act_cols:
                        st.warning(f"Could not find monthly Act columns for the selected years ({', '.join(years)})")
                        return
                
                    month_order = ['Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar']
                
                    def get_sort_key(col_name):
                        col_name = str(col_name).lower()
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', col_name, re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', col_name)
                        month_idx = month_order.index(month_match.group(1).capitalize()) if month_match and month_match.group(1).capitalize() in month_order else 99
                        year = int(year_match.group(1)) if year_match else 0
                        return (year, month_idx)
                
                    act_cols_sorted = sorted(act_cols, key=get_sort_key)
                
                    first_col = table_df.columns[0]
                    exclude_terms = ['Total', 'TOTAL', 'Grand', 'GRAND', 'Total Sales']
                    products_df = table_df[~table_df[first_col].str.contains('|'.join(exclude_terms), na=False, case=False)].copy()
                    monthwise_data = products_df[[first_col] + act_cols_sorted].copy()
                    
                    clean_col_names = []
                    for col in act_cols_sorted:
                        month_match = re.search(r'\b(apr|may|jun|jul|aug|sep|oct|nov|dec|jan|feb|mar)\b', str(col), re.IGNORECASE)
                        year_match = re.search(r'[-–](\d{2})\b', str(col))
                        if month_match and year_match:
                            clean_col_names.append(f"{month_match.group(1).capitalize()}-{year_match.group(1)}")
                        else:
                            clean_col_names.append(str(col))
                    
                    monthwise_data.columns = [first_col] + clean_col_names
                    
                    for col in clean_col_names:
                        monthwise_data[col] = pd.to_numeric(monthwise_data[col].astype(str).str.replace(',', ''), errors='coerce')
                    
                    monthwise_data = monthwise_data.dropna()
                    
                    if monthwise_data.empty:
                        st.warning("No numeric data available for product monthwise performance after filtering")
                        return
                
                    st.write(f"### Product Monthwise Performance ({', '.join(selected_years if selected_years else years)})")
                
                    chart_data = monthwise_data.melt(id_vars=first_col, 
                                                  var_name="Month", 
                                                  value_name="Value")
                    chart_data = make_jsonly_serializable(chart_data)
                    
                    display_visualization_cloud_optimized(tab, "Product Monthwise", chart_data, "Month", "Value", visual_type)
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        total_by_month = monthwise_data[clean_col_names].sum()
                        best_month = total_by_month.idxmax()
                        st.metric("Best Month", best_month, f"{total_by_month[best_month]:,.0f}")
                    with col2:
                        avg_monthly = total_by_month.mean()
                        st.metric("Monthly Average", f"{avg_monthly:,.0f}")
                    with col3:
                        total_performance = total_by_month.sum()
                        st.metric("Total Performance", f"{total_performance:,.0f}")
                
                    ppt_type = 'bar' if visual_type == "Bar Chart" else 'line'
                    ppt_bytes = create_ppt_with_chart(
                        f"Product Monthwise Performance - {selected_sheet}",
                        chart_data,
                        "Month",
                        "Value",
                        ppt_type
                    )
                
                    st.download_button(
                        "⬇️ Download Product Monthwise PPT",
                        ppt_bytes,
                        "product_monthwise.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"download_product_monthwise_ppt_{selected_sheet}_{sheet_index}"
                    )
                    return chart_data

            # Store chart data for master PPT
            # Plot visualizations in respective tabs
            budget_vs_actual_data = plot_budget_vs_actual(tabs_dict["Budget vs Actual"], visual_type)
            budget_data = plot_monthly_comparison(tabs_dict["Budget"], "Budget", visual_type)
            ly_data = plot_monthly_comparison(tabs_dict["LY"], "LY", visual_type)
            act_data = plot_monthly_comparison(tabs_dict["Act"], "Act", visual_type)
            gr_data = plot_monthly_comparison(tabs_dict["Gr"], "Gr", visual_type)
            ach_data = plot_monthly_comparison(tabs_dict["Ach"], "Ach", visual_type)
            ytd_budget_data = plot_ytd_comparison(tabs_dict["YTD Budget"], r'\bBudget\b.*YTD', "Budget", visual_type)
            ytd_ly_data = plot_ytd_comparison(tabs_dict["YTD LY"], r'\bLY\b.*YTD', "LY", visual_type)
            ytd_act_data = plot_ytd_comparison(tabs_dict["YTD Act"], r'\bAct\b.*YTD', "Act", visual_type)
            ytd_gr_data = plot_ytd_comparison(tabs_dict["YTD Gr"], r'\bGr\b.*YTD', "Gr", visual_type)
            ytd_ach_data = plot_ytd_comparison(tabs_dict["YTD Ach"], r'\bAch\b.*YTD', "Ach", visual_type)
            branch_performance_data = plot_branch_performance(tabs_dict["Branch Performance"], visual_type)
            branch_monthwise_data = plot_branch_monthwise(tabs_dict["Branch Monthwise"], visual_type)
            product_performance_data = plot_product_performance(tabs_dict["Product Performance"], visual_type)
            product_monthwise_data = plot_product_monthwise(tabs_dict["Product Monthwise"], visual_type)

            # Generate master PPT for all visualizations
            all_data = [
                ("Budget vs Actual", budget_vs_actual_data),
                ("Budget", budget_data),
                ("LY", ly_data),
                ("Act", act_data),
                ("Gr", gr_data),
                ("Ach", ach_data),
                ("YTD Budget", ytd_budget_data),
                ("YTD LY", ytd_ly_data),
                ("YTD Act", ytd_act_data),
                ("YTD Gr", ytd_gr_data),
                ("YTD Ach", ytd_ach_data),
                ("Branch Performance", branch_performance_data),
                ("Branch Monthwise", branch_monthwise_data),
                ("Product Performance", product_performance_data),
                ("Product Monthwise", product_monthwise_data)
            ]

            # Master PPT generation with fixed column handling and removed clickable titles
            if any(data is not None for _, data in all_data):
                st.sidebar.markdown("---")
                st.sidebar.subheader("📊 Download All Visuals")
                
                master_ppt = Presentation()
                
                for label, data in all_data:
                    if data is not None and (not isinstance(data, pd.DataFrame) or not data.empty):
                        try:
                            chart_data, x_col, y_col = get_chart_data_for_ppt(data, label, first_col, visual_type)
                            
                            if chart_data is None or x_col is None or y_col is None:
                                continue
                                
                            slide = master_ppt.slides.add_slide(master_ppt.slide_layouts[6])  # Use blank layout to avoid clickable titles
                            
                            # Determine color for Act-related charts
                            color_override = '#FF8C00' if 'Act' in label else None
                            
                            # Special handling for YTD charts
                            if "YTD" in label:
                                # Create a figure with straight x-axis labels
                                fig, ax = plt.subplots(figsize=(10, 6), dpi=150)
                                fig.patch.set_facecolor('white')
                                ax.set_facecolor('white')
                                
                                if visual_type == "Bar Chart":
                                    ax.bar(chart_data[x_col], chart_data[y_col], 
                                          color=color_override if color_override else '#2E86AB', 
                                          alpha=0.8, edgecolor='#1B4965', linewidth=1)
                                elif visual_type == "Line Chart":
                                    ax.plot(chart_data[x_col], chart_data[y_col], 
                                           marker='o', linewidth=3, markersize=8, 
                                           color=color_override if color_override else '#2E86AB', 
                                           markerfacecolor='#F18F01', 
                                           markeredgecolor='#1B4965', markeredgewidth=1)
                                    ax.grid(True, alpha=0.3, linestyle='--', linewidth=0.5)
                                elif visual_type == "Pie Chart":
                                    colors = plt.cm.Set3(np.linspace(0, 1, len(chart_data)))
                                    wedges, texts, autotexts = ax.pie(chart_data[y_col], labels=chart_data[x_col], 
                                                                     autopct='%1.1f%%', startangle=90, 
                                                                     colors=colors, textprops={'fontsize': 10, 'weight': 'bold'})
                                    for autotext in autotexts:
                                        autotext.set_color('white')
                                        autotext.set_fontweight('bold')
                                        autotext.set_fontsize(10)
                                
                                ax.set_title(f"{label} Analysis - {table_name}", fontsize=16, fontweight='bold', pad=15)
                                if visual_type != "Pie Chart":
                                    ax.set_xlabel(x_col, fontsize=14, fontweight='bold', labelpad=10)
                                    ax.set_ylabel(y_col, fontsize=14, fontweight='bold', labelpad=10)
                                    ax.tick_params(axis='both', which='major', labelsize=12, width=1, length=4)
                                    
                                    # Format large numbers on y-axis
                                    if chart_data[y_col].max() > 1000:
                                        ax.yaxis.set_major_formatter(plt.FuncFormatter(
                                            lambda x, p: f'{x/1000:.0f}K' if x >= 1000 else f'{x:.0f}'))
                                    
                                    # Straight x-axis labels for YTD
                                    plt.setp(ax.get_xticklabels(), rotation=0, ha='center')
                                    
                                    ax.spines['top'].set_visible(False)
                                    ax.spines['right'].set_visible(False)
                                    ax.spines['left'].set_linewidth(1)
                                    ax.spines['bottom'].set_linewidth(1)
                                
                                plt.tight_layout(pad=2.0)
                            else:
                                # For non-YTD charts, use the standard chart creation
                                fig = create_cloud_optimized_chart(chart_data, x_col, y_col, 'bar', 
                                                                 f"{label} Analysis - {table_name}", color_override)
                            
                            img_buffer = BytesIO()
                            fig.savefig(img_buffer, format='png', dpi=200, bbox_inches='tight')
                            plt.close(fig)
                            img_buffer.seek(0)
                            
                            # Add title as a text box (not clickable)
                            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                            tf = txBox.text_frame
                            tf.text = f"{label} Analysis - {table_name} - {selected_sheet}"
                            tf.paragraphs[0].font.size = Inches(0.25)
                            tf.paragraphs[0].font.bold = True
                            
                            slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5), width=Inches(8))
                            
                        except Exception as e:
                            st.warning(f"Error creating chart for {label}: {e}")
                            continue
                
                master_ppt_bytes = BytesIO()
                master_ppt.save(master_ppt_bytes)
                master_ppt_bytes.seek(0)
                
                st.sidebar.download_button(
                    "⬇️ Download All Charts as PPT",
                    master_ppt_bytes,
                    "all_visuals.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"download_all_charts_ppt_{selected_sheet}_{sheet_index}"
                )

            # Memory cleanup
            optimize_memory()

else:
    st.info("Please upload an Excel file to begin analysis.")
    
# Final memory cleanup
optimize_memory()
